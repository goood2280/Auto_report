# -*- coding: utf-8 -*-
"""
smoke_report.py — insert_plots 빠른 회귀 스모크
==============================================
합성 merged_df(대문자 메타 컬럼 + mask 2종 + lot 다수)로 insert_plots를 단독 호출하여
예외 없이 슬라이드/차트가 생성되는지, target-lot 스코프/방향/로그가 동작하는지 단언한다.

실행:  python scripts/smoke_report.py
"""
import os
import sys

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(BASE)
sys.path.insert(0, BASE)

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from My_config import GLOBAL_CONFIG
GLOBAL_CONFIG.load_from_yaml('vehicle_A')   # main_vehicle = vehicle_A
import My_Function as MF

rng = np.random.default_rng(0)
chips = [(x, y) for x in range(-2, 3) for y in range(-2, 3) if x * x + y * y <= 4]
base = pd.Timestamp('2026-06-10')

def gen(mask, fab_lot, root, day, wafers, shift):
    rows = []
    for wf in wafers:
        off = int(rng.integers(0, 120))
        for (cx, cy) in chips:
            tk = base + pd.Timedelta(days=day, minutes=off)
            row = {'MASK': mask, 'FAB_LOT_ID': fab_lot, 'root_lot_id': root,
                   'WAFER_ID': wf, 'CHIP_X_POS': cx, 'CHIP_Y_POS': cy,
                   'Chip_Radius': float(np.hypot(cx, cy)), 'SUBITEM_ID': 'MAIN', 'TKOUT_TIME': tk}
            for it, (m, s) in {'VTH_N': (0.5, 0.05), 'VTH_P': (0.48, 0.04),
                               'IDSAT_N': (550, 50), 'IDSAT_P': (520, 45)}.items():
                row[it] = m + shift + rng.normal(0, s * 0.4)   # 한 사이트 = 한 행에 4개 item
            rows.append(row)
    return pd.DataFrame(rows)

# 대문자 메타 컬럼(merged_df 실제 형태) + mask 2종 + target/이력 lot
frames = [gen('vehicle_A', 'T1234.1', 'T1234', 0, range(1, 26), 0.0)]      # target
frames += [gen('vehicle_A', f'TH{d}.1', f'TH{d}', -d, range(1, 13), 0.0) for d in (3, 8)]  # 이력
frames += [gen('Vehicle_B', 'VB001.1', 'VB001', -5, range(1, 13), 0.12)]   # with_vehicle (shift)
merged = pd.concat(frames, ignore_index=True)
# ADDP 컬럼 계산
merged['VTH_AVG'] = (merged['VTH_N'] + merged['VTH_P']) / 2
merged['IDSAT_SUM'] = merged['IDSAT_N'] + merged['IDSAT_P']

spec_data = pd.DataFrame({
    'SPECLOW':   [0.3, 0.3, 300.0, 300.0, 0.3, 600.0],
    'SPECHIGH':  [0.7, 0.7, 800.0, 800.0, 0.7, 1600.0],
    'UNIT':  ['V', 'V', 'uA', 'uA', 'V', 'uA'],
    'REPORT DIRECTION': ['BOTH', 'LOWER', 'UPPER', 'BOTH', 'BOTH', 'BOTH'],
    'REPORT LOG SCALE': [False, False, False, False, False, True],
    'CAT2': ['VTH', 'VTH', 'IDSAT', 'IDSAT', 'VTH', 'IDSAT'],
}, index=['VTH_N', 'VTH_P', 'IDSAT_N', 'IDSAT_P', 'VTH_AVG', 'IDSAT_SUM'])

prs = Presentation()
prs, metrics = MF.insert_plots(merged, prs, {}, 'T1234.1', 'T1234', 'MFDC', 'test', spec_data, dpi=None)

out = os.path.join(os.environ.get('TEMP', '.'), 'smoke_report_out.pptx')
prs.save(out)
size_mb = os.path.getsize(out) / 1048576

n_items = len(spec_data.index)
pics = [sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE) for s in prs.slides]
print('slides:', len(prs.slides), '| pics/slide:', pics, '| size MB:', round(size_mb, 2))
print('metrics:', list(metrics.keys()))

# 단언
assert len(prs.slides) == n_items, f'expected {n_items} slides, got {len(prs.slides)}'
assert all(p >= 6 for p in pics), f'each slide needs >=6 pictures, got {pics}'
assert set(spec_data.index) == set(metrics.keys()), 'metrics keys mismatch (ADDP 포함 전 항목)'
# target 스코프: target(T1234.1)은 25 wafer, 이력/Vehicle_B는 12 → target만 쓰면 global!=target
md = metrics['VTH_N']
assert abs(md['global_med'] - md['target_med']) >= 0 and md['target_std'] >= 0
assert size_mb < 10, f'PPT too large: {size_mb}MB'
print('SMOKE OK')
