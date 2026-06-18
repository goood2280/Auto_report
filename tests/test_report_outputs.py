import shutil
import subprocess
import sys
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
HTML_LIMIT_BYTES = 2 * 1024 * 1024
PPT_LIMIT_BYTES = 10 * 1024 * 1024


class ReportOutputTest(unittest.TestCase):
    def setUp(self) -> None:
        subprocess.run(
            [sys.executable, "-B", "scripts/create_temp_fixtures.py"],
            cwd=ROOT,
            check=True,
            stdout=subprocess.DEVNULL,
        )

        report_root = ROOT / "BASE_DIR" / "Report" / "VEHICLE_A"
        for child in ["Mail", "HTML"]:
            path = report_root / child
            if path.exists():
                shutil.rmtree(path)
            path.mkdir(parents=True, exist_ok=True)

    def test_main_generates_mail_html_and_ppt_with_size_limits(self) -> None:
        result = subprocess.run(
            [sys.executable, "-B", "Main.py", "VEHICLE_A"],
            cwd=ROOT,
            text=True,
            capture_output=True,
            timeout=60,
        )
        self.assertEqual(result.returncode, 0, result.stderr + result.stdout)

        ppt_files = sorted((ROOT / "BASE_DIR" / "Report" / "VEHICLE_A" / "Mail").glob("*.pptx"))
        html_files = sorted((ROOT / "BASE_DIR" / "Report" / "VEHICLE_A" / "HTML").glob("*.html"))

        self.assertEqual(len(ppt_files), 1)
        self.assertEqual(len(html_files), 1)
        self.assertLess(ppt_files[0].stat().st_size, PPT_LIMIT_BYTES)
        self.assertLess(html_files[0].stat().st_size, HTML_LIMIT_BYTES)

        prs = Presentation(ppt_files[0])
        self.assertGreaterEqual(len(prs.slides), 5)
        slide_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertIn("Auto Report", slide_text)
        self.assertIn("LOT001", slide_text)
        self.assertIn("Score Board - Index 1", slide_text)
        self.assertIn("Score Board - Index 2", slide_text)
        self.assertIn("Statistical Table - Index 1", slide_text)
        self.assertIn("Box Plot - Index 1", slide_text)
        self.assertIn("Trend - Index 1", slide_text)
        self.assertIn("WF Map - Index 1", slide_text)
        self.assertNotIn("Measurement Data Sample", slide_text)

        html = html_files[0].read_text(encoding="utf-8")
        self.assertIn("LOT001", html)
        self.assertIn("ADDP_ITEM_01", html)
        self.assertIn("Score Board", html)
        self.assertIn("Inline Table", html)
        self.assertIn("History", html)
        self.assertNotIn("Data Sample", html)


if __name__ == "__main__":
    unittest.main()
