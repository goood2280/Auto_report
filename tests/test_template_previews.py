import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


class TemplatePreviewTest(unittest.TestCase):
    def test_template_html_uses_mail_report_sections(self) -> None:
        html = (ROOT / "templates" / "template_report.html").read_text(encoding="utf-8")
        self.assertIn("Score Board", html)
        self.assertIn("Inline Table", html)
        self.assertIn("History", html)
        self.assertNotIn("Data Sample", html)

    def test_readme_embeds_template_preview_images(self) -> None:
        readme = (ROOT / "README.md").read_text(encoding="utf-8")
        expected_images = [
            "templates/previews/template_report_html.png",
            "templates/previews/template_report_ppt_cover.png",
            "templates/previews/template_report_scoreboard_index_1.png",
            "templates/previews/template_report_scoreboard_index_2.png",
            "templates/previews/template_report_statistical_table_index_1.png",
            "templates/previews/template_report_box_plot_index_1.png",
            "templates/previews/template_report_trend_index_1.png",
            "templates/previews/template_report_wafer_map_index_1.png",
        ]
        for relative_path in expected_images:
            self.assertTrue((ROOT / relative_path).exists(), relative_path)
            self.assertIn(f"]({relative_path})", readme)

    def test_template_ppt_contains_score_board_examples(self) -> None:
        prs = Presentation(ROOT / "templates" / "template_report.pptx")
        self.assertGreaterEqual(len(prs.slides), 5)
        slide_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertIn("Score Board - Index 1", slide_text)
        self.assertIn("Score Board - Index 2", slide_text)
        self.assertIn("Statistical Table - Index 1", slide_text)
        self.assertIn("Box Plot - Index 1", slide_text)
        self.assertIn("Trend - Index 1", slide_text)
        self.assertIn("WF Map - Index 1", slide_text)
        self.assertNotIn("Measurement Data Sample", slide_text)


if __name__ == "__main__":
    unittest.main()
