# -*- coding: utf-8 -*-
"""Auto Report 유틸리티 함수 모듈.

이 모듈은 Main.py에서 사용하는 모든 유틸리티 함수를 포함합니다.
- 데이터 쿼리: etdata_query(), inlinedata_query(), wipdata_query()
- PPT 생성: make_title_page(), insert_score_board(), insert_plots()
- 데이터 변환: Reformatize(), apply_style_by_index()
- 유틸리티: log_to_file(), clear_temp_inside_run()

사용법: from My_Function import *
"""

# ===================================================================
#  표준 라이브러리 (Standard Library Imports)
# ===================================================================
import os
import re
import gc
import glob
import time
import traceback
from datetime import datetime, timedelta, date

# ===================================================================
#  서드파티 라이브러리 (Third-party Imports)
# ===================================================================
import numpy as np
import pandas as pd

# ===================================================================
#  프로젝트 내부 모듈 (Project Internal Imports)
# ===================================================================
from My_config import GLOBAL_CONFIG
# NOTE: bigdataquery(getData)는 모듈 최상단에서 import하지 않는다.
#   병렬 렌더링 워커가 My_Function을 재import할 때마다 무거운 bigdataquery 재import·안내문이
#   뜨던 문제 방지. 실제 쿼리 시점(getData_with_retry)에만 지연 import한다.


# ===================================================================
#  유틸리티 함수 (Utility Functions)
# ===================================================================

def log_to_file(message, log_path):
    """타임스탬프와 함께 메시지를 로그 파일에 기록.

    Parameters
    ----------
    message : str
        로그에 기록할 메시지 문자열.
    log_path : str
        로그 파일의 절대 경로. 부모 디렉토리가 없으면 자동 생성.
    """
    try:
        os.makedirs(os.path.dirname(log_path), exist_ok=True)
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {message}\n")
        # 통합 로그 30MB 초과 시 오래된(앞) 내용 제거하고 최신 ~24MB만 유지
        if os.path.getsize(log_path) > 30 * 1024 * 1024:
            with open(log_path, "rb") as f:
                f.seek(-(24 * 1024 * 1024), os.SEEK_END)
                data = f.read()
            _nl = data.find(b'\n')
            if _nl != -1:
                data = data[_nl + 1:]
            with open(log_path, "wb") as f:
                f.write(data)
    except Exception as e:
        print(f"[WARN] log_to_file 실패: {e}")


def get_email_list(file_path, target_group, default_group='HOL', domain="@url"):
    """메일링 리스트 엑셀에서 수신 그룹(시트)의 KNOX_ID를 읽어 수신자 리스트를 반환.

    target_group에 해당하는 시트가 있으면 그 시트를, 없으면 default_group 시트를 사용한다.
    각 KNOX_ID에 '@'가 없으면 domain을 붙여 이메일을 구성하고,
    메일 API용 [{"email", "recipientType":"TO", "seq"}] 형태로 반환한다.
    """
    xls = pd.ExcelFile(file_path)
    sheet_name = target_group if target_group in xls.sheet_names else default_group
    df = xls.parse(sheet_name)
    email_list = df['KNOX_ID'].dropna().drop_duplicates(keep='first').tolist()
    domain = domain.strip()
    final_email_list = [
        {
        "email": email.strip() if "@" in email else f"{email.strip()}{domain}",
        "recipientType": "TO",
        "seq": i
        }
        for i, email in enumerate(email_list, start=1)
    ]
    return final_email_list


def reformatter_verify(reformatter):
    """reformatter(csv) DataFrame의 구조·내용 유효성을 검증.

    ``reformatter/{vehicle}_reformatter.csv`` 가 하위 처리(쿼리·피벗·ADDP 계산·
    spec 표시·차트)에 필요한 컬럼을 모두 갖추고 있는지, REAL/ADDP 행이 올바른지
    확인합니다.

    각 단계는 통과 여부를 ``[CHECK n] ... OK/FAIL`` 형태로 터미널에 출력합니다.

    검증 항목
    ---------
    [CHECK 1] DataFrame 비어 있지 않음.
    [CHECK 2] 필수 컬럼 존재 (CATEGORY, ITEMID, ALIAS, SCALE FACTOR, ABSOLUTE,
              ADDP FORM, REPORT ORDER, SPECLOW, SPECHIGH, REPORT DIRECTION).
    [CHECK 3] 선택 컬럼 존재 (UNIT, TARGET, REPORT LOG SCALE, CAT1, CAT2, PPT_ONLY) — 누락 시 경고만.
    [CHECK 4] ITEMID 중복 없음 (ITEMID는 유일 / ALIAS 중복은 허용).
    [CHECK 5] CATEGORY 값이 {'REAL', 'ADDP'} 범위 내.
    [CHECK 6] REAL 행은 ITEMID 보유.
    [CHECK 7] ADDP 행은 ADDP FORM(수식) 보유.

    Parameters
    ----------
    reformatter : pd.DataFrame
        검증 대상 reformatter DataFrame.

    Returns
    -------
    bool
        모든 필수 검증을 통과하면 True, 아니면 False.
    """
    print("=" * 60)
    print("[reformatter_verify] reformatter 검증 시작")
    print("=" * 60)

    # ── [CHECK 1] DataFrame 비어 있지 않음 ──
    print("[CHECK 1] DataFrame 비어있지 않음 검사...", end=" ")
    if reformatter is None or reformatter.empty:
        print("FAIL")
        print("  -> [ERROR] reformatter가 비어 있습니다.")
        return False
    print(f"OK (행 {len(reformatter)}개)")

    # ── [CHECK 2] 필수 열 존재 검사 ──
    # 하위 코드(Main.py/insert_plots)가 직접 참조하는 필수 컬럼
    required = [
        "CATEGORY", "ITEMID", "ALIAS", "SCALE FACTOR", "ABSOLUTE",
        "ADDP FORM", "REPORT ORDER", "SPECLOW", "SPECHIGH", "REPORT DIRECTION",
    ]
    print("[CHECK 2] 필수 열 존재 검사...", end=" ")
    missing = [c for c in required if c not in reformatter.columns]
    if missing:
        print("FAIL")
        print(f"  -> [ERROR] reformatter에 필수 컬럼이 없습니다: {missing}")
        return False
    print(f"OK ({len(required)}개 모두 존재)")

    # ── [CHECK 3] 선택 열 존재 검사 (없어도 동작하나 차트 표현에 사용) ──
    optional = ["UNIT", "TARGET", "REPORT LOG SCALE", "CAT1", "CAT2", "PPT_ONLY"]
    print("[CHECK 3] 선택 열 존재 검사...", end=" ")
    opt_missing = [c for c in optional if c not in reformatter.columns]
    if opt_missing:
        print(f"WARN (누락: {opt_missing} → 차트 표현 일부 제한)")
    else:
        print(f"OK ({len(optional)}개 모두 존재)")

    # ── [CHECK 4] ITEMID 중복 검사 (ITEMID는 유일 / ALIAS 중복은 허용) ──
    print("[CHECK 4] ITEMID 중복 검사 (ALIAS 중복은 허용)...", end=" ")
    itemid_nonnull = reformatter["ITEMID"].dropna().astype(str)
    itemid_nonnull = itemid_nonnull[itemid_nonnull.str.strip() != ""]
    dup_itemid = itemid_nonnull[itemid_nonnull.duplicated()].unique().tolist()
    if dup_itemid:
        print("FAIL")
        print(f"  -> [ERROR] reformatter ITEMID가 중복되었습니다(유일해야 함): {dup_itemid}")
        return False
    print(f"OK (ITEMID {itemid_nonnull.nunique()}개 모두 유일)")

    # ── [CHECK 5] CATEGORY 값 검사 ({'REAL', 'ADDP'} 범위) ──
    print("[CHECK 5] CATEGORY 값 검사...", end=" ")
    cats = set(reformatter["CATEGORY"].dropna().astype(str).str.upper().unique())
    unknown = cats - {"REAL", "ADDP"}
    if unknown:
        print("FAIL")
        print(f"  -> [ERROR] reformatter CATEGORY에 알 수 없는 값이 있습니다: {unknown}")
        return False
    print(f"OK (CATEGORY={sorted(cats)})")

    cat_upper = reformatter["CATEGORY"].astype(str).str.upper()

    # ── [CHECK 6] REAL 행 ITEMID 보유 검사 ──
    print("[CHECK 6] REAL 행 ITEMID 보유 검사...", end=" ")
    real_rows = reformatter[cat_upper == "REAL"]
    if real_rows["ITEMID"].isna().any():
        bad = real_rows[real_rows["ITEMID"].isna()]["ALIAS"].tolist()
        print("FAIL")
        print(f"  -> [ERROR] REAL 항목에 ITEMID가 비어 있습니다: {bad}")
        return False
    print(f"OK (REAL {len(real_rows)}행)")

    # ── [CHECK 7] ADDP 행 ADDP FORM(수식) 보유 검사 ──
    print("[CHECK 7] ADDP 행 ADDP FORM(수식) 보유 검사...", end=" ")
    addp_rows = reformatter[cat_upper == "ADDP"]
    if addp_rows["ADDP FORM"].isna().any():
        bad = addp_rows[addp_rows["ADDP FORM"].isna()]["ALIAS"].tolist()
        print("FAIL")
        print(f"  -> [ERROR] ADDP 항목에 ADDP FORM(수식)이 비어 있습니다: {bad}")
        return False
    print(f"OK (ADDP {len(addp_rows)}행)")

    print("-" * 60)
    print("[reformatter_verify] [PASS] 모든 검사 통과")
    print("=" * 60)
    return True


def extract_and_sort_numbers(x):
    """리스트 문자열에서 숫자를 추출하여 정렬된 문자열로 반환.

    문자열 형태의 리스트(예: "[3, 1, 2]")를 파싱하여
    오름차순 정렬된 쉼표 구분 문자열(예: "1, 2, 3")로 변환합니다.

    Parameters
    ----------
    x : str or list or tuple
        변환 대상. 문자열이면 eval로 파싱 시도.

    Returns
    -------
    str
        정렬된 숫자의 쉼표 구분 문자열.
    """
    if isinstance(x, str):
        try:
            x = eval(x)
        except Exception:
            return str(x)
    if isinstance(x, (list, tuple)):
        return ", ".join(str(i) for i in sorted(x))
    return str(x)


def remove_brackets(x):
    """문자열에서 대괄호·따옴표를 제거.

    wafer_id 리스트 등의 표시용 정리에 사용됩니다.
    예: "['A', 'B']" → "A, B"

    Parameters
    ----------
    x : str
        정리 대상 문자열.

    Returns
    -------
    str
        괄호와 따옴표가 제거된 문자열.
    """
    if isinstance(x, str):
        return x.replace("[", "").replace("]", "").replace("'", "")
    return str(x)


def replace_negatives_with_0(x):
    """음수값을 0으로 치환.

    Pass Rate 등 음수가 의미 없는 지표에서 사용됩니다.

    Parameters
    ----------
    x : numeric or any
        검사 대상 값.

    Returns
    -------
    numeric or any
        음수이면 0, 그 외에는 원래 값 그대로 반환.
    """
    try:
        if pd.notna(x) and float(x) < 0:
            return 0
    except (ValueError, TypeError):
        pass
    return x


def convert_target_data(x, suffixes_remove, replace_map, prefixes_remove=None):
    """항목명에서 지정 접두사/접미사를 제거하고 문자열을 치환.

    reformatter의 ITEMID나 ALIAS 정리에 활용됩니다.

    Parameters
    ----------
    x : str
        변환 대상 문자열.
    suffixes_remove : list[str] or None
        제거할 접미사 목록. 예: ["_AVG", "_MAX"]
    replace_map : dict or None
        치환 맵. 예: {"OLD": "NEW"}
    prefixes_remove : list[str] or None
        제거할 접두사 목록. 예: ["ET_", "PRE_"]

    Returns
    -------
    str
        변환된 문자열.
    """
    if not isinstance(x, str):
        return x
    for prefix in (prefixes_remove or []):
        if prefix and x.startswith(prefix):
            x = x[len(prefix):]
    for suffix in (suffixes_remove or []):
        if suffix and x.endswith(suffix):
            x = x[: -len(suffix)]
    for old, new in (replace_map or {}).items():
        x = x.replace(old, new)
    return x


def display_name(x):
    """사용자에게 '보여지는' 항목명 = convert_target_data(alias, ...) 후처리 결과.

    My_config의 prefixes_remove/suffixes_remove/replace_map를 적용해 표시용 이름을 만든다.
    데이터 키(merged_df 컬럼·metrics_dict 키·score_color 스케일 등)에는 절대 쓰지 말고
    '표시 문자열'에만 사용한다(내부 키는 원래 alias 유지).
    """
    try:
        return convert_target_data(
            x,
            getattr(GLOBAL_CONFIG, 'suffixes_remove', []) or [],
            getattr(GLOBAL_CONFIG, 'replace_map', {}) or {},
            getattr(GLOBAL_CONFIG, 'prefixes_remove', []) or [])
    except Exception:
        return x


def apply_style_by_index(row):
    """Pass Rate 값에 따라 셀 배경 색상 스타일을 반환 (pandas Styler용).

    색상 기준 (Color Thresholds):
        - ≥ 99.0% : 초록색  (#00B050) — 양호 (Good)
        - ≥ 95.0% : 연초록  (#92D050) — 보통 (Normal)
        - ≥ 90.0% : 노란색  (#FFFF00) — 주의 (Caution)
        - < 90.0% : 빨간색  (#FF0000) — 경고 (Alert)
        - NaN/빈값 : 진회색 (#555555) — 측정 없음 (No Data)

    Parameters
    ----------
    row : pd.Series
        한 행의 데이터 (각 셀 = Pass Rate 값).

    Returns
    -------
    list[str]
        각 셀에 대응하는 CSS 스타일 문자열 리스트.
    """
    styles = []
    for val in row:
        if pd.isna(val) or val == "":
            styles.append("background-color:#555555;color:white")
        else:
            try:
                v = float(val)
                if v >= 99.0:
                    styles.append("background-color:#00B050;color:white")
                elif v >= 95.0:
                    styles.append("background-color:#92D050")
                elif v >= 90.0:
                    styles.append("background-color:#FFFF00")
                else:
                    styles.append("background-color:#FF0000;color:white")
            except (ValueError, TypeError):
                styles.append("background-color:#555555;color:white")
    return styles


def Reformatize(data, ALIAS, FORMULA):
    """ADDP(Arithmetic Derived Data Parameter) 수식 컬럼을 계산하여 DataFrame에 추가.

    reformatter의 ADDP 카테고리 수식(ADDP FORM)을 파싱하여 새 컬럼을 생성합니다.
    수식 내 ``{ALIAS}`` 참조는 ``pivot.get("ALIAS")`` 로 치환되어 해당 컬럼값으로
    계산됩니다.

    ADDP 수식 예시:
        "1.0*(({VTH_N} + {VTH_P}) / 2)"  →  (VTH_N + VTH_P) / 2

    ADDP → ADDP 재귀 참조
    ---------------------
    ADDP 항목의 ``{}`` 안 ALIAS가 또 다른 ADDP일 수 있습니다(ADDP가 ADDP를 참조).
    이 함수는 고정점(fixpoint) 반복 루프로 매 패스마다 "지금 계산 가능한" 항목만
    먼저 풀어 그 결과를 ``data`` 컬럼으로 추가합니다. 따라서 어떤 ADDP가 아직
    계산되지 않은 다른 ADDP를 참조하면 그 패스에서는 실패(에러)로 남겨두었다가,
    참조 대상이 채워진 다음 패스에서 해소됩니다. 결과적으로 ``{ADDP}`` → ``{ADDP}``
    → ... → 최종 real item 까지 체인을 **재귀적으로 따라가** 계산됩니다.
    (체인 깊이만큼 반복하도록 최대 ``max(10, ADDP수 + 1)`` 패스 수행)

    Parameters
    ----------
    data : pd.DataFrame
        피벗 완료된 ET 데이터 (REAL item이 컬럼으로 존재).
    ALIAS : list[str]
        ADDP 항목의 별칭(새 컬럼명) 리스트.
    FORMULA : list[str]
        ADDP 수식 문자열 리스트. ``{COL}`` 형태로 다른 컬럼(REAL/ADDP)을 참조.

    Returns
    -------
    pd.DataFrame
        ADDP 컬럼이 추가된 DataFrame.
    """
    def addpf(formula, data):
        pivot = data
        def LOG(u, df):
            return np.log10(ABS(df))
        def POWER(a, b):
            return np.power(a, b)
        def sqrt(a):
            return np.sqrt(a)
        def ABS(a):
            return np.abs(a)
        def rmax(*args):
            df_max = pd.DataFrame()
            for arg in args:
                df_max = pd.concat([df_max, arg], axis=1)
            return df_max.max(axis=1)
        def rmin(*args):
            df_min = pd.DataFrame()
            for arg in args:
                df_min = pd.concat([df_min, arg], axis=1)
            return df_min.min(axis=1)
        def MA_Window(*args):
            # set data
            x_data, y_data = [], pd.DataFrame()
            spec, compliance = [], 10
            for arg in args:
                if isinstance(arg, list):
                    x_data = np.array(arg)
                elif isinstance(arg, str) or isinstance(arg, int) or isinstance(arg, float):
                    if isinstance(spec, list):
                        spec = np.log10(float(arg))
                    else:
                        compliance = abs(float(arg))
                else:
                    y_data = pd.concat([y_data, arg.apply(lambda a: np.log10(float(a) + 1E-14))], axis=1)
            # check & calculate data
            if x_data.shape[0] == y_data.shape[1]:
                # 계수계산
                df_coeffs = pd.DataFrame(np.polyfit(x_data, y_data.T, 2).T, index=y_data.index, columns=['a2', 'a1', 'a0'])
                # 판별식계산
                df_coeffs['discriminant'] = df_coeffs['a1']**2 - 4 * df_coeffs['a2'] * (df_coeffs['a0'] - spec)
                # 양방향margin계산
                df_coeffs['plus_margin'] = np.where(df_coeffs['discriminant']>=0, (-df_coeffs['a1'] + np.sqrt(df_coeffs['discriminant'])) / (2 * df_coeffs['a2']), np.nan)
                df_coeffs['minus_margin'] = np.where(df_coeffs['discriminant']>=0, (-df_coeffs['a1'] - np.sqrt(df_coeffs['discriminant'])) / (2 * df_coeffs['a2']), np.nan)
                # compliance적용
                df_coeffs['plus_margin'] = np.where(df_coeffs['plus_margin']>compliance, compliance, df_coeffs['plus_margin'])
                df_coeffs['minus_margin'] = np.where(df_coeffs['minus_margin']<-compliance, -compliance, df_coeffs['minus_margin'])
                # 볼록함수예외처리
                df_coeffs.loc[(df_coeffs['discriminant']<0)&(df_coeffs['a2']>0), ['plus_margin','minus_margin']] = 0
                # 오목함수예외처리
                df_coeffs.loc[df_coeffs['a2']<=0, ['plus_margin','minus_margin']] = [compliance,-compliance]
                # ovl_index계산
                df_coeffs['ovl_index'] = -0.5 * (df_coeffs['plus_margin'] + df_coeffs['minus_margin'])
                # ma_window계산
                df_coeffs[''] = df_coeffs['plus_margin'] - df_coeffs['minus_margin']
                # NEW ma_window계산
                df_coeffs['new'] = df_coeffs[['plus_margin', 'minus_margin']].abs().min(axis=1)
                return df_coeffs[['minus_margin','plus_margin','ovl_index','','new']]
            else:
                dummy = [np.nan for _ in range(5)]
                return pd.DataFrame([dummy for _ in range(y_data.shape[0])], index=y_data.index, columns=['minus_margin','plus_margin','ovl_index','','new'])
        def stddev(a):
            return a.groupby([pivot["root_lot_id"], pivot["wafer_id"], pivot["tkout_time"]]).transform(np.std)
        def std(a):
            return a.groupby([pivot["root_lot_id"], pivot["wafer_id"], pivot["tkout_time"]]).transform(np.std)
        def STD(a):
            return a.groupby([pivot["root_lot_id"], pivot["wafer_id"], pivot["tkout_time"]]).transform(np.std)
        def AVG(a):
            return a.groupby([pivot["root_lot_id"], pivot["wafer_id"], pivot["tkout_time"]]).transform(np.mean)
        try:
            a = eval(formula)
        except Exception:
            a = 'error'
        return a

    # pivot 된 data에 Reformatter(ADDP) 적용 — 고정점 반복으로 ADDP→ADDP 재귀 해소
    ALIAS_LEFT = []
    FORMULA_LEFT = []
    max_passes = max(10, len(ALIAS) + 1)   # ADDP 참조 체인 깊이만큼 반복 보장
    for _ in range(max_passes):
        calnum = len(ALIAS)   # 패스 시작 시점의 미계산 항목 수
        for alias, formula in zip(ALIAS, FORMULA):
            formula_tmp = formula
            formula = formula.replace('{', '(pivot.get("')
            formula = formula.replace('}', '"))')
            a = addpf(formula, data)   # ADDP 계산 실행
            if str(type(a)) == "<class 'str'>":
                # 아직 못 푼 항목(참조 ADDP가 아직 미계산 등) → 다음 패스에서 재시도
                ALIAS_LEFT.append(alias)
                FORMULA_LEFT.append(formula_tmp)
                continue
            else:
                if isinstance(a, pd.DataFrame) and a.shape[1] > 1:
                    data[[alias + '_' + col if col != '' else alias for col in a]] = a
                else:
                    data[alias] = a
        ALIAS = ALIAS_LEFT
        FORMULA = FORMULA_LEFT
        ALIAS_LEFT = []
        FORMULA_LEFT = []

        if calnum == len(ALIAS):   # 한 패스 동안 진전이 없으면(더 이상 해소 불가) 종료
            break

    if ALIAS:
        print(f"[WARN] Reformatize: 참조를 해소하지 못한 ADDP 항목: {ALIAS}")
    return data


def clear_temp_inside_run():
    """temp_*.png 임시 파일 삭제."""
    for f in glob.glob("temp_*.png"):
        try:
            os.remove(f)
        except OSError:
            pass


def clear_anomaly_inside_run():
    """anomaly_*.png 임시 파일 삭제."""
    for f in glob.glob("anomaly_*.png"):
        try:
            os.remove(f)
        except OSError:
            pass


#  랏 완료 후 삭제하지 않고 남겨둘 '진단/근거' 파일 접두어.
# RUN/TEMP에서 한 사이클 종료 시 지울 대상 = '임시 그림파일(이미지)'만.
_RUN_TEMP_IMAGE_EXTS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tif', '.tiff', '.webp')


def clear_run_temp_files():
    """RUN/TEMP 폴더의 '임시 그림파일(png 등 이미지)'만 삭제(그 외 파일·폴더는 유지).

    랏 리포트 1건 생성이 끝난 뒤 호출 — Trend PNG({alias}.png) 등 임시 이미지만 비운다.
    (AI 인풋파일은 RUN/AI에 별도 보관하며 삭제하지 않는다. anomaly_basis 등 비이미지
     산출물도 그대로 남긴다.) HTML은 이미지가 base64로 내장돼 있어 삭제 후에도 정상.
    """
    _tdir = os.path.join('RUN', 'TEMP')
    if not os.path.isdir(_tdir):
        return
    for _root, _dirs, _files in os.walk(_tdir):
        for _f in _files:
            if _f.lower().endswith(_RUN_TEMP_IMAGE_EXTS):   # 이미지 파일만 삭제
                try:
                    os.remove(os.path.join(_root, _f))
                except OSError:
                    pass


# ===================================================================
#  PPT 생성 함수 (PowerPoint Generation)
# ===================================================================

def make_title_page(vehicle, lot_id, step_merged):
    """코드로 직접 그리는 표지(첫 페이지) Presentation을 생성하여 반환.

    템플릿 파일을 읽지 않고 빈 16:9 슬라이드에 표지를 직접 구성합니다.
      - 중앙 대형 제목: "{vehicle} HOL Auto Report"
      - 그 아래: "{lot_id} / {step_merged}"  (step_merged = step_desc(step_id))
      - 우상단: 날짜(오늘)
      - 상/하단 네이비 액센트 바 + 중앙 구분선으로 깔끔한 디자인

    Args:
        vehicle, lot_id, step_merged: 표지에 표기할 정보.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from datetime import datetime as _dt

    FONT = GLOBAL_CONFIG.theme_font_family
    NAVY = RGBColor(*GLOBAL_CONFIG.theme_title_color)
    GREY = RGBColor(0x60, 0x6A, 0x7A)
    RECT = MSO_SHAPE.RECTANGLE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    # 배경 흰색
    bg = slide.background
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(255, 255, 255)

    def _bar(top, height):
        b = slide.shapes.add_shape(RECT, Inches(0), top, SW, height)
        b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
        b.shadow.inherit = False
        return b

    # 상단/하단 네이비 액센트 바
    _bar(Inches(0), Inches(0.45))
    _bar(SH - Inches(0.45), Inches(0.45))

    # 우상단 발행 날짜
    today = _dt.now().strftime('%Y-%m-%d')
    date_box = slide.shapes.add_textbox(SW - Inches(3.6), Inches(0.7), Inches(3.3), Inches(0.4))
    dp = date_box.text_frame.paragraphs[0]
    dp.text = f"{today}"
    dp.alignment = PP_ALIGN.RIGHT
    dp.font.size = Pt(13); dp.font.color.rgb = GREY; dp.font.name = FONT

    # 중앙 대형 제목
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), SW - Inches(1.2), Inches(1.5))
    ttf = title_box.text_frame; ttf.word_wrap = True
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = ttf.paragraphs[0]
    tp.text = f"{vehicle} HOL Auto Report"
    tp.alignment = PP_ALIGN.CENTER
    tp.font.size = Pt(46); tp.font.bold = True; tp.font.color.rgb = NAVY; tp.font.name = FONT

    # 중앙 구분선 (제목 아래)
    line = slide.shapes.add_shape(RECT, SW / 2 - Inches(2.2), Inches(4.18), Inches(4.4), Pt(2.5))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()
    line.shadow.inherit = False

    # 서브타이틀: lot_id / step_desc(step_id)
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.35), SW - Inches(1.2), Inches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = f"{lot_id}  /  {step_merged}"
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(24); sp.font.bold = False; sp.font.color.rgb = GREY; sp.font.name = FONT

    return prs


def insert_score_board(VIP_group, prs, lot_id, title, spec_data=None, config=None, item_link_cells=None):
    """Pass Rate Score Board (PPT). 컬럼이 (lot, wafer) MultiIndex이면 lot별로 분리 표기.
       (단일 레벨 wafer 컬럼도 하위호환 동작 — 모두 lot_id 한 그룹으로 처리)."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.xmlchemy import OxmlElement

    NAVY = RGBColor(*GLOBAL_CONFIG.theme_title_color)
    FONT = GLOBAL_CONFIG.theme_font_family
    LOTBG = RGBColor(0xD9, 0xE1, 0xF2)
    WAFBG = RGBColor(0xF0, 0xF0, 0xF0)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    def _set_borders(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        for t in ['lnL', 'lnR', 'lnT', 'lnB']:
            for el in tcPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}' + t):
                tcPr.remove(el)
        _bi = 0
        for t in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = OxmlElement(t); ln.set('w', '12700'); ln.set('cmpd', 'sng')
            sf = OxmlElement('a:solidFill'); sc = OxmlElement('a:srgbClr'); sc.set('val', '333333')
            sf.append(sc); ln.append(sf); tcPr.insert(_bi, ln); _bi += 1

    def _style(cell, text, bg, fg, bold=False, sz=7):
        cell.text = str(text)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.01); cell.margin_right = Inches(0.01)
        cell.margin_top = Inches(0.0); cell.margin_bottom = Inches(0.0)
        cell.text_frame.word_wrap = False
        _set_borders(cell)
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(sz); par.font.bold = bold; par.font.name = FONT
            par.font.color.rgb = fg; par.alignment = PP_ALIGN.CENTER

    # ---- 컬럼 정규화: [(orig_col, lot, wafer), ...] / lot 정렬(target 먼저), wafer 오름차순 ----
    def _as_lw(c):
        if isinstance(c, tuple) and len(c) == 2:
            l, w = c
        else:
            l, w = lot_id, c
        try:
            w = int(float(str(w).replace('#', '')))
        except Exception:
            pass
        return (str(l), w)
    _cols = list(VIP_group.columns)
    _lw = [_as_lw(c) for c in _cols]
    _lots = sorted({x[0] for x in _lw}, key=lambda l: (0 if str(l) == str(lot_id) else 1, str(l)))
    # 측정된 (lot, wafer) → 원본 컬럼 매핑
    _present = {}
    for c, lw in zip(_cols, _lw):
        _present[(lw[0], lw[1])] = c
    # HTML Score Board와 동일하게 '실제 데이터가 있는 wafer'만 표시(#1~25 고정 배열 폐기).
    # lot은 target 먼저(_lots 정렬), lot 내 wafer 오름차순.
    order = []   # (orig_col, lot, wafer)
    for _lot in _lots:
        _wafs = sorted([w for (l, w) in _present if l == _lot],
                       key=lambda w: (0, w) if isinstance(w, int) else (1, str(w)))
        for w in _wafs:
            order.append((_present[(_lot, w)], _lot, w))

    chunk_size = 30
    total_pages = (len(VIP_group) - 1) // chunk_size + 1

    for page in range(total_pages):
        chunk_df = VIP_group.iloc[page * chunk_size: (page + 1) * chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.58))
        header_shape.fill.solid(); header_shape.fill.fore_color.rgb = NAVY
        header_shape.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(0.2), Inches(0.06), Inches(12.7), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        page_suffix = f" ({page+1}/{total_pages})" if total_pages > 1 else ""
        p.text = f"Score Board - {title}{page_suffix}"
        p.font.size = Pt(20); p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255); p.font.name = FONT

        ncols = 1 + len(order)
        nrows = len(chunk_df) + 2   # 헤더 2행(lot / wafer)
        table_height = Inches(min(6.4, 0.5 + len(chunk_df) * 0.18))
        tbl = slide.shapes.add_table(nrows, ncols, Inches(0.22), Inches(0.72),
                                     Inches(12.89), table_height).table

        item_w = 2.25   # ITEM명 폭(약 45자→30자, 15자 축소). 남는 폭은 wafer 열로 분배
        ww = max(0.14, (12.89 - item_w) / max(len(order), 1))
        tbl.columns[0].width = Inches(item_w)
        for j in range(1, ncols):
            tbl.columns[j].width = Inches(ww)   # 모든 wafer 열 동일 너비

        # 헤더: ITEM 세로 병합 + lot 가로 병합 + #wafer
        tbl.cell(0, 0).merge(tbl.cell(1, 0)); _style(tbl.cell(0, 0), "ITEM", NAVY, WHITE, True, 9)
        tbl.cell(0, 0).text_frame.word_wrap = False   # ITEM 헤더 한 줄
        k = 0
        while k < len(order):
            _lot = order[k][1]; k2 = k
            while k2 + 1 < len(order) and order[k2 + 1][1] == _lot:
                k2 += 1
            c0, c1 = 1 + k, 1 + k2
            if c1 > c0:
                tbl.cell(0, c0).merge(tbl.cell(0, c1))
            _style(tbl.cell(0, c0), str(_lot), LOTBG, BLACK, True, 8)
            k = k2 + 1
        for jj, (_oc, _lot, _waf) in enumerate(order):
            _style(tbl.cell(1, 1 + jj), f"#{_waf}", NAVY, WHITE, True, 7)

        # 본문
        for i, (idx, row) in enumerate(chunk_df.iterrows()):
            r = i + 2
            base = RGBColor(245, 247, 250) if i % 2 == 1 else WHITE
            _style(tbl.cell(r, 0), display_name(str(idx)), base, BLACK, False, 7)   # 표시명 후처리, 45자 한 줄 수용 위해 폰트 7pt
            if item_link_cells is not None:   # Item명 → 차트 슬라이드 내부 링크(insert_plots 후 연결)
                item_link_cells[str(idx)] = (slide, tbl.cell(r, 0))
            tbl.cell(r, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            tbl.cell(r, 0).text_frame.word_wrap = False   # ITEM명 한 줄(줄바꿈 방지)
            for jj, (_oc, _lot, _waf) in enumerate(order):
                cell = tbl.cell(r, 1 + jj)
                val = row[_oc] if (_oc is not None and _oc in row.index) else None
                # 연속 색상(HTML과 동일), ITEM별 스케일 override 지원
                if pd.notna(val) and str(val).strip() != "":
                    try:
                        v = float(val); txt = f"{v:.1f}"
                        _bg, _fg = GLOBAL_CONFIG.score_color(v, str(idx))
                        bg = RGBColor(*GLOBAL_CONFIG._hex2rgb(_bg))
                        fg = RGBColor(*GLOBAL_CONFIG._hex2rgb(_fg))
                    except (ValueError, TypeError):
                        txt = str(val); bg = RGBColor(*GLOBAL_CONFIG._hex2rgb(GLOBAL_CONFIG.score_color_na)); fg = WHITE
                else:
                    txt = ""; bg = RGBColor(*GLOBAL_CONFIG._hex2rgb(GLOBAL_CONFIG.score_color_na)); fg = WHITE
                _style(cell, txt, bg, fg, False, 7)

    return prs


def link_scoreboard_items(item_link_cells, item_slide_map):
    """Score Board의 Item명 셀을 해당 아이템 차트 슬라이드로 내부 링크(파랑·밑줄).

    insert_score_board가 insert_plots보다 먼저 실행되므로(차트 슬라이드 미생성),
    insert_plots가 item_slide_map을 반환한 뒤 이 함수로 링크를 건다.
    item_link_cells = {item_name: (score_board_slide, item_cell)}.
    """
    if not item_link_cells or not item_slide_map:
        return
    from pptx.dml.color import RGBColor as _RGB
    for _it, _pair in item_link_cells.items():
        _tgt = item_slide_map.get(_it)
        if _tgt is None:
            continue
        try:
            _sl, _cell = _pair
            _runs = _cell.text_frame.paragraphs[0].runs
            if not _runs:
                continue
            _run = _runs[0]
            _run.font.color.rgb = _RGB(0x00, 0x33, 0xCC); _run.font.underline = True
            _add_internal_slide_link(_run, _sl, _tgt)
        except Exception:
            continue


def _add_internal_slide_link(run, source_slide, target_slide):
    """run 텍스트에 '같은 PPT 내 target_slide로 점프'하는 내부 하이퍼링크를 건다.

    PowerPoint 슬라이드 점프는 a:hlinkClick(action=ppaction://hlinksldjump) + 대상 슬라이드
    파트로의 relationship(r:id)으로 구현한다. 슬라이드 순서가 바뀌어도 relationship이
    파트를 직접 가리키므로 링크는 유지된다.
    """
    try:
        from pptx.oxml.ns import qn
        from pptx.opc.constants import RELATIONSHIP_TYPE as _RT
        rId = source_slide.part.relate_to(target_slide.part, _RT.SLIDE)
        rPr = run._r.get_or_add_rPr()
        for _h in rPr.findall(qn('a:hlinkClick')):
            rPr.remove(_h)
        _hl = rPr.makeelement(qn('a:hlinkClick'), {})
        _hl.set(qn('r:id'), rId)
        _hl.set('action', 'ppaction://hlinksldjump')
        rPr.append(_hl)   # rPr 자식 스키마상 hlinkClick은 fill/latin 뒤 → append로 순서 유지
        return True
    except Exception as _e:
        print(f"[WARN] 내부 슬라이드 링크 실패: {_e}")
        return False


def insert_findings_page(prs, findings, after_index=2, title="■ Anomaly 상세 (통계 자동 분석)",
                         main_vehicle=None, radius_zones=(60, 100), item_slide_map=None,
                         rule_trace=None):
    """코드 통계 분석 Finding 전체를 1개 슬라이드로 만들어 Score Board 뒤에 삽입.

    HTML [0]에는 우선순위 상위 N건만 보이고, 전체 상세는 이 PPT 페이지를 참조.
    after_index 위치(보통 1=title + Score Board 페이지수)로 슬라이드를 이동시킨다.
    상단에는 비교 기준·robust 산포 계산법·radius zone 정의를 '참고사항'으로 1회만 안내한다.
    rule_trace 전달 시 맨 뒤에 '전체 anomaly rule 체크 결과'(매칭/해당없음) 요약 블록을 덧붙인다.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import re as _re

    NAVY = RGBColor(*GLOBAL_CONFIG.theme_title_color)
    FONT = GLOBAL_CONFIG.theme_font_family
    # 신호등 3색(anomaly_engine._SEV_COLOR 미러): 빨강 이상 / 주황 주의 / 노랑 측정이상추정
    sev_dot = {"CRITICAL": RGBColor(0xD6, 0x27, 0x28),
               "WARNING":  RGBColor(0xF5, 0x9E, 0x0B),
               "NOTICE":   RGBColor(0xEA, 0xB3, 0x08),
               "INFO":     RGBColor(0x5D, 0x6D, 0x7E)}
    sev_label = {"CRITICAL": "이상", "WARNING": "주의", "NOTICE": "주의", "INFO": "참고"}

    def _strip(t):
        return _re.sub(r'\s+', ' ', _re.sub(r'<[^>]+>', '', str(t))).strip()

    n_crit = sum(1 for f in findings if f.get("severity") == "CRITICAL")
    n_warn = sum(1 for f in findings if f.get("severity") == "WARNING")
    _WHITE = RGBColor(255, 255, 255); _DIV = RGBColor(0xC9, 0xD2, 0xE0)
    _hcnt = [("CRITICAL", "이상", n_crit), ("WARNING", "주의", n_warn)]

    # ── 참고사항(1페이지에만): 비교 기준·robust 산포 계산법·radius zone·판정/우선순위 ──
    _veh = main_vehicle or '제품'
    try:
        _rc, _rm = float(radius_zones[0]), float(radius_zones[1])
    except Exception:
        _rc, _rm = 60.0, 100.0
    _dsp = getattr(GLOBAL_CONFIG, 'anomaly_lot_dispersion_ratio', 2.0)
    _note_lines = [
        f"※ 참고: 모든 판정은 대상 lot의 'wafer 단위'로 보며, 비교 기준은 제품({_veh}) 전체의 'wafer별' 통계입니다.",
        "· robust 산포 = 값들의 1.4826×MAD(0이면 IQR/1.349, 그래도 0이면 std).",
        "· '보통 wafer 산포' = 제품 전 lot 각 wafer의 '내부 robust 산포'들의 중앙값.   "
        "'N배' = 대상 wafer 내부 산포 / 보통 wafer 산포.",
        f"· 위치(radius zone): Center ≤ {_rc:g}, Middle {_rc:g}~{_rm:g}, Edge > {_rm:g}.",
        "· [판정 기준]",
        "   - 이상(빨강): spec을 벗어난 측정 point가 하나라도 있으면 이상. (median 이동은 판정에 사용하지 않음)",
        f"   - 주의(주황): spec은 모두 만족하지만, 대상 lot의 어떤 wafer 내부 산포가 '보통 wafer 산포'의 {_dsp:g}배를 넘으면 주의.",
        "   - 그 외: 참고.",
        "· [우선순위 P] — 값이 클수록 위에 정렬. R_max=최대 wafer spec-out 비율(out pt/측정 pt), "
        "N_wf=spec-out wafer 수, D=최대 wafer 산포배수.",
        "   - 이상: P = 20000 + 100·R_max + N_wf/100      - 주의: P = 10000 + 100·D      (동점 시 REPORT ORDER 오름차순)",
    ]

    # ── 카테고리(cat2)별 그룹핑 — 우선순위 상 '첫 등장' 순으로 카테고리 배열, 카테고리 내는 우선순위 유지 ──
    def _fcat(_f):
        _c = str(_f.get('cat2', '') or '').strip()
        if _c:
            return _c
        return '지식 판정(규칙)' if _f.get('type') == 'KNOWLEDGE' else '기타'
    _cat_order, _by_cat = [], {}
    for _f in findings:
        _c = _fcat(_f)
        if _c not in _by_cat:
            _by_cat[_c] = []; _cat_order.append(_c)
        _by_cat[_c].append(_f)
    # 렌더 블록: 카테고리 헤더('H') + 그 카테고리 finding('F')들 → 카테고리별로 구분되어 보임
    _blocks = []
    for _c in _cat_order:
        _blocks.append(('H', _c))
        for _f in _by_cat[_c]:
            _blocks.append(('F', _f))
    # 전체 anomaly rule 체크 결과(매칭/해당없음) — Finding 뒤에 요약 블록으로 덧붙임
    _rt = list(rule_trace or [])
    if _rt:
        _rt_hit = sum(1 for _t in _rt if _t.get('matched'))
        _blocks.append(('H', f"Rule Check 결과 (전체 {len(_rt)}개 · 매칭 {_rt_hit} · 해당없음 {len(_rt) - _rt_hit})"))
        for _t in _rt:   # 매칭 먼저, 그다음 해당없음
            if _t.get('matched'):
                _blocks.append(('R', _t))
        for _t in _rt:
            if not _t.get('matched'):
                _blocks.append(('R', _t))
    # 블록(헤더+finding)이 많으면 페이지 분할
    PER_PAGE = 26
    _pages = [_blocks[i:i + PER_PAGE] for i in range(0, len(_blocks), PER_PAGE)] or [[]]
    _total = len(_pages)

    def _render_cat_header(_p, _txt):
        """카테고리 구분 헤더(진한 남색 볼드) — finding 그룹 앞에 배치."""
        _p.space_before = Pt(7)
        _hr = _p.add_run(); _hr.text = f"▍ {_txt}"
        _hr.font.bold = True; _hr.font.size = Pt(12)
        _hr.font.color.rgb = NAVY; _hr.font.name = FONT
    _created = []
    for _pi, _chunk in enumerate(_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _created.append(slide)
        bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(255, 255, 255)

        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.62))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = NAVY; hdr.line.fill.background()
        htf = hdr.text_frame; htf.margin_left = Inches(0.2)
        hp = htf.paragraphs[0]; hp.text = ""

        def _hrun(text, color, sz=18, bold=True, _hp=hp):
            rr = _hp.add_run(); rr.text = text
            rr.font.size = Pt(sz); rr.font.bold = bold; rr.font.name = FONT; rr.font.color.rgb = color

        _sfx = f" ({_pi + 1}/{_total})" if _total > 1 else ""
        _hrun(f"{title}{_sfx}   —   ", _WHITE)
        for _k, (_sev, _lbl, _cnt) in enumerate(_hcnt):
            _hrun("● ", sev_dot.get(_sev, _WHITE))
            _hrun(f"{_lbl} {_cnt}", _WHITE)
            if _k < len(_hcnt) - 1:
                _hrun("  |  ", _DIV)
        hp.font.size = Pt(18); hp.font.bold = True
        hp.font.color.rgb = RGBColor(255, 255, 255); hp.font.name = FONT

        box = slide.shapes.add_textbox(Inches(0.3), Inches(0.72),
                                       prs.slide_width - Inches(0.6), prs.slide_height - Inches(0.95))
        tf = box.text_frame; tf.word_wrap = True
        _firstpara = True

        if _pi == 0:      # 참고사항은 1페이지에만
            for _ln in _note_lines:
                _np = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                _nr = _np.add_run(); _nr.text = _ln
                _nr.font.size = Pt(8); _nr.font.italic = True
                _nr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80); _nr.font.name = FONT
            _spp = tf.add_paragraph(); _spp.text = ""

        if not findings and not _rt:
            p = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
            p.text = "유의미한 통계 이상 없음"; p.font.size = Pt(12); p.font.name = FONT
        else:
            if not findings and _pi == 0:   # finding은 없지만 rule 체크 결과는 표시
                p = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                p.text = "유의미한 통계 이상 없음 (아래는 전체 rule 체크 내역)"
                p.font.size = Pt(11); p.font.name = FONT; p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            # 페이지가 카테고리 중간(finding)부터 시작하면 해당 카테고리 헤더를 상단에 재표기
            if _chunk and _chunk[0][0] == 'F':
                _hp2 = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                _render_cat_header(_hp2, _fcat(_chunk[0][1]) + " (계속)")
            elif _chunk and _chunk[0][0] == 'R':   # Rule Check 섹션이 다음 페이지로 이어짐
                _hp2 = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                _render_cat_header(_hp2, "Rule Check 결과 (계속)")
            for _bk, _bv in _chunk:
                if _bk == 'H':   # 카테고리 헤더
                    p = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                    _render_cat_header(p, _bv)
                    continue
                if _bk == 'R':   # Rule Check 한 줄(매칭 O / 해당없음 ·)
                    _t = _bv
                    p = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                    p.space_before = Pt(1)
                    _ok = bool(_t.get('matched'))
                    _mr = p.add_run(); _mr.text = ("● " if _ok else "○ ")
                    _mr.font.bold = True; _mr.font.size = Pt(10); _mr.font.name = FONT
                    _mr.font.color.rgb = (RGBColor(0xD6, 0x27, 0x28) if _ok else RGBColor(0x9A, 0xA0, 0xA6))
                    _nm = p.add_run(); _nm.text = f"[{_t.get('kind','')}] {_t.get('name','')}"
                    _nm.font.bold = True; _nm.font.size = Pt(10); _nm.font.name = FONT
                    _nm.font.color.rgb = (RGBColor(0x1A, 0x1A, 0x1A) if _ok else RGBColor(0x6B, 0x72, 0x80))
                    _rs = p.add_run(); _rs.text = f"  —  {_t.get('result','')}"
                    _rs.font.size = Pt(9); _rs.font.name = FONT
                    _rs.font.color.rgb = (RGBColor(0x55, 0x55, 0x55) if _ok else RGBColor(0x9A, 0xA0, 0xA6))
                    continue
                f = _bv
                p = tf.paragraphs[0] if _firstpara else tf.add_paragraph(); _firstpara = False
                p.space_before = Pt(2)   # finding 간 약간의 간격
                sev = f.get("severity", "INFO")
                r0 = p.add_run(); r0.text = "● "
                r0.font.bold = True; r0.font.size = Pt(11)
                r0.font.color.rgb = sev_dot.get(sev, RGBColor(0x5D, 0x6D, 0x7E)); r0.font.name = FONT
                # 제목: 아이템 이름 부분에 '해당 차트 슬라이드로 점프'하는 내부 하이퍼링크 부여
                _full_title = f"{sev_label.get(sev, sev)} {f.get('title', '')}"
                _it_raw = f.get('item', '')
                _it_disp = display_name(_it_raw) if _it_raw else ''
                _tgt_slide = item_slide_map.get(_it_raw) if (item_slide_map and _it_raw) else None
                if _tgt_slide is not None and _it_disp and (_it_disp in _full_title):
                    _pos = _full_title.find(_it_disp)
                    for _txt, _islink in ((_full_title[:_pos], False),
                                          (_it_disp, True),
                                          (_full_title[_pos + len(_it_disp):], False)):
                        if not _txt:
                            continue
                        _rr = p.add_run(); _rr.text = _txt
                        _rr.font.bold = True; _rr.font.size = Pt(11); _rr.font.name = FONT
                        if _islink:
                            _rr.font.color.rgb = RGBColor(0x00, 0x33, 0xCC); _rr.font.underline = True
                            _add_internal_slide_link(_rr, slide, _tgt_slide)
                        else:
                            _rr.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                else:
                    r1 = p.add_run()
                    r1.text = _full_title
                    r1.font.bold = True; r1.font.size = Pt(11)
                    r1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A); r1.font.name = FONT
                det = _strip(f.get("detail", ""))
                if det:
                    r2 = p.add_run(); r2.text = "  —  " + det
                    r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55); r2.font.name = FONT
                # ── wafer별 median/σ (이상·주의 항목) — 요청: PPT 상세에 wafer 통계 표시 ──
                _ws = f.get('wafer_stats') or {}
                if _ws:
                    _items_ws = sorted(_ws.items(), key=lambda z: (z[0] is None, z[0]))
                    _cap = 8
                    _cells = []
                    for _w, _v in _items_ws[:_cap]:
                        try:
                            _cells.append(f"#{_w} med={float(_v['median']):.3g}/σ={float(_v['std']):.3g}")
                        except Exception:
                            continue
                    _more = f" 외 {len(_items_ws) - _cap}매" if len(_items_ws) > _cap else ""
                    _hdr = "wafer median/σ"
                    _rm = f.get('rep_median'); _rs = f.get('rep_stddev')
                    _rep = []
                    if _rm is not None:
                        try: _rep.append(f"med={float(_rm):.3g}")
                        except Exception: pass
                    if _rs is not None:
                        try: _rep.append(f"σ={float(_rs):.3g}")
                        except Exception: pass
                    if _rep:
                        _hdr += f" (대표 {', '.join(_rep)})"
                    if _cells:
                        _wp = tf.add_paragraph(); _wp.space_before = Pt(1)
                        _wr = _wp.add_run(); _wr.text = f"     {_hdr}: " + ", ".join(_cells) + _more
                        _wr.font.size = Pt(8); _wr.font.italic = True
                        _wr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80); _wr.font.name = FONT

    # 생성한 페이지(들)를 순서 유지하며 after_index 위치로 이동 (Score Board 바로 뒤)
    try:
        sldIdLst = prs.slides._sldIdLst
        _n = len(_created)
        _tail = list(sldIdLst)[-_n:]
        for _el in _tail:
            sldIdLst.remove(_el)
        _pos = min(after_index, len(list(sldIdLst)))
        for _off, _el in enumerate(_tail):
            sldIdLst.insert(_pos + _off, _el)
    except Exception as e:
        print(f"[WARN] Anomaly 상세 페이지 위치 이동 실패: {e}")
    return prs


def _wafer_circle_params(df, x_col, y_col, rad_col=None, mask_col=None,
                         main_vehicle=None, wafer_radius_mm=150.0):
    """실제 150mm wafer 원 (cx, cy, semi_x, semi_y, aspect)을 '플롯 좌표계(격자)' 단위로 산출.

    ● MASK(vehicle) 매칭 — vehicle별 shot 크기가 다를 수 있으므로, 같은 MASK==main_vehicle
      shot만으로 계산한다(mask_col·main_vehicle이 있고 데이터가 충분할 때). mask_col 미지정 시
      df에서 'MASK'/'mask'를 자동 인식, main_vehicle 미지정 시 GLOBAL_CONFIG.vehicle 사용.

    ● shot 크기(mm) 산출 — 격자 1칸당 물리거리 = shot_width/shot_height. Chip_Radius r 과
      shot 중점 격자좌표(chip_x_adj, chip_y_adj)의 관계
          r² = kx²·(x-cx)² + ky²·(y-cy)²   (kx=shot_width, ky=shot_height, (cx,cy)=wafer 중심)
      를 선형화 r² = A·x² + B·y² + p·x + q·y + C 로 최소자승 fit:
          cx=-p/2A,  cy=-q/2B,  shot_width=√A,  shot_height=√B.
      보통 wafer당 수십~수백 shot이라 5pt 이상 충분. 좌표는 (x,y) 중복 제거해 'shot 레이아웃'만 쓴다
      (측정 밀도 편향 제거). rad_col은 측정프레임의 Chip_Radius(Data Extractor MASK==vehicle 기반).

    ● 150mm 원 — 격자 반축: semi_x = 150/shot_width,  semi_y = 150/shot_height.
      aspect = shot_height/shot_width (= ky/kx, set_aspect 값) → 타원이 정원으로 보인다.
      shot이 원 밖으로 나가도 그대로 둔다(강제 포함하지 않음 — 실제 wafer 경계를 정확히 표시).

    Chip_Radius가 없거나 fit이 degenerate하면 shot 크기를 알 수 없어, 등방(aspect=1) shot-중심
    bounding 원(중심=좌표평균, 반경=최대편차×1.05)으로 폴백한다. 좌표 자체가 무효면 None.
    """
    try:
        import numpy as _np
        if x_col not in df.columns or y_col not in df.columns:
            return None
        # ── MASK(vehicle) 매칭: 같은 vehicle shot만 사용 ──
        if mask_col is None:
            mask_col = 'MASK' if 'MASK' in df.columns else ('mask' if 'mask' in df.columns else None)
        if main_vehicle is None:
            try:
                main_vehicle = GLOBAL_CONFIG.get('vehicle')
            except Exception:
                main_vehicle = None
        _src = df
        if mask_col and mask_col in df.columns and main_vehicle is not None:
            _sub = df[df[mask_col] == main_vehicle]
            if len(_sub) >= 6:
                _src = _sub

        x = pd.to_numeric(_src[x_col], errors='coerce')
        y = pd.to_numeric(_src[y_col], errors='coerce')
        _hasr = bool(rad_col and rad_col in _src.columns)
        if _hasr:
            r = pd.to_numeric(_src[rad_col], errors='coerce')
            m0 = x.notna() & y.notna() & r.notna()
        else:
            m0 = x.notna() & y.notna()
        if int(m0.sum()) < 6:
            return None
        x = x[m0].to_numpy(dtype=float); y = y[m0].to_numpy(dtype=float)
        cxg = float(x.mean()); cyg = float(y.mean())

        # ── Chip_Radius fit → shot_width/height·중심·150mm 원 ──
        if _hasr:
            r = r[m0].to_numpy(dtype=float)
            # (x,y) 좌표별 대표 radius(중앙값) — shot 레이아웃만 남겨 측정 밀도 편향 제거
            _lay = (pd.DataFrame({'x': x, 'y': y, 'r': r})
                    .groupby(['x', 'y'], as_index=False)['r'].median())
            xs = _lay['x'].to_numpy(dtype=float)
            ys = _lay['y'].to_numpy(dtype=float)
            rs = _lay['r'].to_numpy(dtype=float)
            if len(xs) >= 6 and float(_np.nanstd(rs)) > 1e-9:
                # r² = A·x² + B·y² + p·x + q·y + C  (x/y 스케일 분리, 축 정렬 가정)
                M = _np.column_stack([xs ** 2, ys ** 2, xs, ys, _np.ones_like(xs)])
                try:
                    sol, *_ = _np.linalg.lstsq(M, rs ** 2, rcond=None)
                    A, B, p, q, C = [float(v) for v in sol]
                    if _np.isfinite([A, B, p, q]).all() and A > 0 and B > 0:
                        cx = -p / (2.0 * A); cy = -q / (2.0 * B)
                        shot_w = A ** 0.5; shot_h = B ** 0.5      # 격자 1칸당 mm(shot 크기)
                        semi_x = wafer_radius_mm / shot_w         # 150mm 원의 격자 x 반축
                        semi_y = wafer_radius_mm / shot_h         # 150mm 원의 격자 y 반축
                        aspect = shot_h / shot_w                  # ky/kx → set_aspect(정원)
                        # degenerate만 배제: 반축이 측정 shot 편차의 절반보다 작으면(원이 데이터보다
                        # 터무니없이 작음) 폴백. 정상 fit은 semi_x≥max|x-cx| 이므로 통과.
                        _rx = float(_np.abs(xs - cx).max()); _ry = float(_np.abs(ys - cy).max())
                        if (_np.isfinite([cx, cy, semi_x, semi_y, aspect]).all()
                                and semi_x > 0 and semi_y > 0 and 0.1 <= aspect <= 10.0
                                and semi_x >= _rx * 0.5 and semi_y >= _ry * 0.5):
                            return (cx, cy, semi_x, semi_y, aspect)
                except Exception:
                    pass

        # ── 폴백: radius 없음/실패 → 등방 shot-중심 bounding 원(shot 크기 미상) ──
        max_dist = float(_np.sqrt((x - cxg) ** 2 + (y - cyg) ** 2).max())
        if max_dist > 0:
            _r = max_dist * 1.05
            return (cxg, cyg, _r, _r, 1.0)
        return None
    except Exception:
        return None


def _wfmap_aspect(params):
    """set_aspect에 넘길 값 — params가 있으면 ky/kx, 없으면 'equal'."""
    try:
        return params[4] if params else 'equal'
    except Exception:
        return 'equal'


def _add_wafer_circle(ax, params, color='#9aa4b0', lw=0.7, zorder=0):
    """_wafer_circle_params 결과(타원 반축)를 ax에 wafer 경계로 추가(칩 산점 뒤).

    호출부에서 ax.set_aspect(_wfmap_aspect(params))로 aspect를 맞추면 정원(true circle)로 보인다.
    """
    if not params:
        return
    try:
        from matplotlib.patches import Ellipse
        cx, cy, sx, sy = params[0], params[1], params[2], params[3]
        ax.add_patch(Ellipse((cx, cy), width=2.0 * sx, height=2.0 * sy, fill=False,
                             edgecolor=color, linewidth=lw, zorder=zorder))
    except Exception:
        pass


def _wfmap_shot_pitch(vals):
    """인접 shot 센터 간 거리(격자 pitch) — 정렬된 unique 좌표값의 '양수 diff 중앙값'.

    shot 사이 간격이 없다고 보고, 이 pitch를 shot의 가로/세로 크기로 써서 빈틈없이 채운다.
    측정 pt 수(밀도)와 무관하게 격자 간격 자체를 쓰므로 shot 크기가 항목별로 동일하다.
    """
    try:
        import numpy as _np
        u = _np.unique(_np.round(
            pd.to_numeric(pd.Series(vals), errors='coerce').dropna().to_numpy(dtype=float), 6))
        if len(u) < 2:
            return 1.0
        dif = _np.diff(u); dif = dif[dif > 0]
        return float(_np.median(dif)) if len(dif) else 1.0
    except Exception:
        return 1.0


def _draw_wfmap_shots(ax, xs, ys, px, py, values=None, colors=None,
                      cmap=None, norm=None, zorder=1):
    """shot을 인접 센터 간 거리(px,py) 크기의 '사각형'으로 빈틈없이 그린다(shot 사이 gap 제거).

    데이터 좌표계에서 한 변이 pitch(px/py)인 사각형을 각 center에 배치하면, aspect가 어떻든
    인접 사각형의 변이 맞닿아 위/아래·좌우 줄(간격)이 생기지 않는다.
    values+cmap+norm → 연속 컬러맵(정상 WF MAP) / colors → 개별 색(spec-out). PatchCollection 반환.
    """
    try:
        import numpy as _np
        from matplotlib.collections import PatchCollection
        from matplotlib.patches import Rectangle
        xs = _np.asarray(xs, dtype=float); ys = _np.asarray(ys, dtype=float)
        if len(xs) == 0:
            return None
        _hx, _hy = float(px) / 2.0, float(py) / 2.0
        _rects = [Rectangle((x - _hx, y - _hy), px, py) for x, y in zip(xs, ys)]
        if values is not None:
            pc = PatchCollection(_rects, cmap=cmap, norm=norm,
                                 edgecolors='none', linewidths=0, zorder=zorder)
            pc.set_array(_np.asarray(values, dtype=float))
        else:
            pc = PatchCollection(_rects, facecolors=colors,
                                 edgecolors='none', linewidths=0, zorder=zorder)
        ax.add_collection(pc)
        return pc
    except Exception:
        return None


def _wfmap_axis_limits(gx0, gx1, gy0, gy1, xpad, ypad, circ, margin_frac=0.08):
    """WF MAP 축 범위 = 데이터 범위 + wafer 경계(타원)를 '모두' 포함하는 (xlo,xhi,ylo,yhi).

    경계가 데이터 바깥으로 나가도 잘리지 않도록 반축 기준 작은 여백(margin_frac)을 둔다.
    """
    xlo, xhi = gx0 - xpad, gx1 + xpad
    ylo, yhi = gy0 - ypad, gy1 + ypad
    if circ:
        _cx, _cy, _sx, _sy = circ[0], circ[1], circ[2], circ[3]
        _mx = abs(_sx) * margin_frac; _my = abs(_sy) * margin_frac
        xlo = min(xlo, _cx - _sx - _mx); xhi = max(xhi, _cx + _sx + _mx)
        ylo = min(ylo, _cy - _sy - _my); yhi = max(yhi, _cy + _sy + _my)
    return xlo, xhi, ylo, yhi


def _wfmap_png_bytes(fig, dpi, colors=64):
    """WF MAP figure를 팔레트(PNG-8) 양자화+최적화로 인코딩해 PNG bytes 반환.

    WF MAP은 표시 크기가 작고(≈60~90px) 색 종류가 적어(칩 컬러맵 + 원 테두리 + 배경)
    팔레트 양자화 시 화질 손실이 사실상 없으면서 파일 크기는 절반 이하로 줄어든다.
    → score board 80 index·wafer map 30종 + anomaly 5종이 모두 full로 나와도 HTML<2MB.
    Pillow 없거나 실패 시 원본 PNG bytes로 폴백한다.
    """
    import io as _io
    buf = _io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', pad_inches=0.02, facecolor='white')
    try:
        from PIL import Image as _PILImage
        buf.seek(0)
        _im = _PILImage.open(buf).convert('RGB')
        _pal = _im.quantize(colors=int(colors), dither=_PILImage.Dither.NONE)
        _out = _io.BytesIO()
        _pal.save(_out, format='PNG', optimize=True)
        return _out.getvalue()
    except Exception:
        return buf.getvalue()


def render_wafer_wfmaps_b64(df, item, min_pts=50, lot_prefix=None,
                            direction='BOTH', size_in=0.85, dpi=None,
                            spec_low=None, spec_high=None, by_lot=False):
    """index의 'wafer별 첫 측정(가장 이른 tkout) WF MAP'을 {wafer_id_str: base64}로 반환.

    by_lot=True이면 같은 root_lot_id의 형제 lot을 구분하기 위해 (lot, wafer)별로
    그려 키를 'f"{FAB_LOT_ID}|{int(wafer)}"' 형태로 반환한다(Score Board lot 분리 표시용).

    각 wafer에서 측정 point 수가 min_pts 이상인 경우만 포함(sparse 측정은 제외).
    Score Board에서 wafer 열 아래 행으로 정렬해 넣기 위한 용도.
    """
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import base64, io

    if dpi is None:   # HTML용 WF MAP 해상도는 My_config.html_wfmap_dpi로 관리
        dpi = getattr(GLOBAL_CONFIG, 'html_wfmap_dpi', 110)

    if item not in df.columns:
        return {}

    def _pick(*names):
        for n in names:
            if n in df.columns:
                return n
        return None
    cx = _pick('CHIP_X_ADJ', 'CHIP_X_POS', 'chip_x_pos')
    cy = _pick('CHIP_Y_ADJ', 'CHIP_Y_POS', 'chip_y_pos')
    cw = _pick('WAFER_ID', 'wafer_id')
    ct = _pick('TKOUT_TIME', 'tkout_time')
    cl = _pick('FAB_LOT_ID', 'fab_lot_id')
    crad = _pick('Chip_Radius', 'chip_radius')
    cm = _pick('MASK', 'mask')                        # vehicle 매칭용(150mm 원 shot 크기 계산)
    if not (cx and cy and cw and ct):
        return {}

    d = df[[c for c in [cx, cy, cw, ct, cl, crad, cm, item] if c]].dropna(subset=[item]).copy()
    if lot_prefix and cl:
        d = d[d[cl].astype(str).str.startswith(str(lot_prefix))]
    if len(d) == 0:
        return {}
    d[ct] = pd.to_datetime(d[ct], errors='coerce')
    _circ = _wafer_circle_params(d, cx, cy, crad, mask_col=cm)   # 실제 150mm wafer 원(vehicle 매칭)

    norm = _wfmap_norm(direction, spec_low, spec_high, d[item])
    cmap = _wfmap_cmap(direction)
    # 전체 wafer 격자가 잘리지 않도록 공통 축범위 + 반칩 여백 (가장자리 칩 보존)
    _gx0, _gx1 = float(d[cx].min()), float(d[cx].max())
    _gy0, _gy1 = float(d[cy].min()), float(d[cy].max())
    _xr = (_gx1 - _gx0) or 1.0; _yr = (_gy1 - _gy0) or 1.0
    _xpad = _xr / max(int(d[cx].nunique()) - 1, 1) * 0.7 if d[cx].nunique() > 1 else _xr * 0.1
    _ypad = _yr / max(int(d[cy].nunique()) - 1, 1) * 0.7 if d[cy].nunique() > 1 else _yr * 0.1
    # 축 표시범위(원 포함) — 모든 wafer 공통. shot은 인접 센터 간격(pitch) 크기 사각형으로 그림(gap 제거)
    _xlo, _xhi, _ylo, _yhi = _wfmap_axis_limits(_gx0, _gx1, _gy0, _gy1, _xpad, _ypad, _circ)
    _pit_x = _wfmap_shot_pitch(d[cx]); _pit_y = _wfmap_shot_pitch(d[cy])

    out = {}
    # by_lot: (lot, wafer)별로 그려 형제 lot을 분리. 아니면 wafer별(lot은 첫 측정 선택).
    group_keys = [cl, cw] if (by_lot and cl) else [cw]
    for gkey, wdf in d.groupby(group_keys):
        if by_lot and cl:
            lot_val, wid = gkey
        else:
            wid = gkey if not isinstance(gkey, tuple) else gkey[0]
            lot_val = None
        # 같은 (lot,)wafer가 여러 tkout이면 가장 이른 측정 그룹 선택
        gk = [k for k in [cl, ct] if k]
        if gk:
            sizes = wdf.groupby(gk).size().reset_index(name='n')
            sizes = sizes[sizes['n'] >= min_pts]
            if len(sizes) == 0:
                continue
            first = sizes.sort_values(ct).iloc[0]
            sel = pd.Series(True, index=wdf.index)
            for k in gk:
                sel &= (wdf[k] == first[k])
            g = wdf[sel]
        else:
            if len(wdf) < min_pts:
                continue
            g = wdf
        if len(g) == 0:
            continue
        fig, ax = plt.subplots(figsize=(size_in, size_in))
        _draw_wfmap_shots(ax, g[cx].astype(float).values, g[cy].astype(float).values,
                          _pit_x, _pit_y, values=g[item].astype(float).values, cmap=cmap, norm=norm)
        # 원 테두리: radius 기반 원점·150mm 원을 '진한 검정'으로. 내부 배경색 없음(흰색)
        #  — 칩(shot) 컬러가 내부를 채우므로 별도 배경 fill 불필요.
        _add_wafer_circle(ax, _circ, color='#000000', lw=1.0)
        # aspect=ky/kx로 맞춰 wafer 경계가 정원으로 보이게(찌그러짐 방지)
        ax.set_xticks([]); ax.set_yticks([]); ax.set_aspect(_wfmap_aspect(_circ), adjustable='box'); ax.set_facecolor('white')
        # 방향: 왼쪽=chip_x_adj 작은 쪽(정방향), 위쪽=chip_y_adj 작은 쪽(y축 반전). 경계까지 포함(안 잘리게)
        ax.set_xlim(_xlo, _xhi)
        ax.set_ylim(_yhi, _ylo)
        for sp in ax.spines.values():
            sp.set_visible(False)
        fig.subplots_adjust(left=0.02, right=0.98, top=0.98, bottom=0.02)
        _png = _wfmap_png_bytes(fig, dpi)   # 팔레트 양자화(PNG-8)로 용량 최소화
        plt.close(fig)
        try:
            _wkey = str(int(wid))
        except (ValueError, TypeError):
            _wkey = str(wid)
        if by_lot and lot_val is not None:
            out[f"{lot_val}|{_wkey}"] = base64.b64encode(_png).decode('utf-8')
        else:
            out[_wkey] = base64.b64encode(_png).decode('utf-8')
    return out


def render_specout_wfmaps_b64(merged_df, item, spec_low=None, spec_high=None,
                              target_lot=None, max_maps=25,
                              size_in=0.62, dpi=None):
    """spec-out(=flier) 칩맵을 측정(lot, wafer, tkout) 단위로 그려 [(label, b64), ...]로 반환.

    [0] Anomaly Trend Chart 의 SPEC OUT 항목 우측에 붙이는 용도.
      - 칩 색: spec 통과=회색(#bdbdbd), spec 이탈=빨강(#d32f2f).
      - 대상: spec-out 칩이 1개 이상 있는 측정(tkout)만.
      - 정렬/상한: ① target_lot(FAB_LOT_ID)의 spec-out wafer는 wafer_id 오름차순으로 '모두' 표시
                   (lot 25매가 다 spec이면 25장 다 나옴 — max_maps와 무관하게 전량 보장),
                   ② 남는 칸은 다른 lot을 TKOUT_TIME 최신순으로 max_maps 총개수까지 채움.
      - label = f"{ROOT_LOT_ID} #{WAFER_ID}".
    spec_low/spec_high는 호출부에서 REPORT DIRECTION을 이미 반영한 값(UPPER=하한 None,
    LOWER=상한 None)을 넘겨야 Trend의 SPEC OUT 판정과 일치한다.
    """
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    import base64, io

    if dpi is None:   # HTML용 WF MAP 해상도는 My_config.html_wfmap_dpi로 관리
        dpi = getattr(GLOBAL_CONFIG, 'html_wfmap_dpi', 110)

    if item not in merged_df.columns:
        return []
    if spec_low is None and spec_high is None:
        return []

    def _pick(*names):
        for n in names:
            if n in merged_df.columns:
                return n
        return None
    cx = _pick('CHIP_X_ADJ', 'CHIP_X_POS', 'chip_x_pos')
    cy = _pick('CHIP_Y_ADJ', 'CHIP_Y_POS', 'chip_y_pos')
    cw = _pick('WAFER_ID', 'wafer_id')
    ct = _pick('TKOUT_TIME', 'tkout_time')
    croot = _pick('ROOT_LOT_ID', 'root_lot_id')
    clot = _pick('FAB_LOT_ID', 'fab_lot_id')
    crad = _pick('Chip_Radius', 'chip_radius')
    cm = _pick('MASK', 'mask')                        # vehicle 매칭용(150mm 원 shot 크기 계산)
    if not (cx and cy and cw and ct):
        return []

    keep = [c for c in [cx, cy, cw, ct, croot, clot, crad, cm, item] if c]
    d = merged_df[keep].dropna(subset=[item]).copy()
    if len(d) == 0:
        return []
    d[ct] = pd.to_datetime(d[ct], errors='coerce')
    _circ = _wafer_circle_params(d, cx, cy, crad, mask_col=cm)   # 실제 150mm wafer 원(vehicle 매칭)
    vals = pd.to_numeric(d[item], errors='coerce')
    lo = None if spec_low is None else float(spec_low)
    hi = None if spec_high is None else float(spec_high)
    out_mask = pd.Series(False, index=d.index)
    if lo is not None:
        out_mask = out_mask | (vals < lo)
    if hi is not None:
        out_mask = out_mask | (vals > hi)
    d['_specout'] = out_mask.values

    # 측정 단위 그룹: (root_lot, fab_lot, wafer, tkout) — spec-out 칩이 있는 측정만
    gcols = [c for c in [croot, clot, cw, ct] if c]
    groups = []
    for gkey, g in d.groupby(gcols, dropna=False):
        if not g['_specout'].any():
            continue
        gvals = gkey if isinstance(gkey, tuple) else (gkey,)
        gd = dict(zip(gcols, gvals))
        groups.append((gd, g))
    if not groups:
        return []

    def _is_tgt(gd):
        return target_lot is not None and clot is not None and str(gd.get(clot)) == str(target_lot)

    def _waf(gd):
        try:
            return int(gd.get(cw))
        except (ValueError, TypeError):
            return 10 ** 9
    tgt = sorted([x for x in groups if _is_tgt(x[0])], key=lambda x: _waf(x[0]))
    rest = sorted([x for x in groups if not _is_tgt(x[0])],
                  key=lambda x: (x[0].get(ct) if pd.notna(x[0].get(ct)) else pd.Timestamp.min),
                  reverse=True)
    # target lot의 spec-out wafer는 max_maps와 무관하게 전량 표시, 남는 칸만 나머지로 채움
    if max_maps and max_maps > 0:
        _n_rest = max(0, int(max_maps) - len(tgt))
        ordered = tgt + rest[:_n_rest]
    else:
        ordered = tgt + rest

    # 전체 wafer 격자가 잘리지 않도록 공통 축범위 + 반칩 여백 (모든 소형맵 동일 범위)
    gx0, gx1 = float(d[cx].min()), float(d[cx].max())
    gy0, gy1 = float(d[cy].min()), float(d[cy].max())
    xr = (gx1 - gx0) or 1.0; yr = (gy1 - gy0) or 1.0
    xpad = xr / max(int(d[cx].nunique()) - 1, 1) * 0.7 if d[cx].nunique() > 1 else xr * 0.1
    ypad = yr / max(int(d[cy].nunique()) - 1, 1) * 0.7 if d[cy].nunique() > 1 else yr * 0.1
    # 축 표시범위(원 포함) — 모든 맵 공통. shot은 인접 센터 간격(pitch) 크기 사각형으로 그림(gap 제거)
    _xlo, _xhi, _ylo, _yhi = _wfmap_axis_limits(gx0, gx1, gy0, gy1, xpad, ypad, _circ)
    _pit_x = _wfmap_shot_pitch(d[cx]); _pit_y = _wfmap_shot_pitch(d[cy])

    res = []
    for gd, g in ordered:
        fig, ax = plt.subplots(figsize=(size_in, size_in))
        colors = np.where(g['_specout'].values, '#d32f2f', '#bdbdbd')
        _draw_wfmap_shots(ax, g[cx].astype(float).values, g[cy].astype(float).values,
                          _pit_x, _pit_y, colors=colors)
        _add_wafer_circle(ax, _circ, color='#000000', lw=1.0)
        # 흰 배경(회색 격자 제거) + 눈금/스파인 없음 → 경계 + shot map만 표시. aspect로 정원 유지
        ax.set_xticks([]); ax.set_yticks([]); ax.set_aspect(_wfmap_aspect(_circ), adjustable='box'); ax.set_facecolor('white')
        # 방향: 왼쪽=chip_x_adj 작은 쪽, 위쪽=chip_y_adj 작은 쪽(y축 반전). 경계까지 포함(안 잘리게)
        ax.set_xlim(_xlo, _xhi)
        ax.set_ylim(_yhi, _ylo)
        for sp in ax.spines.values():
            sp.set_visible(False)
        fig.subplots_adjust(left=0.02, right=0.98, top=0.98, bottom=0.02)
        _png = _wfmap_png_bytes(fig, dpi)   # 팔레트 양자화(PNG-8)로 용량 최소화
        plt.close(fig)
        _root = gd.get(croot, '') if croot else ''
        _wv = gd.get(cw, '')
        try:
            _wv = int(_wv)
        except (ValueError, TypeError):
            pass
        # (label, base64, is_target_lot) — target lot WF MAP은 HTML에서 라벨을 진한 파란색 강조
        res.append((f"{_root} #{_wv}", base64.b64encode(_png).decode('utf-8'), _is_tgt(gd)))
    return res


def _wfmap_cmap(direction):
    """REPORT DIRECTION별 WF MAP 컬러맵.
      - LOWER(하한 관리): 낮은값=빨강, 높은값=파랑  → 'coolwarm_r'
      - UPPER/BOTH:       낮은값=파랑, 높은값=빨강  → 'coolwarm'
    """
    return 'coolwarm_r' if str(direction).strip().upper() == 'LOWER' else 'coolwarm'


def _wfmap_norm(direction, spec_low, spec_high, values):
    """WF MAP diverging 컬러 normalization (PPT/HTML 공통, 동일 규칙).

    spec line 기준으로 색을 고정한다:
      - speclow → 컬러맵 0(파랑끝) · spechigh → 1(빨강끝) · median → 0.5(연회색 center)
      - cmap(_wfmap_cmap)이 LOWER에서 반전되므로 색의 빨강/파랑은 방향에 맞게 자동 적용
        (UPPER/BOTH: spechigh=빨강·speclow=파랑 / LOWER: 반대).
    한쪽 spec만 있으면(UPPER/LOWER) median 기준 대칭 확장하고, 둘 다 없으면 데이터 1~99% 선형 fallback.
    """
    import numpy as np
    import matplotlib.colors as mcolors

    def _f(x):
        try:
            x = float(x)
            return x if np.isfinite(x) else None
        except Exception:
            return None

    v = pd.to_numeric(pd.Series(values), errors='coerce').dropna().astype(float)
    med = float(v.median()) if len(v) else None
    lo, hi = _f(spec_low), _f(spec_high)
    if med is not None:
        if lo is None and hi is not None:
            lo = med - (hi - med)
        elif hi is None and lo is not None:
            hi = med + (med - lo)
    if med is None or lo is None or hi is None or not (hi > lo):
        # spec 결손 → 데이터 1~99% 선형
        if len(v):
            vmin, vmax = float(v.quantile(0.01)), float(v.quantile(0.99))
        else:
            vmin, vmax = 0.0, 1.0
        if not (vmax > vmin):
            vmax = vmin + 1e-9
        return mcolors.Normalize(vmin=vmin, vmax=vmax)
    # median(연회색 center)이 [lo,hi] 안에 오도록 clamp (TwoSlopeNorm 요구: lo<vc<hi)
    eps = (hi - lo) * 1e-6
    vc = min(max(med, lo + eps), hi - eps)
    return mcolors.TwoSlopeNorm(vmin=lo, vcenter=vc, vmax=hi)


def render_index_wfmap_b64(df, item, min_pts=50, lot_prefix=None,
                           cmap=None, size_in=1.15, dpi=None, direction='BOTH'):
    """특정 index의 '첫 측정(가장 이른 tkout_time) WF MAP'을 base64 PNG로 반환.

    - min_pts 이상 측정된 (lot, wafer, tkout) 그룹만 대상으로 하고, 그 중 tkout_time이
      가장 이른 그룹(=첫 측정)의 단일 wafer 칩맵을 그린다.
    - 조건을 만족하는 측정이 없으면 None (Score Board에 썸네일 미표기).
    """
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import base64, io

    if dpi is None:   # HTML용 WF MAP 해상도는 My_config.html_wfmap_dpi로 관리
        dpi = getattr(GLOBAL_CONFIG, 'html_wfmap_dpi', 120)

    if item not in df.columns:
        return None

    def _pick(*names):
        for n in names:
            if n in df.columns:
                return n
        return None
    cx = _pick('CHIP_X_ADJ', 'CHIP_X_POS', 'chip_x_pos')
    cy = _pick('CHIP_Y_ADJ', 'CHIP_Y_POS', 'chip_y_pos')
    cw = _pick('WAFER_ID', 'wafer_id')
    ct = _pick('TKOUT_TIME', 'tkout_time')
    cl = _pick('FAB_LOT_ID', 'fab_lot_id')
    crad = _pick('Chip_Radius', 'chip_radius')
    cm = _pick('MASK', 'mask')                        # vehicle 매칭용(150mm 원 shot 크기 계산)
    if not (cx and cy and cw and ct):
        return None

    cols = [c for c in [cx, cy, cw, ct, cl, crad, cm, item] if c]
    d = df[cols].dropna(subset=[item]).copy()
    if lot_prefix and cl:
        d = d[d[cl].astype(str).str.startswith(str(lot_prefix))]
    if len(d) == 0:
        return None
    d[ct] = pd.to_datetime(d[ct], errors='coerce')
    _circ = _wafer_circle_params(d, cx, cy, crad, mask_col=cm)   # 실제 150mm wafer 원(vehicle 매칭)

    grp_keys = [k for k in [cl, cw, ct] if k]
    sizes = d.groupby(grp_keys).size().reset_index(name='n')
    sizes = sizes[sizes['n'] >= min_pts]
    if len(sizes) == 0:
        return None
    first = sizes.sort_values(ct).iloc[0]          # 가장 이른 tkout 그룹
    sel = pd.Series(True, index=d.index)
    for k in grp_keys:
        sel &= (d[k] == first[k])
    g = d[sel]
    if len(g) == 0:
        return None

    vals = d[item].astype(float)
    vmin = float(vals.quantile(0.01)); vmax = float(vals.quantile(0.99))
    if not (vmax > vmin):
        vmax = vmin + 1e-9
    _gx0, _gx1 = float(d[cx].min()), float(d[cx].max())
    _gy0, _gy1 = float(d[cy].min()), float(d[cy].max())
    _xr = (_gx1 - _gx0) or 1.0; _yr = (_gy1 - _gy0) or 1.0
    _xpad = _xr / max(int(d[cx].nunique()) - 1, 1) * 0.7 if d[cx].nunique() > 1 else _xr * 0.1
    _ypad = _yr / max(int(d[cy].nunique()) - 1, 1) * 0.7 if d[cy].nunique() > 1 else _yr * 0.1
    # 축 표시범위(원 포함) — shot은 인접 센터 간격(pitch) 크기 사각형으로 그림(gap 제거)
    _xlo, _xhi, _ylo, _yhi = _wfmap_axis_limits(_gx0, _gx1, _gy0, _gy1, _xpad, _ypad, _circ)
    _pit_x = _wfmap_shot_pitch(d[cx]); _pit_y = _wfmap_shot_pitch(d[cy])
    import matplotlib.colors as _mcolors
    _norm_idx = _mcolors.Normalize(vmin=vmin, vmax=vmax)

    fig, ax = plt.subplots(figsize=(size_in, size_in))
    _draw_wfmap_shots(ax, g[cx].astype(float).values, g[cy].astype(float).values,
                      _pit_x, _pit_y, values=g[item].astype(float).values,
                      cmap=(cmap or _wfmap_cmap(direction)), norm=_norm_idx)
    _add_wafer_circle(ax, _circ, color='#000000', lw=1.0)
    ax.set_xticks([]); ax.set_yticks([]); ax.set_aspect(_wfmap_aspect(_circ), adjustable='box')
    ax.set_facecolor('white')
    # 방향: 왼쪽=chip_x_adj 작은 쪽, 위쪽=chip_y_adj 작은 쪽(y축 반전). 경계까지 포함(안 잘리게)
    ax.set_xlim(_xlo, _xhi)
    ax.set_ylim(_yhi, _ylo)
    for sp in ax.spines.values():
        sp.set_visible(False)
    _png = _wfmap_png_bytes(fig, dpi)   # 팔레트 양자화(PNG-8)로 용량 최소화
    plt.close(fig)
    return base64.b64encode(_png).decode('utf-8')


def calcaulate_description_image_info_dict(description_ppt_path, img_quality=20):
    """설명 PPT(python-pptx)를 열어 슬라이드별 좌상단 텍스트(Category)→슬라이드 매핑 반환.

    PNG 변환/중간파일(win32com Export) 의존을 제거했다. HOL_Auto_Report_Description.pptx
    파일만 있으면 되고, 이후 insert_plots가 매칭된 슬라이드를 python-pptx로 **직접 복사**해
    삽입한다. (PowerPoint 설치·COM 불필요, RUN/TEMP/desc_slide_*.png 미생성)

    반환: {category_text_lower: {'slide': slide, 'w': src_slide_width, 'h': src_slide_height}}.
    source slide/크기를 함께 넘겨, insert_plots가 현재 PPT 크기에 맞게 스케일 복사한다.
    """
    from pptx import Presentation

    desc_dict = {}
    if not description_ppt_path or not os.path.exists(description_ppt_path):
        print(f"[WARN] Description PPT를 찾을 수 없습니다: {description_ppt_path}")
        return desc_dict

    try:
        src_prs = Presentation(description_ppt_path)
        _src_w, _src_h = int(src_prs.slide_width), int(src_prs.slide_height)
        for slide in src_prs.slides:
            best_text, min_dist = None, float('inf')
            # 슬라이드 내 모든 도형 스캔 후 (0,0) 좌상단에 가장 가까운 텍스트를 Category로
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = (shape.text_frame.text or "").strip()
                if not text:
                    continue
                top = shape.top if shape.top is not None else 0
                left = shape.left if shape.left is not None else 0
                dist = int(top) ** 2 + int(left) ** 2
                if dist < min_dist:
                    min_dist = dist
                    best_text = text.split('\r')[0].split('\n')[0].strip()
            if best_text:
                key = best_text.replace('\x0b', '').strip().lower()
                if key:
                    desc_dict[key] = {'slide': slide, 'w': _src_w, 'h': _src_h}
    except Exception as e:
        print(f"[ERROR] Description PPT 파싱 중 에러 발생: {e}")

    return desc_dict


def _copy_slide_into(dest_prs, src_slide, src_w=None, src_h=None):
    """src_slide(다른 Presentation의 슬라이드)를 dest_prs에 새 슬라이드로 직접 복사.

    python-pptx로 도형 XML을 복사하고, 이미지 등 관계(rel)를 dest로 옮기며 rId를
    재매핑한다. 원본 슬라이드 크기(src_w/src_h)가 현재 PPT와 다르면 각 도형의
    위치/크기를 스케일하여 **현재 PPT 크기에 맞게 조정**한다(PNG 변환 없이 이미지처럼 꽉 차게).
    """
    import copy as _copy
    import io as _io
    _R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

    # ── 설명 이미지 재압축 설정 (My_config) ── 원본 고해상도 이미지를 다운스케일+JPEG 재압축해 용량 절감
    _recompress = getattr(GLOBAL_CONFIG, 'description_image_recompress', True)
    _max_px = int(getattr(GLOBAL_CONFIG, 'description_image_max_px', 2000) or 0)
    _jpeg_q = int(getattr(GLOBAL_CONFIG, 'description_image_jpeg_quality', 80))

    def _shrink_image_blob(blob):
        """이미지 blob을 max_px 이하로 축소 + JPEG 재압축한 bytes 반환(실패/무효 시 None)."""
        try:
            from PIL import Image
            im = Image.open(_io.BytesIO(blob))
            w, h = im.size
            if _max_px and max(w, h) > _max_px:
                _s = _max_px / float(max(w, h))
                im = im.resize((max(1, int(w * _s)), max(1, int(h * _s))), Image.LANCZOS)
            # 투명도 있으면 흰 배경 합성 후 RGB로(JPEG는 알파 미지원)
            if im.mode in ('RGBA', 'LA', 'P'):
                _rgba = im.convert('RGBA')
                _bg = Image.new('RGB', _rgba.size, (255, 255, 255))
                _bg.paste(_rgba, mask=_rgba.split()[-1])
                im = _bg
            else:
                im = im.convert('RGB')
            _out = _io.BytesIO()
            im.save(_out, format='JPEG', quality=_jpeg_q, optimize=True)
            return _out.getvalue()
        except Exception:
            return None

    dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[6])
    # 빈 레이아웃의 기본 placeholder 제거(빈 상자 방지)
    for ph in list(dest_slide.shapes):
        try:
            if ph.is_placeholder:
                ph._element.getparent().remove(ph._element)
        except Exception:
            pass

    # 관계(이미지 등) 복사 → old rId → new rId 매핑 (레이아웃/노트 관계는 제외)
    rid_map = {}
    try:
        for rId, rel in src_slide.part.rels.items():
            if 'slideLayout' in rel.reltype or 'notesSlide' in rel.reltype:
                continue
            if rel.is_external:
                new_rId = dest_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            elif _recompress and 'image' in rel.reltype:
                # 이미지: 다운스케일+JPEG 재압축한 새 이미지 파트로 추가(용량↓). 실패 시 원본 그대로.
                new_rId = None
                try:
                    _orig = rel.target_part.blob
                    _small = _shrink_image_blob(_orig)
                    if _small is not None and len(_small) < len(_orig):
                        _img_part, new_rId = dest_slide.part.get_or_add_image_part(_io.BytesIO(_small))
                except Exception:
                    new_rId = None
                if new_rId is None:
                    new_rId = dest_slide.part.relate_to(rel.target_part, rel.reltype)
            else:
                new_rId = dest_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rId
    except Exception as _re:
        print(f"[WARN] description 슬라이드 관계 복사 일부 실패: {_re}")

    # 원본 도형의 '해결된' 위치/크기를 먼저 캡처한다.
    #   placeholder 도형은 자기 XML에 위치(<a:off>/<a:ext>)가 없고 슬라이드 레이아웃에서
    #   상속받는다. blank 레이아웃(6번)으로 복사하면 상속처가 없어 가운데로 몰리므로,
    #   python-pptx가 상속을 해석해 돌려주는 절대 좌표를 복사본에 명시적으로 새겨 넣는다.
    src_geoms = []
    for shape in src_slide.shapes:
        try:
            src_geoms.append((shape.left, shape.top, shape.width, shape.height))
        except Exception:
            src_geoms.append((None, None, None, None))

    # 도형 XML 복사
    for shape in src_slide.shapes:
        dest_slide.shapes._spTree.append(_copy.deepcopy(shape._element))

    # 복사된 XML의 old rId 참조를 new rId로 치환 (r:embed, r:link, r:id 등)
    if rid_map:
        for el in dest_slide.shapes._spTree.iter():
            for attr in list(el.attrib):
                if attr.startswith(_R) and el.attrib[attr] in rid_map:
                    el.attrib[attr] = rid_map[el.attrib[attr]]

    # 캡처한 절대 위치/크기를 복사본에 명시적으로 기입(placeholder 상속 → 절대좌표 고정)
    # + 원본↔현재 슬라이드 크기 차이만큼 스케일하여 현재 PPT에 맞춘다.
    try:
        _sw = int(src_w) if src_w else int(dest_prs.slide_width)
        _sh = int(src_h) if src_h else int(dest_prs.slide_height)
        sx = dest_prs.slide_width / _sw if _sw else 1.0
        sy = dest_prs.slide_height / _sh if _sh else 1.0
        for (l, t, w, h), sh in zip(src_geoms, list(dest_slide.shapes)):
            try:
                if l is not None:
                    sh.left = int(round(l * sx))
                if t is not None:
                    sh.top = int(round(t * sy))
                if w is not None:
                    sh.width = int(round(w * sx))
                if h is not None:
                    sh.height = int(round(h * sy))
            except Exception:
                pass
    except Exception as _se:
        print(f"[WARN] description 슬라이드 크기 스케일 실패: {_se}")
    return dest_slide


def _fmt_stat_value(v):
    """PPT 통계표 셀 숫자 표기 규칙.

    - |v| ≤ 0.001         → 유효숫자 2개 scientific (예: 5.0e-4)
    - 0.001 < |v| < 1     → 소수점 셋째자리(0.123)
    - 1 ≤ |v| < 1000      → 소수점 첫째자리(12.3)
    - |v| ≥ 1000          → 유효숫자 2개 scientific(1.1e5)
    - NaN/None            → '-'
    """
    try:
        if v is None or pd.isna(v):
            return "-"
    except Exception:
        return "-"
    v = float(v)
    av = abs(v)
    if av == 0:
        return "0"

    def _sci(x):
        _s = f"{x:.1e}"                    # 예: '1.1e+05' / '5.0e-04'
        _mant, _exp = _s.split('e')
        return f"{_mant}e{int(_exp)}"      # '1.1e5' / '5.0e-4'

    if av <= 1e-3 or av >= 1000:
        return _sci(v)
    if av < 1:
        return f"{v:.3f}"
    return f"{v:.1f}"


# ===================================================================
#  병렬 렌더링 인프라 (Parallel chart rendering infrastructure)
# ===================================================================
# 발행 속도 최대화를 위해 matplotlib 차트 렌더링(CPU 바운드)을 워커 프로세스로
# 분리한다. 워커 수는 실행 환경의 CPU 코어 수와 '가용' 메모리를 보고 자동 결정:
#   workers = min(코어수, parallel_max_workers, (가용GB - reserve) / per_worker)
#   예) 4코어/50GB → 4워커, 2코어/10GB → 2워커, 가용 메모리 부족 → 1(직렬 폴백)
# 관련 설정: My_config.parallel_workers(0=자동)/parallel_max_workers/
#            parallel_mem_per_worker_gb/parallel_reserve_gb
#
# ⚠️ Windows의 프로세스 spawn은 워커가 __main__(Main.py)을 다시 import 하므로
#    Main.py 실행 본문은 반드시 `if __name__ == "__main__":` 가드 안에 있어야 한다.

# Wafer ID별 고정 컬러 팔레트 (25색) — 모든 차트/워커가 공통 사용.
# wafer #1 → [0], #2 → [1] ... 동일 wafer 번호는 항상 같은 색.
_WAFER_PALETTE = [
    '#e6194b', '#3cb44b', '#4363d8', '#f58231', '#911eb4',
    '#42d4f4', '#f032e6', '#bfef45', '#fabed4', '#469990',
    '#dcbeff', '#9a6324', '#bcbd22', '#800000', '#aaffc3',
    '#808000', '#ffd8b1', '#000075', '#a9a9a9', '#ffe119',
    '#1a1aff', '#ff4dd2', '#00cca3', '#b35900', '#5c5c8a',
]


def _get_available_mem_gb():
    """현재 '가용' 물리 메모리(GB)를 반환. 측정 불가 시 None.

    psutil → (Windows) GlobalMemoryStatusEx → (Linux) /proc/meminfo 순으로 시도.
    """
    try:
        import psutil
        return psutil.virtual_memory().available / (1024 ** 3)
    except Exception:
        pass
    if os.name == 'nt':
        try:
            import ctypes

            class _MEMORYSTATUSEX(ctypes.Structure):
                _fields_ = [
                    ("dwLength", ctypes.c_ulong), ("dwMemoryLoad", ctypes.c_ulong),
                    ("ullTotalPhys", ctypes.c_ulonglong), ("ullAvailPhys", ctypes.c_ulonglong),
                    ("ullTotalPageFile", ctypes.c_ulonglong), ("ullAvailPageFile", ctypes.c_ulonglong),
                    ("ullTotalVirtual", ctypes.c_ulonglong), ("ullAvailVirtual", ctypes.c_ulonglong),
                    ("ullAvailExtendedVirtual", ctypes.c_ulonglong)]

            st = _MEMORYSTATUSEX()
            st.dwLength = ctypes.sizeof(_MEMORYSTATUSEX)
            if ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(st)):
                return st.ullAvailPhys / (1024 ** 3)
        except Exception:
            pass
    try:
        with open('/proc/meminfo') as f:
            for line in f:
                if line.startswith('MemAvailable:'):
                    return int(line.split()[1]) / (1024 ** 2)
    except Exception:
        pass
    return None


def get_parallel_workers(config=None):
    """실행 환경(CPU 코어/가용 메모리)을 보고 렌더링 워커 프로세스 수를 결정.

    - My_config.parallel_workers > 0 이면 그 값을 그대로 사용(강제 지정).
    - 자동: min(코어수, parallel_max_workers,
              (가용GB - parallel_reserve_gb) // parallel_mem_per_worker_gb)
    - 가용 메모리 측정 불가 시 보수적으로 min(코어수, 2).
    - 결과 1이면 프로세스 풀 없이 직렬 렌더링(저사양 환경 안전 동작).
    """
    cfg = config if config is not None else GLOBAL_CONFIG
    try:
        forced = int(getattr(cfg, 'parallel_workers', 0) or 0)
    except (TypeError, ValueError):
        forced = 0
    if forced > 0:
        return forced
    cores = os.cpu_count() or 2
    cap = int(getattr(cfg, 'parallel_max_workers', 8) or 8)
    per_gb = float(getattr(cfg, 'parallel_mem_per_worker_gb', 1.2) or 1.2)
    reserve = float(getattr(cfg, 'parallel_reserve_gb', 3.0) or 3.0)
    avail = _get_available_mem_gb()
    if avail is None:
        mem_cap = min(cores, 2)
    else:
        mem_cap = int(max(0.0, avail - reserve) // per_gb)
    return max(1, min(cores, cap, mem_cap))


_CHART_POOL = None
_CHART_POOL_N = 0


def _get_chart_pool(n):
    """렌더링용 프로세스 풀(모듈 전역 1개)을 생성/재사용. n<=1이면 None(직렬).

    풀은 한 번 만들면 리포트/랏이 바뀌어도 재사용한다(워커 spawn 비용 1회만).
    """
    global _CHART_POOL, _CHART_POOL_N
    if n is None or n <= 1:
        return None
    if _CHART_POOL is not None and _CHART_POOL_N >= n:
        return _CHART_POOL
    shutdown_chart_pool()
    try:
        from concurrent.futures import ProcessPoolExecutor
        _CHART_POOL = ProcessPoolExecutor(max_workers=n)
        _CHART_POOL_N = n
    except Exception as e:
        print(f"[WARN] 렌더링 워커 풀 생성 실패(직렬로 진행): {e}")
        _CHART_POOL = None
        _CHART_POOL_N = 0
    return _CHART_POOL


def shutdown_chart_pool():
    """렌더링 워커 풀 종료(atexit에도 등록됨 — Main 종료 시 명시 호출 권장)."""
    global _CHART_POOL, _CHART_POOL_N
    if _CHART_POOL is not None:
        try:
            _CHART_POOL.shutdown(wait=False, cancel_futures=True)
        except TypeError:      # Python<3.9: cancel_futures 미지원
            _CHART_POOL.shutdown(wait=False)
        except Exception:
            pass
        _CHART_POOL = None
        _CHART_POOL_N = 0


import atexit as _atexit
_atexit.register(shutdown_chart_pool)


def _select_target_lot_frame(frame, target_lot_id, target_root_lot_id, target_DC_step_id):
    """대상 lot 선택: match_key(root+step) → fab_lot_id 정확일치 → root prefix 순.

    (insert_plots 내부 closure였던 _select_target_lot을 워커에서도 쓰도록 모듈화.
     리포트 단위 = root_lot_id + step. 같은 root의 형제 lot_id를 모두 포함한다.)
    """
    if 'match_key' in frame.columns and target_root_lot_id and target_DC_step_id:
        mk = f"{target_root_lot_id}_{target_DC_step_id}"
        sel = frame[frame['match_key'].astype(str) == mk]
        if len(sel) > 0:
            return sel
    if 'fab_lot_id' not in frame.columns:
        return frame
    flid = frame['fab_lot_id'].astype(str)
    sel = frame[flid == str(target_lot_id)]
    if len(sel) == 0 and target_root_lot_id:
        sel = frame[flid.str.startswith(str(target_root_lot_id))]
    return sel


def _render_wafer_legend_bytes(cfg):
    """Wafer 색 범례(모든 슬라이드 공통·정적)를 1회 렌더링해 jpg bytes로 반환."""
    import io as _io
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plt.rcParams['font.family'] = cfg['FONT']
    w_colors = {str(i): _WAFER_PALETTE[(i - 1) % len(_WAFER_PALETTE)] for i in range(1, 26)}
    fig_leg, ax_leg = plt.subplots(figsize=(8.0, 0.2))
    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax_leg.axis('off')
    ax_leg.text(0, 0, "Wafer Color:", va='center', ha='left', fontsize=8, fontweight='bold', color=cfg['C_NAVY'])
    for i in range(1, 26):
        ax_leg.scatter([i * 0.85 + 2.5], [0], color=w_colors.get(str(i)), s=25)
        ax_leg.text(i * 0.85 + 2.5 + 0.15, 0, f"#{i}", va='center', ha='left', fontsize=6, color=cfg['C_NEUTRAL'])
    ax_leg.set_xlim(0, 25)
    ax_leg.set_ylim(-1, 1)
    buf = _io.BytesIO()
    fig_leg.savefig(buf, format='jpg', dpi=cfg['dpi'], facecolor='white', pil_kwargs={'quality': cfg['jpg_q']})
    plt.close(fig_leg)
    return buf.getvalue()


def _render_item_charts(task):
    """[워커/직렬 공용] 한 index의 차트 5종(Box/WF MAP/Trend/Radius/CDF)을 렌더링.

    insert_plots의 per-item 렌더링을 프로세스 워커로 분리한 것. 차트 모양/색/
    판정 로직은 종전과 동일하며, PPT 조립(슬라이드/표)은 메인 프로세스가 한다.
    GLOBAL_CONFIG를 참조하지 않고 task['cfg']만 사용한다(스폰된 워커에는
    vehicle yaml 설정이 로드되어 있지 않기 때문).

    반환 dict:
      status('ok'|'skip'|'error'), reason, warnings[],
      imgs{box,map,trend,rad,cum: jpg bytes}, map_ratio(WF MAP h/w 비율),
      stat_rows(통계 행 라벨 4개), stat_cells{1..25: [셀 문자열 4개] | None},
      metrics(항목 지표 dict), summary_row(Index Aggregation Table 행)
    """
    import io as _io
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates

    cfg = task['cfg']
    item_name = task['item_name']
    spec_name = task['spec_name']
    target_lot_id = task['target_lot_id']
    target_root_lot_id = task['target_root_lot_id']
    target_DC_step_id = task['target_DC_step_id']

    sp = task['spec']
    spec_low = sp['spec_low']
    spec_high = sp['spec_high']
    wfmap_spec_low = sp['wfmap_spec_low']
    wfmap_spec_high = sp['wfmap_spec_high']
    direction = sp['direction']
    log_scale = sp['log_scale']
    y_label = sp['y_label']

    dpi = cfg['dpi']; jpg_q = cfg['jpg_q']; map_q = cfg['map_q']
    C_NAVY = cfg['C_NAVY']; C_ACCENT = cfg['C_ACCENT']; C_NEUTRAL = cfg['C_NEUTRAL']
    C_GRID = cfg['C_GRID']; C_SPINE = cfg['C_SPINE']; C_VEHICLE = cfg['C_VEHICLE']
    C_BAND = cfg['C_BAND']; WV_PALETTE = cfg['WV_PALETTE']; FONT = cfg['FONT']
    main_vehicle = cfg['main_vehicle']

    w_colors = {str(i): _WAFER_PALETTE[(i - 1) % len(_WAFER_PALETTE)] for i in range(1, 26)}
    fixed_wafers = [str(i) for i in range(1, 26)]

    out = {'item_name': item_name, 'status': 'ok', 'reason': '', 'warnings': [],
           'imgs': {}, 'map_ratio': None, 'stat_rows': [], 'stat_cells': {},
           'metrics': None, 'summary_row': None}

    def _remove_spines(ax):
        for spine in ['top', 'right']:
            ax.spines[spine].set_visible(False)
        ax.spines['left'].set_color(C_SPINE)
        ax.spines['bottom'].set_color(C_SPINE)
        ax.spines['left'].set_linewidth(0.6)
        ax.spines['bottom'].set_linewidth(0.6)

    def _label_axes(ax, xlabel=None, ylabel=None, ylabel_size=7, xlabel_size=7):
        if xlabel is not None:
            ax.set_xlabel(xlabel, fontsize=xlabel_size, color=C_NAVY, fontname=FONT)
        if ylabel is not None:
            ax.set_ylabel(ylabel, fontsize=ylabel_size, color=C_NAVY, fontname=FONT)

    try:
        cols = task['cols']
        col_x = cols['col_x']; col_y = cols['col_y']; col_rad = cols['col_rad']
        col_sub = cols['col_sub']; col_time = cols['col_time']
        col_lot = cols['col_lot']; col_mask = cols['col_mask']
        w_col_src = cols['w_col_src']

        item_df = task['df']

        # 메타 컬럼명을 canonical 소문자로 통일 (merged_df가 대문자여도 안전)
        rename_map = {}
        if col_time and col_time != 'tkout_time': rename_map[col_time] = 'tkout_time'
        if col_lot and col_lot != 'fab_lot_id': rename_map[col_lot] = 'fab_lot_id'
        if col_mask and col_mask != 'mask': rename_map[col_mask] = 'mask'
        if w_col_src != 'wafer_id': rename_map[w_col_src] = 'wafer_id'
        if rename_map:
            item_df = item_df.rename(columns=rename_map)
        w_col = 'wafer_id'

        # 누락 컬럼 대체 (Fallback for missing columns)
        if not col_x:
            item_df['chip_x_pos'] = np.random.rand(len(item_df))
            item_df['chip_y_pos'] = np.random.rand(len(item_df))
            col_x, col_y = 'chip_x_pos', 'chip_y_pos'
        if not col_rad:
            item_df['chip_radius'] = np.sqrt(item_df[col_x]**2 + item_df[col_y]**2)
            col_rad = 'chip_radius'
        if not col_sub:
            item_df['subitem_id'] = 'ALL'
            col_sub = 'subitem_id'
        if 'tkout_time' not in item_df.columns:
            item_df['tkout_time'] = pd.to_datetime('today')

        item_df[w_col] = item_df[w_col].astype(str).str.replace('#', '')
        item_df['tkout_time'] = pd.to_datetime(item_df['tkout_time'])
        item_df = item_df.sort_values(by='tkout_time')

        # ---- 데이터 스코프 분리 ----
        # item_df_full: 전 vehicle·전 lot (Trend 비교 전용)
        # item_df(=target_df): 리포팅 대상 lot만 (Box/WF MAP/통계표/Radius/CDF/통계)
        item_df_full = item_df
        target_df = _select_target_lot_frame(item_df, target_lot_id, target_root_lot_id, target_DC_step_id)
        item_df = target_df
        if len(item_df) == 0:
            out['status'] = 'skip'
            out['reason'] = f"대상 lot '{target_lot_id}' 데이터 없음"
            return out

        measured_wafers = sorted(item_df[w_col].unique(), key=lambda x: int(x) if str(x).isdigit() else x)
        grouped = item_df.groupby(w_col)[item_name]

        # ---- 다중 lot_id 구분 (같은 root_lot_id에 여러 lot_id가 묶일 때) ----
        lot_col = 'fab_lot_id' if 'fab_lot_id' in item_df.columns else None
        target_lots = (sorted(item_df[lot_col].astype(str).unique()) if lot_col else [])
        multi_lot = len(target_lots) > 1
        LOT_MARKERS = ['o', '^', 's', 'D', 'v', 'P', 'X', '*']
        LOT_LINE_COLORS = ['#d62728', '#1f77b4', '#2ca02c', '#9467bd', '#ff7f0e', '#8c564b', '#e377c2', '#17becf']
        lot_marker = {lot: LOT_MARKERS[i % len(LOT_MARKERS)] for i, lot in enumerate(target_lots)}
        lot_color = {lot: LOT_LINE_COLORS[i % len(LOT_LINE_COLORS)] for i, lot in enumerate(target_lots)}

        # ---- Index Aggregation Table용 (lot, wafer)별 집계값 ----
        # REPORT DIRECTION → BOTH=Median / UPPER=P90 / LOWER=P10. lot_id별 분리해 wafer별 1값.
        try:
            if direction == 'UPPER':
                _q, _sstat = 0.90, 'P90'
            elif direction == 'LOWER':
                _q, _sstat = 0.10, 'P10'
            else:
                _q, _sstat = None, 'Median'

            def _aggv(s):
                s = pd.to_numeric(s, errors='coerce').dropna()
                if len(s) == 0:
                    return None
                return float(s.median()) if _q is None else float(s.quantile(_q))

            _lc = 'fab_lot_id' if 'fab_lot_id' in item_df.columns else None
            _wv = {}   # {(lot, wafer_int): value}
            if _lc:
                for (_lot, _waf), _g in item_df.groupby([_lc, w_col]):
                    _v = _aggv(_g[item_name])
                    if _v is None:
                        continue
                    try:
                        _waf = int(float(_waf))
                    except (ValueError, TypeError):
                        pass
                    _wv[(str(_lot), _waf)] = _v
            else:
                _v = _aggv(item_df[item_name])
                if _v is not None:
                    _wv[(str(target_lot_id), 0)] = _v
            if _wv:
                out['summary_row'] = {'index': item_name, 'stat': _sstat, 'wafer_vals': _wv}
        except Exception as _se:
            out['warnings'].append(f"summary 집계 실패: {_se}")

        # ---- 통계 테이블 값 (REPORT DIRECTION별 4행 × wafer #1~25) ----
        #   UPPER: Max P90 Med P10 / LOWER: P90 Med P10 Min / BOTH: P90 Med P10 Std
        if direction == 'UPPER':
            stat_defs = [("Max", lambda s: s.max()), ("P90", lambda s: s.quantile(0.90)),
                         ("Med", lambda s: s.median()), ("P10", lambda s: s.quantile(0.10))]
        elif direction == 'LOWER':
            stat_defs = [("P90", lambda s: s.quantile(0.90)), ("Med", lambda s: s.median()),
                         ("P10", lambda s: s.quantile(0.10)), ("Min", lambda s: s.min())]
        else:
            stat_defs = [("P90", lambda s: s.quantile(0.90)), ("Med", lambda s: s.median()),
                         ("P10", lambda s: s.quantile(0.10)),
                         ("Std", lambda s: s.std() if len(s) > 1 else float('nan'))]
        out['stat_rows'] = [lbl for lbl, _ in stat_defs]
        for w_idx in range(1, 26):
            w_str = str(w_idx)
            if w_str in grouped.groups:
                grp_data = pd.to_numeric(grouped.get_group(w_str), errors='coerce').dropna()
                vals = []
                for _lbl, _fn in stat_defs:
                    try:
                        vals.append(_fmt_stat_value(_fn(grp_data)))
                    except Exception:
                        vals.append("-")
                out['stat_cells'][w_idx] = vals
            else:
                out['stat_cells'][w_idx] = None

        # 통계표 (lot_id 포함) — HTML score board처럼 (lot, wafer)별 값. 측정된 것만.
        out['stat_cells_lw'] = {}   # {(lot_str, wafer_int): [v1..v4]}
        _lc_s = 'fab_lot_id' if 'fab_lot_id' in item_df.columns else None
        if _lc_s:
            for (_lot, _waf), _g in item_df.groupby([_lc_s, w_col]):
                try:
                    _wi = int(str(_waf).replace('#', ''))
                except (ValueError, TypeError):
                    continue
                _d = pd.to_numeric(_g[item_name], errors='coerce').dropna()
                _vals = []
                for _lbl, _fn in stat_defs:
                    try:
                        _vals.append(_fmt_stat_value(_fn(_d)))
                    except Exception:
                        _vals.append("-")
                out['stat_cells_lw'][(str(_lot), _wi)] = _vals

        # ---- 차트 임시 버퍼 (디스크 파일 없이 메모리에서 처리) ----
        tmp_box = _io.BytesIO()
        tmp_map = _io.BytesIO()
        tmp_trend = _io.BytesIO()
        tmp_rad = _io.BytesIO()
        tmp_cum = _io.BytesIO()

        plt.rcParams['axes.linewidth'] = 0.6
        plt.rcParams['font.size'] = 7.5
        plt.rcParams['font.family'] = FONT
        plt.rcParams['axes.facecolor'] = '#ffffff'
        plt.rcParams['figure.facecolor'] = '#ffffff'
        plt.rcParams['grid.color'] = C_GRID
        plt.rcParams['grid.linestyle'] = '-'
        plt.rcParams['grid.linewidth'] = 0.5

        # ---- 3. BOX Plot (측정된 웨이퍼만 플롯하되, X축은 1~25 고정) ----
        fig_box, ax_box = plt.subplots(figsize=(7.8, 1.9))
        box_data = []
        box_positions = []
        for w in measured_wafers:
            try:
                w_int = int(w)
                box_positions.append(w_int)
                box_data.append(grouped.get_group(w).values)
            except ValueError: pass

        if len(box_data) > 0:
            bplot = ax_box.boxplot(
                box_data, positions=box_positions, patch_artist=True,
                showfliers=True,
                flierprops={'marker': 'o', 'markerfacecolor': C_ACCENT, 'markeredgecolor': 'none', 'markersize': 3.5},
                medianprops={'color': C_NEUTRAL, 'linewidth': 1.2},
                whiskerprops={'color': '#94a3b8', 'linewidth': 0.8},
                capprops={'color': '#94a3b8', 'linewidth': 0.8}
            )
            for i, patch in enumerate(bplot['boxes']):
                w_str = str(box_positions[i])
                patch.set_facecolor(w_colors.get(w_str, '#ffffff'))
                patch.set_alpha(0.85)
                patch.set_edgecolor(C_NEUTRAL)
                patch.set_linewidth(0.8)

        # Spec line은 그리되 범례(Spec Limit)는 표시하지 않음
        if spec_low is not None:
            ax_box.axhline(y=float(spec_low), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        if spec_high is not None:
            ax_box.axhline(y=float(spec_high), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        # 대상 lot_id(리포트 subject)에 속한 wafer 번호 → x축 라벨에 '*' 표기 + 우하단 범례
        _tgt_wafnums = set()
        if 'fab_lot_id' in item_df.columns:
            _tl_box = item_df[item_df['fab_lot_id'].astype(str) == str(target_lot_id)]
            for _w in _tl_box[w_col].unique():
                try:
                    _tgt_wafnums.add(int(_w))
                except (ValueError, TypeError):
                    pass
        ax_box.set_xticks(range(1, 26))
        ax_box.set_xticklabels([f"#{i}" + ("*" if i in _tgt_wafnums else "") for i in range(1, 26)])
        ax_box.set_xlim(0.5, 25.5)
        ax_box.tick_params(axis='x', rotation=45, labelsize=7)
        if _tgt_wafnums:
            # box plot '바깥 아래'(축 아래, x라벨 밑)에 '*: lot_id' 범례를 우측 정렬로 배치
            #  (y<0 = 축 영역 밖 아래, 오른쪽 끝을 box plot 우측에 맞춤)
            ax_box.text(1.0, -0.5, f"*: {target_lot_id}", transform=ax_box.transAxes,
                        ha='right', va='top', fontsize=6, color=C_NEUTRAL, fontstyle='italic')
        _label_axes(ax_box, xlabel="Wafer #", ylabel=y_label)
        _remove_spines(ax_box)
        ax_box.set_axisbelow(True)
        if log_scale:
            # 로그 스케일: 데이터 범위가 한 decade 미만이어도 10의 거듭제곱마다 y축 선이 나오도록
            # major(10^n) + minor(2~9×10^n) locator를 명시하고 둘 다 grid 표시.
            from matplotlib.ticker import LogLocator, NullFormatter
            ax_box.set_yscale('log')
            ax_box.yaxis.set_major_locator(LogLocator(base=10.0, numticks=15))
            ax_box.yaxis.set_minor_locator(
                LogLocator(base=10.0, subs=(2, 3, 4, 5, 6, 7, 8, 9), numticks=15))
            ax_box.yaxis.set_minor_formatter(NullFormatter())
            ax_box.grid(True, which='major', axis='both', color=C_GRID, linestyle='-', linewidth=0.6)
            ax_box.grid(True, which='minor', axis='y', color=C_GRID, linestyle='-', linewidth=0.35, alpha=0.6)
        else:
            ax_box.minorticks_off()  # minor tick(세부선) 제거 — major만 표시
            ax_box.grid(True, which='major', axis='both', color=C_GRID, linestyle='-', linewidth=0.5)
        # 용량 다이어트를 위한 JPG 포맷 저장 및 quality 옵션 적용
        fig_box.savefig(tmp_box, format='jpg', dpi=dpi, bbox_inches="tight", facecolor='white', pil_kwargs={'quality': jpg_q})
        plt.close(fig_box)
        out['imgs']['box'] = tmp_box.getvalue()

        # ---- 4. WF MAP (행=PGM(pt), 열=wafer #1~25 고정) ----
        # Wafer 좌표: flat-zone 회전이 반영된 보정 좌표(CHIP_X_ADJ/CHIP_Y_ADJ)가
        # 있으면 우선 사용해 실제 웨이퍼 배치대로 그린다. 없으면 raw chip_x/y로 fallback.
        map_x = 'CHIP_X_ADJ' if 'CHIP_X_ADJ' in item_df.columns else col_x
        map_y = 'CHIP_Y_ADJ' if 'CHIP_Y_ADJ' in item_df.columns else col_y

        # WF MAP 행 = PGM(pt)별 한 줄. PGM(pt) 컬럼이 있으면 그것으로 그룹핑(없으면 subitem).
        _wfmap_grp_col = 'PGM(pt)' if 'PGM(pt)' in item_df.columns else col_sub
        sub_groups = list(item_df.groupby(_wfmap_grp_col))
        n_pgm = len(sub_groups) if len(sub_groups) > 0 else 1

        global_x_min = item_df[map_x].min()
        global_x_max = item_df[map_x].max()
        global_y_min = item_df[map_y].min()
        global_y_max = item_df[map_y].max()

        x_range = global_x_max - global_x_min if global_x_max > global_x_min else 1
        y_range = global_y_max - global_y_min if global_y_max > global_y_min else 1
        # 가장자리 칩(사각 마커)이 잘리지 않도록 반칩 이상 여백 확보 (칩 간격 기준)
        _nxu = max(int(item_df[map_x].nunique()), 1)
        _nyu = max(int(item_df[map_y].nunique()), 1)
        x_pad = (x_range / max(_nxu - 1, 1)) * 0.7 if _nxu > 1 else x_range * 0.1
        y_pad = (y_range / max(_nyu - 1, 1)) * 0.7 if _nyu > 1 else y_range * 0.1
        global_x_min -= x_pad
        global_x_max += x_pad
        global_y_min -= y_pad
        global_y_max += y_pad

        # ---- WF MAP 공유 컬러 스케일 (spec line 기준 diverging, HTML과 동일 규칙) ----
        wfmap_norm = _wfmap_norm(direction, wfmap_spec_low, wfmap_spec_high, item_df[item_name])
        # 배경 wafer 원(150mm) — radius(mm) 기반 중심·x/y 스케일 fit (모든 셀 공통)
        _circ_ppt = _wafer_circle_params(item_df, map_x, map_y, col_rad,
                                         mask_col=col_mask, main_vehicle=main_vehicle)   # 실제 150mm 원(vehicle 매칭)
        # 셀 공통 축범위: 데이터(패딩) + wafer 경계(타원)를 모두 포함(경계 안 잘리게)
        _px_lo, _px_hi = global_x_min, global_x_max
        _py_lo, _py_hi = global_y_min, global_y_max
        if _circ_ppt:
            _ccx, _ccy, _csx, _csy = _circ_ppt[0], _circ_ppt[1], _circ_ppt[2], _circ_ppt[3]
            _mx = abs(_csx) * 0.08; _my = abs(_csy) * 0.08
            _px_lo = min(_px_lo, _ccx - _csx - _mx); _px_hi = max(_px_hi, _ccx + _csx + _mx)
            _py_lo = min(_py_lo, _ccy - _csy - _my); _py_hi = max(_py_hi, _ccy + _csy + _my)

        # 칩 격자 차원(예: 13x13)
        nx = max(int(item_df[map_x].nunique()), 1)
        ny = max(int(item_df[map_y].nunique()), 1)
        gdim = max(nx, ny)

        # ---- WF MAP 배치: 행=PGM(pt), 열=wafer #1~25 고정 ----
        FIXED_N_WAF = 25
        grid_rows, grid_cols = n_pgm, FIXED_N_WAF
        _multi_pgm = n_pgm > 1
        cell_map = {(i, c): (sub_grp, str(c + 1), sub_name)
                    for i, (sub_name, sub_grp) in enumerate(sub_groups)
                    for c in range(FIXED_N_WAF)}

        render_cell = 0.55                             # 셀(=wafer 1칸) 크기(인치) — 25칸 고정
        # shot: 인접 센터 간격(pitch) 크기 사각형으로 그려 gap 없이(측정 pt 수와 무관하게 동일 크기)
        _shot_px = _wfmap_shot_pitch(item_df[map_x]); _shot_py = _wfmap_shot_pitch(item_df[map_y])
        fig_disp_w = grid_cols * render_cell           # 서브플롯 영역 폭(항상 25*cell)
        fig_h = grid_rows * render_cell + 0.15         # 하단 wafer 번호 라벨 여백

        # 컬러바는 별도 이미지로 분리(슬라이드에서 항상 같은 크기로 우측에 길게 배치) → 맵은 grid만
        fig_map, axes_map = plt.subplots(grid_rows, grid_cols,
                                         figsize=(fig_disp_w, fig_h), squeeze=False,
                                         gridspec_kw={'wspace': 0.06, 'hspace': 0.18})
        sc = None
        for r in range(grid_rows):
            for c in range(grid_cols):
                ax = axes_map[r, c]
                sub_grp, w, sub_name = cell_map[(r, c)]
                w_grp = sub_grp[sub_grp[w_col] == w]
                if not w_grp.empty:   # 측정된 wafer만 die 산점, 없으면 빈 칸
                    sc = _draw_wfmap_shots(ax, w_grp[map_x].astype(float).values,
                                           w_grp[map_y].astype(float).values, _shot_px, _shot_py,
                                           values=w_grp[item_name].astype(float).values,
                                           cmap=_wfmap_cmap(direction), norm=wfmap_norm)
                    _add_wafer_circle(ax, _circ_ppt, color='#000000', lw=1.0)   # 배경 wafer 원 (HTML과 동일)
                ax.set_facecolor('white')
                # 마지막(맨 아래) PGM 행에만 wafer 번호 표기
                if r == grid_rows - 1:
                    ax.set_xlabel(f"#{c + 1}", fontsize=5, labelpad=1, color=C_NEUTRAL)
                # 다중 PGM이면 첫 열 좌측에 PGM(pt) 라벨(90도 회전, 작은 글씨)
                # 표기 형식: step_seq(pt수)  예) DC_01(3)  — '_1.0' 같은 접미사·'pt' 문자는 제거
                if c == 0 and _multi_pgm:
                    _pl = str(sub_name)
                    if 'PGM(pt)' in sub_grp.columns:
                        _u = sub_grp['PGM(pt)'].dropna()
                        if len(_u): _pl = str(_u.iloc[0])
                    _pl = re.sub(r'_[0-9.]+$', '', _pl)   # Duplicate_Count 접미사(_1.0 등) 제거
                    _pl = re.sub(r'(\d+)\s*pt\)', r'\1)', _pl)  # "(137pt)" → "(137)"
                    ax.set_ylabel(_pl, fontsize=4, rotation=90, labelpad=4, color=C_NEUTRAL)
                ax.set_xticks([]); ax.set_yticks([])
                ax.set_aspect(_wfmap_aspect(_circ_ppt), adjustable='box')   # aspect=ky/kx → 정원 유지
                for spine in ax.spines.values():
                    spine.set_visible(False)
                # 방향: 왼쪽=chip_x_adj 작은 쪽, 위쪽=chip_y_adj 작은 쪽(y축 반전). HTML과 동일. 경계 포함
                ax.set_xlim(_px_lo, _px_hi)
                ax.set_ylim(_py_hi, _py_lo)

        fig_map.subplots_adjust(left=0.02, right=0.99, top=0.97, bottom=0.06)
        # 칩 격자가 선명하도록 해상도 상향
        fig_map.savefig(tmp_map, format='jpg', dpi=max(int(dpi), 200), bbox_inches="tight",
                        facecolor='white', pil_kwargs={'quality': map_q})
        plt.close(fig_map)
        out['imgs']['map'] = tmp_map.getvalue()
        # 고정 크기 배치용 비율(25칸 고정 → wafer 수와 무관하게 동일). 조립부에서 사용.
        out['map_ratio'] = fig_h / fig_disp_w

        # ---- WF MAP 컬러바(별도 이미지) : 슬라이드에서 항상 같은 크기로 우측에 세로로 길게 배치 ----
        if sc is not None:
            import matplotlib.cm as _cm
            tmp_cbar = _io.BytesIO()
            _sm = _cm.ScalarMappable(norm=wfmap_norm, cmap=_wfmap_cmap(direction))
            _sm.set_array([])
            fig_cb = plt.figure(figsize=(0.62, 3.2))
            _cax = fig_cb.add_axes([0.04, 0.04, 0.30, 0.92])
            _cb = fig_cb.colorbar(_sm, cax=_cax)
            _cb.ax.tick_params(labelsize=7)
            fig_cb.savefig(tmp_cbar, format='png', dpi=max(int(dpi), 150),
                           bbox_inches='tight', facecolor='white')
            plt.close(fig_cb)
            out['imgs']['cbar'] = tmp_cbar.getvalue()

        # ---- 5. Trend Chart (vehicle/with_vehicle/target 비교 + vehicle 1~99% 구름대) ----
        tdf = item_df_full.copy()
        tdf['tkout_time'] = pd.to_datetime(tdf['tkout_time'])
        has_mask = 'mask' in tdf.columns
        has_lot = 'fab_lot_id' in tdf.columns

        # ---- 특정 항목: site별 모든 값 대신 tkout_time 기준 집계점으로 Trend 표시 ----
        #   config.trend_tkout_agg = {ALIAS: 'P10'/'P90'/'MEDIAN'/'MEAN'}. base ALIAS 키('MAWIN')로
        #   파생 컬럼(MAWIN_*)까지 매칭 — 이상/주의 판정(anomaly_engine)과 '동일한' 규칙(trend_agg_spec).
        from anomaly_engine import trend_agg_spec as _trend_agg_spec
        _agg_map = cfg.get('trend_tkout_agg') or {}
        _agg_spec = _trend_agg_spec(item_name, _agg_map, spec_name)
        if _agg_spec:
            _s = str(_agg_spec).strip().upper()
            if _s in ('MEAN', 'AVG'):
                _aggf = 'mean'
            elif _s in ('MEDIAN', 'P50'):
                _aggf = 'median'
            else:
                _m = re.match(r'P(\d+(?:\.\d+)?)$', _s)
                _aggf = (lambda s, _q=min(max(float(_m.group(1)) / 100.0, 0.0), 1.0): s.quantile(_q)) if _m else 'median'
            _gk = [k for k in ['mask', 'fab_lot_id', 'root_lot_id', 'wafer_id', 'match_key', 'tkout_time']
                   if k in tdf.columns]
            if 'tkout_time' in _gk and len(_gk) > 1:
                tdf = (tdf.dropna(subset=[item_name])
                          .groupby(_gk, as_index=False).agg({item_name: _aggf}))
                tdf['tkout_time'] = pd.to_datetime(tdf['tkout_time'])

        veh_df = tdf[tdf['mask'] == main_vehicle] if has_mask else tdf

        # ---- Trend: root_lot_id 그룹별 (색·마커) 배정 ----
        #   · 같은 root_lot_id 형제 lot은 '동일 색·모양' → 그룹으로 묶여 보임
        #   · 초록(vehicle 구름대)·빨강(리포트 lot 전용)은 팔레트에서 제외
        #   · 색(12)·마커(7) 길이를 서로소로 두어 root가 많아도 (색,모양) 조합이 오래 안 겹침
        _root_col = 'root_lot_id' if 'root_lot_id' in tdf.columns else lot_col
        # 초록(vehicle)·빨강(리포트 lot)과 그 유사색은 제외 — 파랑/보라/갈색/청록/자홍/주황/회색 계열
        _ROOT_COLORS = ['#1f77b4', '#9467bd', '#8c564b', '#17becf', '#e377c2', '#7f7f7f',
                        '#ff7f0e', '#bcbd22', '#393b79', '#00868b', '#5254a3', '#636363']
        _ROOT_MARKERS = ['o', 's', '^', 'D', 'v', 'P', 'X']
        root_style = {}
        if _root_col:
            for _i, _rt in enumerate(sorted(tdf[_root_col].astype(str).dropna().unique())):
                root_style[_rt] = (_ROOT_COLORS[_i % len(_ROOT_COLORS)],
                                   _ROOT_MARKERS[_i % len(_ROOT_MARKERS)])

        def _draw_trend(ax):
            # vehicle 기준 1~99% 구름대 + median (main vehicle 데이터만)
            if len(veh_df) > 0:
                b = veh_df[['tkout_time', item_name]].dropna().sort_values('tkout_time')
                if len(b) > 0:
                    b = b.assign(date=b['tkout_time'].dt.date)
                    daily = b.groupby('date')[item_name].agg(
                        median='median',
                        q01=lambda x: x.quantile(0.01),
                        q99=lambda x: x.quantile(0.99)).reset_index()
                    daily['date'] = pd.to_datetime(daily['date'])
                    daily = daily.set_index('date').sort_index()
                    roll = daily.rolling('3D', min_periods=1).mean()
                    ax.fill_between(roll.index, roll['q01'], roll['q99'], color=C_BAND, alpha=0.6, label=f'{main_vehicle} 1~99%', zorder=1)
                    # 3일 기준 rolling median 라인 (검정색)
                    ax.plot(roll.index, roll['median'], color='black', linewidth=1.5, alpha=1.0, zorder=6, label='3-day median')
            # 색상/모양 규칙:
            #   · 리포트 lot_id = 빨강 원(o), 최상단 zorder — 다른 모든 마커 위에 확실히 표시
            #   · 형제 lot(리포트와 같은 root_lot_id) = 리포트 root 그룹의 색·모양(빨강 아님)
            #   · 그 외 main-vehicle lot = 각자 root_lot_id 그룹의 색·모양(초록/빨강 제외)으로 묶여 보임
            #   · with_vehicle(다른 vehicle) = WV_PALETTE(회색 계열)
            #   모든 마커는 얇은 검정 테두리(edgecolors='black')
            if has_lot:
                tgt = _select_target_lot_frame(tdf, target_lot_id, target_root_lot_id, target_DC_step_id)
                tgt_idx = set(tgt.index)
                if has_mask:
                    # 같은 vehicle이면서 대상(리포트 root)이 아닌 데이터
                    veh_other = tdf[(tdf['mask'] == main_vehicle) & (~tdf.index.isin(tgt_idx))]
                    # with_vehicle(다른 vehicle) 데이터
                    wv = tdf[tdf['mask'] != main_vehicle]
                else:
                    veh_other = tdf[~tdf.index.isin(tgt_idx)]; wv = tdf.iloc[0:0]
                # (배경) 타 root의 main-vehicle lot — root_lot_id 그룹별 색·모양, 범례는 1개로 통합
                if len(veh_other) > 0:
                    ax.scatter(veh_other['tkout_time'], veh_other[item_name], s=10, alpha=0.5,
                               color=C_VEHICLE, label=str(main_vehicle),
                               edgecolors='black', linewidths=0.3, zorder=2)
                if len(wv) > 0:
                    # with_vehicle은 mask(=실제 vehicle 명)별로 분리하여 각각 다른 색 + 개별 범례
                    for _wi, _wv_name in enumerate(sorted(wv['mask'].dropna().unique()) if has_mask else []):
                        _wv_grp = wv[wv['mask'] == _wv_name]
                        if len(_wv_grp) == 0: continue
                        ax.scatter(_wv_grp['tkout_time'], _wv_grp[item_name], s=10, alpha=0.5, color=C_WV, label=str(_wv_name), edgecolors='black', linewidths=0.3, zorder=3)
                if len(tgt) > 0:
                    _report_lot = str(target_lot_id)
                    # 리포트 root 그룹의 색·모양(형제 lot에 사용)
                    _tgt_root = (str(target_root_lot_id) if target_root_lot_id
                                 else (str(tgt[_root_col].iloc[0]) if (_root_col and len(tgt)) else None))
                    _tc, _tm = root_style.get(str(_tgt_root), ('#1f77b4', 's'))
                    if lot_col:
                        _sib = tgt[tgt[lot_col].astype(str) != _report_lot]
                        _rep = tgt[tgt[lot_col].astype(str) == _report_lot]
                    else:
                        _sib = tgt.iloc[0:0]; _rep = tgt
                    # 형제 lot(같은 root) → 그룹 색·모양, 리포트 lot 바로 아래 zorder
                    if len(_sib) > 0:
                        ax.scatter(_sib['tkout_time'], _sib[item_name], s=28, alpha=1.0,
                                   color=_tc, marker=_tm,
                                   label=f"형제 lot ({_tgt_root})" if _tgt_root else "형제 lot",
                                   edgecolors='black', linewidths=0.5, zorder=11)
                    # 리포트 lot → 빨강 원, 최상단
                    if len(_rep) > 0:
                        ax.scatter(_rep['tkout_time'], _rep[item_name], s=34, alpha=1.0, color='red',
                                   marker='o', label=f"{_report_lot}_{target_DC_step_id}",
                                   edgecolors='black', linewidths=0.6, zorder=13)
            else:
                for w in measured_wafers:
                    grp = tdf[tdf[w_col] == w] if w_col in tdf.columns else tdf.iloc[0:0]
                    ax.scatter(grp['tkout_time'], grp[item_name], s=10, alpha=0.7, color=w_colors.get(str(w), 'blue'), edgecolors='black', linewidths=0.3, zorder=2)
            # spec line(s) — 방향(REPORT DIRECTION) 반영된 spec_low/high (범례에는 표시하지 않음)
            if spec_low is not None:
                ax.axhline(y=float(spec_low), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
            if spec_high is not None:
                ax.axhline(y=float(spec_high), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)

            # ── y축 범위 자동 조정: 해당 lot + spec line 중심 ──
            #   이전(모집단) high flier로 y축이 과도하게 커져 대상 lot이 안 보이는 문제 →
            #   target lot·모집단 모두 'robust 범위(1~99%)'만 반영하고 spec line은 반드시 포함.
            #   범위를 벗어난 high flier는 잘려도 무방(대상 lot과 spec을 명확히 보이게 함).
            if not log_scale:
                try:
                    _tgt_y = _select_target_lot_frame(tdf, target_lot_id, target_root_lot_id, target_DC_step_id)
                except Exception:
                    _tgt_y = tdf.iloc[0:0]
                _yv = []
                for _d in (_tgt_y, veh_df):
                    if len(_d):
                        _s = pd.to_numeric(_d[item_name], errors='coerce').dropna()
                        if len(_s):
                            _yv += [float(_s.quantile(0.01)), float(_s.quantile(0.99))]
                if spec_low is not None:
                    _yv.append(float(spec_low))
                if spec_high is not None:
                    _yv.append(float(spec_high))
                _yv = [v for v in _yv if pd.notna(v)]
                if _yv:
                    _ylo, _yhi = min(_yv), max(_yv)
                    # 데이터/ spec 범위 위아래로 12% 여유(spec 초과 값도 빡빡하지 않게)
                    _pad = (_yhi - _ylo) * 0.12 if _yhi > _ylo else (abs(_yhi) * 0.12 or 1.0)
                    ax.set_ylim(_ylo - _pad, _yhi + _pad)
            ax.set_title("")
            ax.xaxis.set_major_formatter(mdates.DateFormatter('%m/%d'))
            ax.tick_params(axis='x', rotation=0, labelsize=7)
            ax.xaxis.set_major_locator(mdates.AutoDateLocator(maxticks=6))
            _label_axes(ax, xlabel="DC tkout_time", ylabel=y_label, ylabel_size=5.5)  # y축명 잘림 방지 위해 축소
            if log_scale: ax.set_yscale('log')
            # Trend 범례 좌상단 고정 — 흰 배경 + 옅은 테두리, 글자·항목 간격 축소(컴팩트)
            _leg = ax.legend(fontsize=6, loc='upper left', frameon=True, facecolor='white',
                             edgecolor='#b0b0b0', framealpha=0.95, labelspacing=0.18,
                             handletextpad=0.3, borderpad=0.3, borderaxespad=0.3)
            if _leg is not None:
                _leg.set_zorder(20)
                _leg.get_frame().set_linewidth(0.6)
            _remove_spines(ax)
            ax.minorticks_off()  # minor tick(세부선) 제거 — major만 표시
            ax.grid(True, which='major', color=C_GRID, linestyle='-', linewidth=0.5)

        fig_trend, ax_trend = plt.subplots(figsize=(4.55, 1.75))
        _draw_trend(ax_trend)
        fig_trend.savefig(tmp_trend, format='jpg', dpi=dpi, bbox_inches="tight", facecolor='white', pil_kwargs={'quality': jpg_q})
        plt.close(fig_trend)
        out['imgs']['trend'] = tmp_trend.getvalue()

        # index(alias) Trend scatter 차트만 RUN/TEMP에 alias명.png로 저장 (Anomaly/HTML 재사용)
        safe_name = re.sub(r'[\\/:*?"<>|]', '_', str(item_name))
        try:
            run_temp = cfg.get('run_temp_dir') or os.path.join('RUN', 'TEMP')
            os.makedirs(run_temp, exist_ok=True)
            fig_trend_png, ax_trend_png = plt.subplots(figsize=(4.55, 2.0))
            _draw_trend(ax_trend_png)
            fig_trend_png.savefig(os.path.join(run_temp, f"{safe_name}.png"),
                                  dpi=cfg.get('html_chart_dpi', 100), bbox_inches="tight")
            plt.close(fig_trend_png)
        except Exception as e:
            out['warnings'].append(f"Failed to save RUN/TEMP/{safe_name}.png: {e}")

        # ---- 지표 (Metrics) 계산 ----
        try:
            # global = vehicle 모집단(veh_df), target = 리포팅 대상 lot(item_df=target_df)
            _gsrc = veh_df if len(veh_df) > 0 else item_df
            g_med = float(_gsrc[item_name].median())
            g_std = float(_gsrc[item_name].std())

            t_df = item_df
            t_med = float(t_df[item_name].median()) if len(t_df) > 0 else g_med
            t_std = float(t_df[item_name].std()) if len(t_df) > 1 else g_std

            dev = round(abs(t_med - g_med) / g_std, 2) if g_std > 0 else 0.0

            # spec-out(이상) 분류는 '해당 lot_id'에만 한정 — 같은 root의 형제 lot은 이상으로 분류하지 않음
            if 'fab_lot_id' in t_df.columns:
                _t_spec = t_df[t_df['fab_lot_id'].astype(str) == str(target_lot_id)]
                if len(_t_spec) == 0:
                    _t_spec = t_df
            else:
                _t_spec = t_df

            # spec 이탈 칩 목록 (기존 행별 iterrows → 벡터화)
            _v = pd.to_numeric(_t_spec[item_name], errors='coerce')
            _m = pd.Series(False, index=_t_spec.index)
            if spec_low is not None:
                _m = _m | (_v < float(spec_low))
            if spec_high is not None:
                _m = _m | (_v > float(spec_high))
            _xcol = 'CHIP_X_ADJ' if 'CHIP_X_ADJ' in _t_spec.columns else (col_x if col_x in _t_spec.columns else None)
            _ycol = 'CHIP_Y_ADJ' if 'CHIP_Y_ADJ' in _t_spec.columns else (col_y if col_y in _t_spec.columns else None)
            _ov = _v[_m]
            _ox = (pd.to_numeric(_t_spec.loc[_m, _xcol], errors='coerce').fillna(0.0)
                   if _xcol else pd.Series(0.0, index=_ov.index))
            _oy = (pd.to_numeric(_t_spec.loc[_m, _ycol], errors='coerce').fillna(0.0)
                   if _ycol else pd.Series(0.0, index=_ov.index))
            s_outs = [{"val": float(v), "x": float(x), "y": float(y)}
                      for v, x, y in zip(_ov, _ox, _oy)]

            out['metrics'] = {
                "item": item_name,
                "global_med": g_med,
                "global_std": g_std,
                "target_med": t_med,
                "target_std": t_std,
                "deviation": dev,
                "spec_outs": s_outs,
                "spec_out_count": len(s_outs)
            }
        except Exception as e:
            out['warnings'].append(f"Failed to calculate metrics: {e}")

        # ---- 6. Radius Plot (반경별 분포 + 3차 다항식 근사) ----
        fig_rad, ax_rad = plt.subplots(figsize=(4.55, 1.75))
        for w in measured_wafers:
            w_df = item_df[item_df[w_col] == w].dropna(subset=[col_rad, item_name]).sort_values(by=col_rad)
            if len(w_df) == 0: continue
            p_color = w_colors.get(str(w), 'blue')
            # 색=wafer, marker=lot_id (multi_lot일 때만 lot별 모양 구분)
            if multi_lot and lot_col:
                for _lot in target_lots:
                    _wl = w_df[w_df[lot_col].astype(str) == _lot]
                    if len(_wl) == 0: continue
                    ax_rad.scatter(_wl[col_rad], _wl[item_name], s=14, alpha=0.85, color=p_color,
                                   marker=lot_marker[_lot], edgecolors='black', linewidths=0.2)
            else:
                ax_rad.scatter(w_df[col_rad], w_df[item_name], s=12, alpha=0.85, color=p_color)
        # 피팅선: wafer별 '3차(cubic)' 근사선을 각각 그리고, 선 색은 해당 wafer 색과 동일하게 맞춘다.
        for w in measured_wafers:
            try:
                _wf = item_df[item_df[w_col] == w].dropna(subset=[col_rad, item_name])
                if len(_wf) < 4:
                    continue   # cubic(3차) fit 최소 4점 필요
                _rx = pd.to_numeric(_wf[col_rad], errors='coerce').astype(float).values
                _ry = pd.to_numeric(_wf[item_name], errors='coerce').astype(float).values
                _m = np.isfinite(_rx) & np.isfinite(_ry)
                _rx, _ry = _rx[_m], _ry[_m]
                if len(_rx) < 4 or _rx.min() >= _rx.max():
                    continue
                _z = np.polyfit(_rx, _ry, 3)               # wafer별 3차 다항식(cubic)
                _pp = np.poly1d(_z)
                _xs = np.linspace(_rx.min(), _rx.max(), 100)
                ax_rad.plot(_xs, _pp(_xs), color=w_colors.get(str(w), 'blue'),
                            alpha=0.9, linewidth=1.3)
            except Exception:
                continue
        # lot_id marker 범례 (multi_lot)
        if multi_lot:
            from matplotlib.lines import Line2D
            _lh = [Line2D([0], [0], marker=lot_marker[_lot], color='#444444', linestyle='none',
                          markersize=6, label=str(_lot)) for _lot in target_lots]
            ax_rad.legend(handles=_lh, fontsize=6, loc='upper left', bbox_to_anchor=(1.0, 1.0),
                          frameon=False, title='Lot', title_fontsize=6)   # 차트 밖 우상단 고정
        # Spec line은 그리되 범례(Spec Limit)는 표시하지 않음
        if spec_low is not None:
            ax_rad.axhline(y=float(spec_low), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        if spec_high is not None:
            ax_rad.axhline(y=float(spec_high), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        ax_rad.set_title("")
        _label_axes(ax_rad, xlabel="Chip Radius", ylabel=y_label, ylabel_size=5.5)  # y축명 잘림 방지 위해 축소
        if log_scale: ax_rad.set_yscale('log')
        _remove_spines(ax_rad)
        ax_rad.minorticks_off()  # minor tick(세부선) 제거 — major만 표시
        ax_rad.grid(True, which='major', color=C_GRID, linestyle='-', linewidth=0.5)
        fig_rad.savefig(tmp_rad, format='jpg', dpi=dpi, bbox_inches="tight", facecolor='white', pil_kwargs={'quality': jpg_q})
        plt.close(fig_rad)
        out['imgs']['rad'] = tmp_rad.getvalue()

        # ---- 7. Cumulative Plot (누적 분포) ----
        fig_cum, ax_cum = plt.subplots(figsize=(4.55, 1.85))
        for w in fixed_wafers:
            p_color = w_colors.get(str(w), 'blue')
            if multi_lot and lot_col:
                # 색=wafer, marker=lot_id. 같은 wafer라도 lot별로 모양이 달라짐
                for _lot in target_lots:
                    _sub = item_df[(item_df[w_col] == w) & (item_df[lot_col].astype(str) == _lot)][item_name].dropna()
                    if len(_sub) == 0: continue
                    _sd = np.sort(_sub.values)
                    _yv = np.arange(len(_sd)) / float(len(_sd) - 1) if len(_sd) > 1 else [1.0]
                    ax_cum.plot(_sd, _yv, marker=lot_marker[_lot], linestyle='-', markersize=4,
                                linewidth=0.9, color=p_color, alpha=0.8)
            elif w in grouped.groups:
                grp = grouped.get_group(w)
                sorted_data = np.sort(grp.values)
                yvals = np.arange(len(sorted_data)) / float(len(sorted_data) - 1) if len(sorted_data) > 1 else [1.0]
                # 같은 wafer_id 데이터를 선으로 연결 (마커 + 라인)
                ax_cum.plot(sorted_data, yvals, marker='.', linestyle='-', markersize=4, linewidth=1.0, color=p_color, alpha=0.8)
        if multi_lot:
            from matplotlib.lines import Line2D
            _lh = [Line2D([0], [0], marker=lot_marker[_lot], color='#444444', linestyle='none',
                          markersize=6, label=str(_lot)) for _lot in target_lots]
            ax_cum.legend(handles=_lh, fontsize=6, loc='upper left', bbox_to_anchor=(1.0, 1.0),
                          frameon=False, title='Lot', title_fontsize=6)   # 차트 밖 우상단 고정
        # Spec line은 그리되 범례(Spec Limit)는 표시하지 않음
        if spec_low is not None:
            ax_cum.axvline(x=float(spec_low), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        if spec_high is not None:
            ax_cum.axvline(x=float(spec_high), color=C_ACCENT, ls="--", lw=1.2, alpha=0.7)
        ax_cum.set_title("")
        _label_axes(ax_cum, xlabel=y_label, ylabel="Cumulative Prob.")
        if log_scale: ax_cum.set_xscale('log')
        _remove_spines(ax_cum)
        ax_cum.minorticks_off()  # minor tick(세부선) 제거 — major만 표시
        ax_cum.grid(True, which='major', color=C_GRID, linestyle='-', linewidth=0.5)
        fig_cum.savefig(tmp_cum, format='jpg', dpi=dpi, bbox_inches="tight", facecolor='white', pil_kwargs={'quality': jpg_q})
        plt.close(fig_cum)
        out['imgs']['cum'] = tmp_cum.getvalue()

        # ---- 임시 버퍼 해제 ----
        for _b in [tmp_box, tmp_map, tmp_trend, tmp_rad, tmp_cum]:
            try: _b.close()
            except Exception: pass
        gc.collect()

    except Exception as e:
        import traceback as _tb
        out['status'] = 'error'
        out['reason'] = f"{e}\n{_tb.format_exc()}"
    return out


def _wfmap_batch_task(payload):
    """[워커] Score Board wafer WF MAP을 항목 묶음(chunk) 단위로 렌더링.

    payload = {'df': 컬럼 슬라이스 df, 'items': [{'item','direction','spec_low','spec_high'}],
               'common': {'min_pts','lot_prefix','dpi','by_lot'}}
    """
    df = payload['df']
    common = payload['common']
    out = {}
    for it in payload['items']:
        try:
            m = render_wafer_wfmaps_b64(df, it['item'],
                                        min_pts=common['min_pts'], lot_prefix=common['lot_prefix'],
                                        direction=it.get('direction', 'BOTH'),
                                        spec_low=it.get('spec_low'), spec_high=it.get('spec_high'),
                                        dpi=common['dpi'], by_lot=common['by_lot'])
        except Exception:
            m = {}
        if m:
            out[it['item']] = m
    return out


def render_wafer_wfmaps_batch(df, item_specs, min_pts=50, lot_prefix=None,
                              dpi=None, by_lot=True, workers=None):
    """여러 index의 wafer WF MAP(render_wafer_wfmaps_b64)을 병렬 렌더링.

    item_specs: [{'item', 'direction', 'spec_low', 'spec_high'}, ...]
    반환: {item: {wafer_key: b64}} (WF MAP 없는 item은 미포함)
    워커 수는 get_parallel_workers()로 자동 결정(1이면 직렬). 워커당 1 chunk로
    묶어 df 전송(pickle)을 항목당이 아닌 워커당 1회로 최소화한다.
    """
    if dpi is None:
        dpi = getattr(GLOBAL_CONFIG, 'html_wfmap_dpi', 110)
    if workers is None:
        workers = get_parallel_workers()

    common = {'min_pts': min_pts, 'lot_prefix': lot_prefix, 'dpi': dpi, 'by_lot': by_lot}

    def _serial():
        out = {}
        for it in item_specs:
            m = render_wafer_wfmaps_b64(df, it['item'], min_pts=min_pts, lot_prefix=lot_prefix,
                                        direction=it.get('direction', 'BOTH'),
                                        spec_low=it.get('spec_low'), spec_high=it.get('spec_high'),
                                        dpi=dpi, by_lot=by_lot)
            if m:
                out[it['item']] = m
        return out

    if workers <= 1 or len(item_specs) <= 1:
        return _serial()

    pool = _get_chart_pool(workers)
    if pool is None:
        return _serial()

    meta_cols = [c for c in ['CHIP_X_ADJ', 'CHIP_Y_ADJ', 'CHIP_X_POS', 'CHIP_Y_POS',
                             'chip_x_pos', 'chip_y_pos', 'WAFER_ID', 'wafer_id',
                             'TKOUT_TIME', 'tkout_time', 'FAB_LOT_ID', 'fab_lot_id',
                             'Chip_Radius', 'chip_radius']   # 원(150mm) radius fit용 — 빠지면 fallback 원
                 if c in df.columns]
    chunks = [item_specs[i::workers] for i in range(workers)]
    try:
        futs = []
        for ch in chunks:
            if not ch:
                continue
            _icols = [it['item'] for it in ch if it['item'] in df.columns]
            _sub = df[list(dict.fromkeys(meta_cols + _icols))]
            futs.append(pool.submit(_wfmap_batch_task,
                                    {'df': _sub, 'items': ch, 'common': common}))
        out = {}
        for f in futs:
            out.update(f.result())
        return out
    except Exception as e:
        print(f"[WARN] WF MAP 병렬 렌더링 실패 → 직렬 폴백: {e}")
        return _serial()


def insert_plots(merged_df, prs, description_image_info_dict,
                 target_lot_id, target_root_lot_id,
                 target_DC_step, target_DC_step_id,
                 spec_data, img_quality=12, ref=False, reformatter=None, dpi=None):
    """최고급 다중 차트 및 통계표 대시보드를 생성하여 PPT에 삽입.

    각 측정 항목(item)에 대해 하나의 슬라이드를 생성하며, 슬라이드 구성:
        - 좌측 상단: 통계 테이블 (Mean/Std/Min/Max, 전체 + 웨이퍼별)
        - 좌측 중앙: BOX Plot (웨이퍼별 분포)
        - 좌측 하단: WF MAP (PGM × Wafer 산점도)
        - 우측 상단: Trend Chart (시계열 + 3일 구름대)
        - 우측 중앙: Radius Plot (반경별 분포 + 3차 근사)
        - 우측 하단: Cumulative Plot (누적 분포)

    렌더링(matplotlib, CPU 바운드)은 _render_item_charts를 통해 워커 프로세스에서
    병렬 수행하고, 이 함수는 결과 이미지를 순서대로 PPT에 조립한다.
    워커 수는 get_parallel_workers()가 환경(CPU/가용 메모리)을 보고 자동 결정하며
    1이면 종전과 같은 직렬 렌더링으로 동작한다.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    import io as _io
    import gc
    import numpy as np
    import re

    # ---- 차트 해상도/압축 품질 설정 (My_config.py에서 조정) ----
    # dpi가 명시되지 않으면 GLOBAL_CONFIG.ppt_chart_dpi 사용 (엔지니어가 직접 변경 가능)
    if dpi is None:
        dpi = GLOBAL_CONFIG.ppt_chart_dpi
    jpg_q = GLOBAL_CONFIG.ppt_chart_jpg_quality   # 라인/박스/트렌드/CDF/레전드 JPEG 품질
    map_q = GLOBAL_CONFIG.ppt_map_jpg_quality     # WF MAP JPEG 품질
    # 차트 공통 색 팔레트 (config에서 단일 관리)
    C_NAVY = GLOBAL_CONFIG.chart_navy
    C_ACCENT = GLOBAL_CONFIG.chart_accent
    C_NEUTRAL = GLOBAL_CONFIG.chart_neutral
    C_GRID = GLOBAL_CONFIG.chart_grid
    C_SPINE = GLOBAL_CONFIG.chart_spine
    C_VEHICLE = GLOBAL_CONFIG.chart_vehicle          # main vehicle scatter (초록)
    C_WV = GLOBAL_CONFIG.chart_with_vehicle          # with_vehicle scatter (회색, 단일 fallback)
    # with_vehicle이 여러 개일 때 각각 다른 색으로 구분하기 위한 팔레트
    # (첫번째는 기존 회색 유지, 이후 항목은 구분되는 색)
    WV_PALETTE = [C_WV, '#9467bd', '#8c564b', '#17becf', '#bcbd22', '#e377c2', '#ff7f0e', '#1f77b4']
    C_BAND = GLOBAL_CONFIG.chart_band                # vehicle 1~99% 구름대 (연초록)
    FONT = GLOBAL_CONFIG.theme_font_family
    NAVY_RGB = GLOBAL_CONFIG.theme_title_color    # (R,G,B) 헤더/타이틀 통일 색
    main_vehicle = GLOBAL_CONFIG.get('vehicle')   # mask 비교 기준 (main vehicle 이름)

    # --- 컬럼명 대소문자 방어 (Column Name Case-Insensitive Lookup) ---
    def _pick(*names):
        for n in names:
            if n in merged_df.columns:
                return n
        return None
    col_x = _pick('chip_x_pos', 'CHIP_X_POS')
    col_y = _pick('chip_y_pos', 'CHIP_Y_POS')
    col_rad = _pick('chip_radius', 'Chip_Radius')
    col_sub = _pick('subitem_id', 'SUBITEM_ID')
    col_time = _pick('tkout_time', 'TKOUT_TIME')   # merged_df는 대문자(TKOUT_TIME)일 수 있음
    col_lot = _pick('fab_lot_id', 'FAB_LOT_ID')
    col_mask = _pick('mask', 'MASK')

    current_cat = None
    _pending_desc = None   # (cat2, entry): 데이터가 실제로 찍힌 CAT2에만 간지 삽입하기 위한 대기 슬롯
    item_slide_map = {}    # {item_name: chart slide} — Anomaly 상세의 아이템명 내부 링크 대상
    metrics_dict = {}

    import os
    os.makedirs("RUN/TEMP", exist_ok=True)

    # --- 페이지 index 목록 구성 (spec_data 항목 + Reformatize 파생 항목 포함) ---
    # Reformatize는 ADDP FORM으로 단일/다중 컬럼 파생 index를 만든다.
    #  - 단일컬럼 파생(예: VTH_AVG)  → merged_df에 'VTH_AVG' 컬럼 그대로 존재
    #  - 다중컬럼 파생(예: window)   → 'WINDOW_ovl_index', 'WINDOW_new' 등 접두어_컬럼으로 존재
    # 따라서 spec_data의 각 ALIAS에 대해, 그 ALIAS로 시작하는 모든 merged_df 컬럼을
    # 페이지 대상에 포함시켜 파생 항목도 빠짐없이 PPT 페이지가 생성되도록 한다.
    # plot_items: (데이터 컬럼명, spec 메타 출처 ALIAS)
    # REPORT ORDER 오름차순(작은 값 먼저)으로 index 정렬
    if 'REPORT ORDER' in spec_data.columns:
        # REPORT ORDER에 '숫자'가 있는 항목만 페이지 생성(비어있거나 숫자가 아니면 제외)
        _ro = pd.to_numeric(spec_data['REPORT ORDER'], errors='coerce')
        ordered_index = _ro.dropna().sort_values(kind='stable').index
    else:
        ordered_index = spec_data.index
    plot_items = []
    seen_cols = set()
    for nm in ordered_index:
        # alias 자체가 컬럼이면 포함(MA_Window의 '' 출력 = alias) + alias_로 시작하는 다중컬럼 파생 모두 포함
        cols_for_nm = []
        if nm in merged_df.columns:
            cols_for_nm.append(nm)
        cols_for_nm += [c for c in merged_df.columns if str(c).startswith(str(nm) + "_")]
        if cols_for_nm:
            for d in cols_for_nm:
                if d not in seen_cols:
                    plot_items.append((d, nm))
                    seen_cols.add(d)
        else:
            plot_items.append((nm, nm))  # 데이터 없음 → 루프에서 skip

    # ── ADDP 파생 index 제목 옆 'real 항목' 표시용 준비 ──
    # (1) 대상 lot에서 실제로 데이터가 있는(비어있지 않은) 항목 컬럼 집합
    try:
        _tl_df = _select_target_lot_frame(merged_df, target_lot_id, target_root_lot_id, target_DC_step_id)
        _lot_present_cols = {c for c in _tl_df.columns if _tl_df[c].notna().any()}
    except Exception:
        _lot_present_cols = set(merged_df.columns)
    # (2) reformatter ALIAS → (CATEGORY, ADDP FORM, ITEMID) 맵 (ADDP 재귀 해소·real item 표시용)
    _ref_cat, _ref_formula, _ref_itemid = {}, {}, {}
    if reformatter is not None and 'ALIAS' in getattr(reformatter, 'columns', []):
        for _, _rr in reformatter.iterrows():
            _a = _rr.get('ALIAS')
            if pd.isna(_a):
                continue
            _ref_cat[str(_a)] = str(_rr.get('CATEGORY', '')).strip().upper()
            _ref_formula[str(_a)] = str(_rr.get('ADDP FORM', '') or '')
            _iid = _rr.get('ITEMID')
            _ref_itemid[str(_a)] = '' if pd.isna(_iid) else str(_iid).strip()

    # 다중출력 ADDP(예: MA_Window)는 {alias}_new / {alias}_ovl_index 등 파생 컬럼을 만든다.
    # 다른 ADDP(rmax 등)가 이 파생 컬럼을 {MAWIN_new} 처럼 참조하면 그 자체는 ALIAS가
    # 아니라서 leaf로 잡히므로, 이름 접두가 일치하는 base ADDP ALIAS를 찾아 재귀 해소한다.
    _addp_alias_list = [a for a, c in _ref_cat.items() if c == 'ADDP']

    def _base_addp_of(ref):
        """ref가 어떤 ADDP의 파생 컬럼({base}_{subcol})이면 그 base ADDP ALIAS를 반환(가장 긴 접두 우선)."""
        _cands = [a for a in _addp_alias_list if ref == a or ref.startswith(str(a) + '_')]
        return max(_cands, key=lambda a: len(str(a))) if _cands else None

    def _addp_leaf_refs(alias, _seen=None):
        """ADDP FORM의 {ALIAS} 참조를 재귀 전개해 leaf(비-ADDP) real item 집합 반환.
        함수명(rmax/ABS/MA_Window 등)은 {} 밖이라 제외된다. 다중출력 ADDP의 파생
        컬럼 참조({MAWIN_new} 등)도 base ADDP를 찾아 최종 real item까지 재귀 해소한다."""
        if _seen is None:
            _seen = set()
        out = set()
        for ref in re.findall(r'\{([^}]+)\}', _ref_formula.get(str(alias), '')):
            ref = ref.strip()
            if not ref or ref in _seen:
                continue
            _seen.add(ref)
            if _ref_cat.get(ref) == 'ADDP':
                out |= _addp_leaf_refs(ref, _seen)
                continue
            # ALIAS가 아니면(파생 컬럼일 수 있음) base ADDP를 찾아 재귀
            _base = _base_addp_of(ref) if ref not in _ref_cat else None
            if _base is not None:
                out |= _addp_leaf_refs(_base, _seen)
            else:
                out.add(ref)
        return out

    # ================================================================
    # 병렬 렌더링 준비 — 렌더링(워커)과 PPT 조립(메인)을 겹쳐 수행
    # ================================================================
    total_items = len(plot_items)
    summary_rows = []   # 마지막 summary 페이지용 (index별 REPORT DIRECTION 기준 집계값)

    n_workers = get_parallel_workers()
    _avail = _get_available_mem_gb()
    print("=" * 60)
    print(f"[insert_plots] 차트 생성 시작 - 총 {total_items}개 index(파생 포함) 처리 예정")
    print(f"[insert_plots] 렌더링 워커 {n_workers}개 "
          f"(cores={os.cpu_count()}, 가용메모리={f'{_avail:.1f}GB' if _avail is not None else '측정불가'})")
    print("=" * 60)

    # 워커에 넘길 경량 설정(dict) — 워커는 GLOBAL_CONFIG(yaml 미로드)를 참조하지 않는다
    cfg_task = {
        'dpi': dpi, 'jpg_q': jpg_q, 'map_q': map_q,
        'C_NAVY': C_NAVY, 'C_ACCENT': C_ACCENT, 'C_NEUTRAL': C_NEUTRAL,
        'C_GRID': C_GRID, 'C_SPINE': C_SPINE, 'C_VEHICLE': C_VEHICLE,
        'C_BAND': C_BAND, 'WV_PALETTE': WV_PALETTE, 'FONT': FONT,
        'main_vehicle': main_vehicle,
        'trend_tkout_agg': getattr(GLOBAL_CONFIG, 'trend_tkout_agg', {}) or {},
        'html_chart_dpi': getattr(GLOBAL_CONFIG, 'html_chart_dpi', 100),
        'run_temp_dir': os.path.abspath(os.path.join('RUN', 'TEMP')),
    }
    w_col_src = _pick('wafer_id', 'WAFER_ID') or 'wafer_id'
    cols_task = {'col_x': col_x, 'col_y': col_y, 'col_rad': col_rad, 'col_sub': col_sub,
                 'col_time': col_time, 'col_lot': col_lot, 'col_mask': col_mask,
                 'w_col_src': w_col_src}

    # 모든 항목 공통 메타 컬럼 (per-item 데이터 슬라이스용)
    req_cols = [w_col_src]
    if col_x: req_cols += [col_x, col_y]
    if col_rad: req_cols.append(col_rad)
    if col_time: req_cols.append(col_time)
    if col_sub: req_cols.append(col_sub)
    if col_lot: req_cols.append(col_lot)
    if col_mask: req_cols.append(col_mask)
    if 'CHIP_X_ADJ' in merged_df.columns: req_cols.append('CHIP_X_ADJ')
    if 'CHIP_Y_ADJ' in merged_df.columns: req_cols.append('CHIP_Y_ADJ')
    if 'match_key' in merged_df.columns: req_cols.append('match_key')  # 형제 lot 묶음(root+step) 선택용
    if 'PGM(pt)' in merged_df.columns: req_cols.append('PGM(pt)')      # WF MAP 행 좌측 PGM(pt) 라벨용

    def _build_task(item_name, spec_name):
        """항목 1개의 렌더링 task 생성. 데이터가 없으면 None."""
        item_df = merged_df[list(set(req_cols)) + [item_name]].dropna(subset=[item_name])
        if len(item_df) == 0:
            return None

        # ---- 스펙 범위 추출 (Spec Limits) + 방향/로그/단위 ----
        spec_low = spec_data.loc[spec_name, "SPECLOW"]
        spec_high = spec_data.loc[spec_name, "SPECHIGH"]
        if pd.isna(spec_low): spec_low = None
        if pd.isna(spec_high): spec_high = None
        # WF MAP 컬러 anchor용 원본 spec(방향 nulling 전): speclow/spechigh 양끝 모두 보존
        wfmap_spec_low, wfmap_spec_high = spec_low, spec_high

        # REPORT DIRECTION: UPPER=상한만, LOWER=하한만, BOTH=둘 다
        direction = 'BOTH'
        if 'REPORT DIRECTION' in spec_data.columns:
            _d = str(spec_data.loc[spec_name, 'REPORT DIRECTION']).strip().upper()
            if _d in ('UPPER', 'LOWER', 'BOTH'):
                direction = _d
        if direction == 'UPPER':
            spec_low = None
        elif direction == 'LOWER':
            spec_high = None

        # REPORT LOG SCALE: 값 축 log10 적용 여부
        log_scale = False
        if 'REPORT LOG SCALE' in spec_data.columns:
            _ls = spec_data.loc[spec_name, 'REPORT LOG SCALE']
            log_scale = str(_ls).strip().lower() in ('true', '1', '1.0', 'yes')

        # UNIT: 축 라벨 단위
        unit = ''
        if 'UNIT' in spec_data.columns:
            _u = str(spec_data.loc[spec_name, 'UNIT']).strip()
            if _u and _u.lower() != 'nan':
                unit = _u
        _disp_item = display_name(item_name)   # 축 라벨은 후처리된 표시명 사용
        y_label = f"{_disp_item} [{unit}]" if unit else _disp_item

        return {
            'item_name': item_name, 'spec_name': spec_name,
            'df': item_df, 'cols': cols_task, 'cfg': cfg_task,
            'target_lot_id': target_lot_id,
            'target_root_lot_id': target_root_lot_id,
            'target_DC_step_id': target_DC_step_id,
            'spec': {'spec_low': spec_low, 'spec_high': spec_high,
                     'wfmap_spec_low': wfmap_spec_low, 'wfmap_spec_high': wfmap_spec_high,
                     'direction': direction, 'log_scale': log_scale, 'y_label': y_label},
        }

    # task는 지연 생성 + 슬라이딩 윈도(워커수×2)로 제출 — 항목별 데이터 슬라이스가
    # 한꺼번에 메모리에 쌓이지 않도록 하여 저사양(10GB)에서도 안전하게 동작한다.
    pool = _get_chart_pool(n_workers) if n_workers > 1 else None
    task_cache = {}      # i -> task dict | '_no_col' | '_empty'
    futures = {}         # i -> Future
    _state = {'next': 0, 'pool_ok': pool is not None}
    _WINDOW = (n_workers * 2 + 1) if pool is not None else 0

    def _get_task(i):
        if i in task_cache:
            return task_cache[i]
        _in, _sn = plot_items[i]
        if _in not in merged_df.columns:
            t = '_no_col'
        else:
            t = _build_task(_in, _sn) or '_empty'
        task_cache[i] = t
        return t

    def _pump():
        """윈도가 빌 때까지 다음 항목들을 풀에 제출(렌더링과 조립을 겹치게 유지)."""
        if pool is None or not _state['pool_ok']:
            return
        while _state['next'] < total_items:
            _busy = sum(1 for f in futures.values() if not f.done())
            if _busy >= _WINDOW:
                break
            i = _state['next']
            _state['next'] = i + 1
            t = _get_task(i)
            if isinstance(t, dict):
                try:
                    futures[i] = pool.submit(_render_item_charts, t)
                except Exception as _pe:
                    print(f"[WARN] 병렬 렌더링 제출 실패 → 직렬로 폴백: {_pe}")
                    _state['pool_ok'] = False
                    break

    _pump()

    # 슬라이드 공통 정적 요소: Wafer 색 범례는 항목과 무관 → 1회만 렌더링해 재사용
    leg_bytes = _render_wafer_legend_bytes(cfg_task)

    # --- 항목별 슬라이드 조립 루프 (렌더링 결과를 REPORT ORDER 순서로 소비) ---
    for idx, (item_name, spec_name) in enumerate(plot_items, start=1):
        i = idx - 1
        _pump()
        t = _get_task(i)
        if t == '_no_col':
            print(f"[{idx}/{total_items}] {item_name} 건너뜀 (merged_df에 데이터 없음)")
            continue
        print(f"[{idx}/{total_items}] {item_name} 처리 중...")

        # ---- 카테고리 간지(Description) '대기' 등록 (My_config.use_description_page로 on/off) ----
        #   실제 데이터가 찍힌 CAT2에만 간지를 붙이기 위해, 여기서는 대기(_pending_desc)로만 잡아둔다.
        #   아래에서 이 CAT2의 item이 '실제 chart 슬라이드'를 만들 때(=데이터 존재) 그 직전에 삽입한다.
        #   데이터가 하나도 없는 CAT2는 다음 CAT2로 넘어가며 대기가 덮어써져 간지가 붙지 않는다.
        if getattr(GLOBAL_CONFIG, 'use_description_page', True) and 'CAT2' in spec_data.columns:
            cat2 = str(spec_data.loc[spec_name, 'CAT2']).strip()
            if cat2 != current_cat and cat2.lower() != 'nan':
                current_cat = cat2
                matched_entry = None
                ck = cat2.lower().strip()
                # 1) CAT2 글자와 정확히 일치하는 description 슬라이드 우선
                for key, _entry in description_image_info_dict.items():
                    if key.lower().strip() == ck:
                        matched_entry = _entry
                        break
                # 2) 없으면 CAT2 글자가 포함(부분일치)된 description 슬라이드 탐색
                if matched_entry is None and ck:
                    for key, _entry in description_image_info_dict.items():
                        kl = key.lower().strip()
                        if ck in kl or kl in ck:
                            matched_entry = _entry
                            break
                _pending_desc = (cat2, matched_entry)   # 데이터가 실제 찍히면 chart 슬라이드 앞에 삽입

        if t == '_empty':
            task_cache.pop(i, None)
            continue

        # ---- 렌더링 결과 수령 (병렬이면 풀 Future, 아니면 이 자리에서 직렬 렌더) ----
        res = None
        if i in futures:
            try:
                res = futures.pop(i).result()
            except Exception as _fe:
                print(f"[WARN] 워커 렌더링 실패({item_name}) → 메인에서 직렬 재시도: {_fe}")
                _state['pool_ok'] = False
                res = None
        if res is None:
            res = _render_item_charts(t)
        task_cache.pop(i, None)   # task(df 슬라이스) 메모리 조기 해제

        for _w in res.get('warnings', []):
            print(f"[WARN] {item_name}: {_w}")
        if res['status'] == 'skip':
            print(f"[{idx}/{total_items}] {item_name} 건너뜀 ({res['reason']})")
            continue
        if res['status'] == 'error':
            print(f"[ERROR] {item_name} 차트 생성 중 에러 발생: {res['reason']}")
            continue

        if res.get('metrics'):
            metrics_dict[item_name] = res['metrics']
        if res.get('summary_row'):
            summary_rows.append(res['summary_row'])

        # 이 item은 실제 데이터가 있어 chart 슬라이드를 만든다 → 대기 중인 CAT2 간지를 이 앞에 삽입
        if _pending_desc is not None:
            _pd_cat, _pd_entry = _pending_desc
            _pending_desc = None
            if _pd_entry is not None:
                try:
                    if isinstance(_pd_entry, dict):
                        _copy_slide_into(prs, _pd_entry.get('slide'),
                                         src_w=_pd_entry.get('w'), src_h=_pd_entry.get('h'))
                    else:
                        _copy_slide_into(prs, _pd_entry)
                except Exception as _de:
                    print(f"[WARN] CAT2 '{_pd_cat}' description 슬라이드 복사 실패: {_de}")
            else:
                print(f"[WARN] CAT2 '{_pd_cat}' description 슬라이드를 찾지 못해 간지 생략")

        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            item_slide_map[item_name] = slide   # Anomaly 상세에서 이 아이템명 → 이 차트 슬라이드로 링크

            # 슬라이드 배경색 설정 (흰색)
            slide_bg = slide.background
            slide_bg.fill.solid()
            slide_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # ---- 레이아웃 상수 (와이드 16:9 13.333x7.5, 여백 최소화·꽉 채움) ----
            # 좌/우 2열을 슬라이드 폭(13.333")에 꽉 차게 배치, 여백 최소화
            LX, RX = 0.12, 8.50          # 좌/우 열 X 좌표 (테이블·그림·타이틀 공통)
            LW, RW = 8.30, 4.70          # 좌/우 열 폭 (LX+LW=8.42, RX+RW=13.20)
            TITLE_GAP = 0.22             # 타이틀과 바로 아래 그림 사이 간격
            SLIDE_BOTTOM = 7.40          # 차트가 채울 슬라이드 하단 경계 (7.5" - 하단여백)
            # 그림 top Y 좌표 (좌열: 통계표/레전드/Box/WFMAP, 우열: Trend/Radius/CDF)
            # 통계표를 키운 만큼 레전드/Box/WFMAP을 아래로 이동
            Y_TABLE, Y_LEG, Y_BOX, Y_MAP = 0.68, 2.34, 2.94, 5.14
            Y_TREND, Y_RAD, Y_CUM = 0.92, 3.08, 5.24

            # 타이틀 카드 타이틀 생성 헬퍼 함수 (그림 top - TITLE_GAP 위치에 배치)
            def add_card_title(text, l, pic_top):
                tx = slide.shapes.add_textbox(Inches(l), Inches(pic_top - TITLE_GAP), Inches(RW), Inches(0.3))
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(10.5)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*NAVY_RGB)
                p.font.name = FONT

            # INDEX 옆에 관련 REAL ITEM 이름 파싱해서 기입
            # ADDP는 ADDP FORM의 {ALIAS} 참조를 재귀 전개(ADDP 안의 ADDP도 풀어) 최종 leaf
            # alias를 얻고, 각 leaf alias의 reformatter ITEMID(=실제 real item)를 모두 표시한다.
            # (함수명 rmax/ABS 등은 {} 밖이라 제외, 다중출력 파생 {WINDOW_new}는 base ADDP로 재귀)
            def _leaf_to_itemids(leaf_aliases):
                _out = []
                for _r in leaf_aliases:
                    _r2 = re.sub(r'_new$', '', _r)                 # 다중출력 파생 접미사 제거
                    _iid = _ref_itemid.get(_r2) or _ref_itemid.get(_r)
                    _out.append(_iid if _iid else _r2)             # ITEMID 없으면 alias명 표시
                return sorted(set(_out))

            real_items_str = ""
            if reformatter is not None and spec_name in reformatter['ALIAS'].values:
                row = reformatter[reformatter['ALIAS'] == spec_name].iloc[0]
                cat = str(row.get('CATEGORY', '')).strip().upper()
                if cat == 'REAL':
                    real_items_str = display_name(str(row.get('ITEMID', '')))   # 표시명 후처리(suffix/prefix 제거)
                elif cat == 'ADDP':
                    _leaf = _addp_leaf_refs(spec_name)             # 최종 leaf alias(재귀 해소)
                    real_items_str = ", ".join(display_name(_x) for _x in _leaf_to_itemids(_leaf))

            # ---- 1. HEADER & INDEX NAME (Title) ----
            header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.58))
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(*NAVY_RGB)
            header_shape.line.fill.background()

            txBox = slide.shapes.add_textbox(Inches(0.2), Inches(0.06), Inches(12.9), Inches(0.45))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f" {display_name(item_name)}"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = FONT

            if real_items_str:
                run = p.add_run()
                run.text = f"  ({real_items_str})"
                run.font.size = Pt(11)
                run.font.bold = False
                run.font.color.rgb = RGBColor(200, 208, 224) # 연회색 소프트 폰트
                run.font.name = FONT

            # 네이티브 타이틀 추가 (각 그림 top 좌표 기준으로 정렬)
            add_card_title("■ Wafer-level Distribution (Box Plot)", LX, Y_BOX)
            add_card_title("■ WF MAP", LX, Y_MAP)
            add_card_title("■ Trend", RX, Y_TREND)
            add_card_title("■ Radius plot", RX, Y_RAD)
            add_card_title("■ Cumulative plot", RX, Y_CUM)

            # ---- 2. Statistical Table (lot_id 포함, 측정 wafer만, wafer 열폭 일정) ----
            _stat_rows = res['stat_rows']
            _slw = res.get('stat_cells_lw') or {}
            _NAVY = RGBColor(*NAVY_RGB); _WH = RGBColor(255, 255, 255); _BK = RGBColor(0, 0, 0)
            _LOTBG = RGBColor(0xD9, 0xE1, 0xF2)

            def _set_cell_border(cell, color="808080", w=6350):
                """PPT 표 셀 4변 테두리(구분선) 설정 — python-pptx 기본 API에 없어 XML로 직접 추가.
                tcPr 자식 순서상 ln* 은 fill 앞에 와야 하므로 fill/기타 요소 앞에 삽입."""
                from pptx.oxml import parse_xml
                from pptx.oxml.ns import qn, nsdecls
                tcPr = cell._tc.get_or_add_tcPr()
                _ref = None
                for _ch in list(tcPr):
                    if _ch.tag in (qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'),
                                   qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'),
                                   qn('a:headers'), qn('a:extLst')):
                        _ref = _ch; break
                for _tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                    _old = tcPr.find(qn(_tag))
                    if _old is not None:
                        tcPr.remove(_old)
                    _ln = parse_xml(
                        f'<{_tag} {nsdecls("a")} w="{w}" cap="flat" cmpd="sng" algn="ctr">'
                        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                        f'<a:prstDash val="solid"/></{_tag}>')
                    if _ref is not None:
                        _ref.addprevious(_ln)
                    else:
                        tcPr.append(_ln)

            def _stat_style(cell, text, bg, fg, bold, sz=8, align=PP_ALIGN.CENTER):
                cell.text = str(text)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.margin_left = Inches(0.0); cell.margin_right = Inches(0.0)
                cell.margin_top = Inches(0.0); cell.margin_bottom = Inches(0.0)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = False
                for par in cell.text_frame.paragraphs:
                    par.font.size = Pt(sz); par.font.bold = bold
                    par.font.color.rgb = fg; par.font.name = FONT
                    par.alignment = align
                _set_cell_border(cell)   # 셀 구분선(테두리) 표시

            if _slw:
                # 컬럼 순서: target lot 먼저, lot 내 wafer 오름차순 (HTML score board와 동일)
                _lots_s = sorted({k[0] for k in _slw},
                                 key=lambda l: (0 if str(l) == str(target_lot_id) else 1, str(l)))
                order_lw = []
                for _lot in _lots_s:
                    for _w in sorted([w for (l, w) in _slw if l == _lot]):
                        order_lw.append((_lot, _w))
                n_waf = max(len(order_lw), 1)
                # wafer 열폭: 기본 0.36" 고정(적게 찍히면 표가 줄어듦). 단 wafer가 많아도
                # 표가 우측 열(Trend/Radius/CDF)을 침범하지 않도록 좌측 열 폭(LW) 안으로
                # 균일 축소(표 내 모든 wafer 동일 폭 유지). 25매도 좌열에 다 들어가게 함.
                STAT_W = 0.62
                WAFER_W = min(0.36, (LW - STAT_W) / n_waf)
                ncol = 1 + n_waf
                nrow = 2 + len(_stat_rows)      # lot 헤더행 + wafer# 헤더행 + stat 행들
                tbl_w = STAT_W + n_waf * WAFER_W
                table_shape = slide.shapes.add_table(
                    nrow, ncol, Inches(LX), Inches(Y_TABLE),
                    Inches(tbl_w), Inches(min(1.75, 0.26 * nrow))).table
                table_shape.columns[0].width = Inches(STAT_W)
                for _j in range(1, ncol):
                    table_shape.columns[_j].width = Inches(WAFER_W)
                # (0,0)+(1,0) 병합 → 라벨 없이 빈 칸(요청: 'Stat' 표기 제거)
                table_shape.cell(0, 0).merge(table_shape.cell(1, 0))
                _stat_style(table_shape.cell(0, 0), "", _NAVY, _WH, True, 9)
                # 헤더행0: lot_id 가로 병합
                _k = 0
                while _k < len(order_lw):
                    _lot = order_lw[_k][0]; _k2 = _k
                    while _k2 + 1 < len(order_lw) and order_lw[_k2 + 1][0] == _lot:
                        _k2 += 1
                    _c0, _c1 = 1 + _k, 1 + _k2
                    if _c1 > _c0:
                        table_shape.cell(0, _c0).merge(table_shape.cell(0, _c1))
                    _stat_style(table_shape.cell(0, _c0), str(_lot), _LOTBG, _BK, True, 8)
                    _k = _k2 + 1
                # 헤더행1: wafer #
                for _j, (_lot, _w) in enumerate(order_lw):
                    _stat_style(table_shape.cell(1, 1 + _j), f"#{_w}", _NAVY, _WH, True, 8)
                # 데이터행: stat 라벨 + (lot,wafer)별 값
                for _r, _lbl in enumerate(_stat_rows):
                    _rr = 2 + _r
                    _bg = RGBColor(245, 247, 250) if _r % 2 == 0 else _WH
                    _stat_style(table_shape.cell(_rr, 0), _lbl, _bg, _BK, False, 8, PP_ALIGN.LEFT)
                    for _j, (_lot, _w) in enumerate(order_lw):
                        _vals = _slw.get((_lot, _w))
                        _txt = _vals[_r] if (_vals and _r < len(_vals)) else "-"
                        # 값이 8자 이상이면 6pt, 그 외 7pt
                        _stat_style(table_shape.cell(_rr, 1 + _j), _txt, _bg, _BK, False,
                                    6 if len(str(_txt)) >= 8 else 7)
            else:
                # (fallback) fab_lot_id 없을 때: 기존 wafer #1~25 고정 26칸 표
                cols = 26
                table_shape = slide.shapes.add_table(5, cols, Inches(LX), Inches(Y_TABLE), Inches(LW), Inches(1.2)).table
                headers = ["Stat"] + [f"#{w}" for w in range(1, 26)]
                for c_idx, h in enumerate(headers):
                    _stat_style(table_shape.cell(0, c_idx), h, _NAVY, _WH, True, 7)
                stat_cells = res['stat_cells']
                for r_idx, label in enumerate(_stat_rows):
                    table_shape.cell(r_idx + 1, 0).text = label
                for w_idx in range(1, 26):
                    vals = stat_cells.get(w_idx)
                    for r_idx in range(1, 5):
                        _v = vals[r_idx - 1] if (vals and r_idx - 1 < len(vals)) else "-"
                        _bg = RGBColor(240, 240, 240) if vals is None else (
                            RGBColor(245, 247, 250) if r_idx % 2 == 1 else _WH)
                        _stat_style(table_shape.cell(r_idx, w_idx), _v, _bg, _BK, False,
                                    6 if len(str(_v)) >= 8 else 7)
                for r_idx in range(1, 5):
                    _bg = RGBColor(245, 247, 250) if r_idx % 2 == 1 else _WH
                    _stat_style(table_shape.cell(r_idx, 0), _stat_rows[r_idx - 1] if r_idx - 1 < len(_stat_rows) else "",
                                _bg, _BK, False, 7, PP_ALIGN.LEFT)

            # ---- 3. 렌더링된 차트 이미지 삽입 (워커가 만든 jpg bytes) ----
            imgs = res['imgs']
            # Wafer Color Legend (전 슬라이드 공통 1회 렌더 재사용)
            slide.shapes.add_picture(_io.BytesIO(leg_bytes), Inches(LX), Inches(Y_LEG), width=Inches(LW))
            # BOX Plot
            slide.shapes.add_picture(_io.BytesIO(imgs['box']), Inches(LX), Inches(Y_BOX), Inches(LW), Inches(1.95))
            # WF MAP(grid) — 우측 컬러바 슬롯(_CBAR_SLOT)만큼 뺀 폭에 배치, 세로가 넘칠 때만 축소
            _CBAR_SLOT = 0.55
            map_avail = SLIDE_BOTTOM - Y_MAP
            map_area_w = LW - _CBAR_SLOT
            _ratio = res['map_ratio'] or 0.1
            pic_w, pic_h = map_area_w, map_area_w * _ratio
            if pic_h > map_avail:
                pic_h = map_avail
                pic_w = pic_h / _ratio
            slide.shapes.add_picture(_io.BytesIO(imgs['map']), Inches(LX), Inches(Y_MAP), Inches(pic_w), Inches(pic_h))
            # Color bar — 항상 우측에 같은 크기로 세로로 길게: top=Y_MAP(현재 위치), bottom≈Cumulative x축
            if imgs.get('cbar'):
                _cb_h = SLIDE_BOTTOM - Y_MAP
                _cb_x = LX + LW - _CBAR_SLOT + 0.02
                slide.shapes.add_picture(_io.BytesIO(imgs['cbar']), Inches(_cb_x), Inches(Y_MAP), height=Inches(_cb_h))
            # Trend / Radius / Cumulative
            slide.shapes.add_picture(_io.BytesIO(imgs['trend']), Inches(RX), Inches(Y_TREND), Inches(RW), Inches(1.95))
            slide.shapes.add_picture(_io.BytesIO(imgs['rad']), Inches(RX), Inches(Y_RAD), Inches(RW), Inches(1.95))
            slide.shapes.add_picture(_io.BytesIO(imgs['cum']), Inches(RX), Inches(Y_CUM), Inches(RW), Inches(2.15))

        except Exception as e:
            import traceback
            print(f"[ERROR] {item_name} 차트 생성 중 에러 발생: {e}")
            traceback.print_exc()

    print("=" * 60)
    print(f"[insert_plots] 차트 생성 완료 - 총 {total_items}개 index 처리")
    print("=" * 60)

    # ==================== 마지막 Index Aggregation Table 페이지(들) ====================
    # 값 산출 기준: REPORT DIRECTION → BOTH=Median / UPPER=P90 / LOWER=P10
    # index가 많으면 30개씩 페이지 분할. 모든 페이지가 동일한 (lot,wafer) 컬럼 구성을
    # 공유하여 칸 크기가 페이지마다 바뀌지 않도록 한다.
    if summary_rows:
        try:
            def _set_cell_borders(cell):
                from pptx.oxml.xmlchemy import OxmlElement
                tcPr = cell._tc.get_or_add_tcPr()
                for t in ['lnL', 'lnR', 'lnT', 'lnB']:
                    for el in tcPr.findall(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{t}'):
                        tcPr.remove(el)
                _bi = 0
                for t in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                    ln = OxmlElement(t); ln.set('w', '12700'); ln.set('cmpd', 'sng')
                    sf = OxmlElement('a:solidFill'); sc = OxmlElement('a:srgbClr'); sc.set('val', '333333')
                    sf.append(sc); ln.append(sf); tcPr.insert(_bi, ln); _bi += 1

            def _fmt_num(v):
                try:
                    f = float(v)
                except (ValueError, TypeError):
                    return str(v)
                if f != 0 and (abs(f) > 10000 or abs(f) <= 0.0001):
                    mant, exp = f"{f:.1E}".split('E')
                    return f"{mant}E{int(exp)}"
                if abs(f) >= 100:          # 100 이상은 소수점 첫째자리까지 (예: 1234.5)
                    return f"{f:.1f}"
                return f"{f:.3f}"          # 100 미만은 기존대로(소수점 3자리)

            def _style(cell, text, bg, fg, bold=False, sz=8, wrap=False):
                cell.text = str(text)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = wrap
                _set_cell_borders(cell)
                cell.margin_left = Inches(0.01); cell.margin_right = Inches(0.01)
                cell.margin_top = Inches(0.0); cell.margin_bottom = Inches(0.0)
                for par in cell.text_frame.paragraphs:
                    par.font.size = Pt(sz); par.font.bold = bold; par.font.name = FONT
                    par.font.color.rgb = fg; par.alignment = PP_ALIGN.CENTER

            # ---- (lot, wafer) 컬럼 구성(전체 index 기준, 모든 페이지 공통) ----
            _all_keys = set()
            for _r in summary_rows:
                _all_keys.update(_r.get('wafer_vals', {}).keys())
            _lots = sorted({k[0] for k in _all_keys},
                           key=lambda l: (0 if str(l) == str(target_lot_id) else 1, str(l)))
            def _wkey(k):
                return k[1] if isinstance(k[1], int) else 10 ** 9
            _lot_groups = []   # (lot, [(lot,wafer), ...])
            _ordered = []
            for _lot in _lots:
                _wc = sorted([k for k in _all_keys if k[0] == _lot], key=_wkey)
                if _wc:
                    _lot_groups.append((_lot, _wc))
                    _ordered.extend(_wc)

            NAVY = RGBColor(*NAVY_RGB)
            LOTBG = RGBColor(0xD9, 0xE1, 0xF2)
            WAFBG = RGBColor(0xF0, 0xF0, 0xF0)
            white = RGBColor(255, 255, 255); black = RGBColor(0, 0, 0)

            # 고정 열 너비: Index 넓게(잘림 방지) + wafer 좁게(상한 캡으로 일정 유지)
            _idx_w, _stat_w = 3.85, 0.55   # Index명 폭 확대(2.25→3.85, wafer 약 4칸분 추가)
            _ww = min(0.40, max(0.14, (13.333 - 0.24 - _idx_w - _stat_w) / max(len(_ordered), 1)))
            _tbl_w = min(13.10, _idx_w + _stat_w + _ww * len(_ordered))

            _CHUNK = 30
            _total_pages = (len(summary_rows) - 1) // _CHUNK + 1
            for _pg in range(_total_pages):
                _chunk = summary_rows[_pg * _CHUNK:(_pg + 1) * _CHUNK]
                s_slide = prs.slides.add_slide(prs.slide_layouts[6])
                s_slide.background.fill.solid()
                s_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # 상단 헤더 바
                hdr = s_slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.62))
                hdr.fill.solid(); hdr.fill.fore_color.rgb = NAVY; hdr.line.fill.background()
                htf = hdr.text_frame; htf.word_wrap = True; htf.margin_left = Inches(0.2)
                hp = htf.paragraphs[0]
                _sfx = f"  ({_pg + 1}/{_total_pages})" if _total_pages > 1 else ""
                hp.text = f"Index Aggregation Table{_sfx}"
                hp.font.size = Pt(18); hp.font.bold = True
                hp.font.color.rgb = RGBColor(255, 255, 255); hp.font.name = FONT
                hp.alignment = PP_ALIGN.LEFT

                ncol = 2 + len(_ordered)
                nrow = 2 + len(_chunk)
                tbl_h = Inches(min(6.5, max(0.26 * nrow, 0.6)))
                table = s_slide.shapes.add_table(nrow, ncol, Inches(0.12), Inches(0.78),
                                                 Inches(_tbl_w), tbl_h).table
                table.columns[0].width = Inches(_idx_w)
                table.columns[1].width = Inches(_stat_w)
                for _ci in range(2, ncol):
                    table.columns[_ci].width = Inches(_ww)

                # 헤더: Index/Stat 세로 병합, lot 가로 병합 + 그 아래 #wafer
                table.cell(0, 0).merge(table.cell(1, 0)); _style(table.cell(0, 0), "Index", NAVY, white, True, 9, wrap=False)
                table.cell(0, 1).merge(table.cell(1, 1)); _style(table.cell(0, 1), "Stat", NAVY, white, True, 8)
                _ci = 2
                for _lot, _wc in _lot_groups:
                    _c0 = _ci; _c1 = _ci + len(_wc) - 1
                    if _c1 > _c0:
                        table.cell(0, _c0).merge(table.cell(0, _c1))
                    _style(table.cell(0, _c0), str(_lot), LOTBG, black, True, 8)
                    for (_lot2, _waf) in _wc:
                        _style(table.cell(1, _ci), f"#{_waf}", WAFBG, black, True, 7)
                        _ci += 1

                # 본문: Index | Stat | (lot,wafer)별 값
                for _ri, _r in enumerate(_chunk, start=2):
                    _style(table.cell(_ri, 0), display_name(str(_r['index'])), white, black, True, 7, wrap=False)  # 표시명 후처리, Index명 45자 한 줄(줄바꿈 방지)
                    _style(table.cell(_ri, 1), str(_r.get('stat', '')), white, black, False, 8)
                    _wv = _r.get('wafer_vals', {})
                    for _cj, _key in enumerate(_ordered, start=2):
                        _val = _wv.get(_key)
                        _style(table.cell(_ri, _cj), '' if _val is None else _fmt_num(_val), white, black, False, 7)

            print(f"[insert_plots] Index Aggregation Table 생성 완료 ({len(summary_rows)}개 index, {_total_pages}페이지)")
        except Exception as _e:
            print(f"[WARN] Index Aggregation Table 생성 실패: {_e}")

    return prs, metrics_dict, item_slide_map


def getData_with_retry(params, custom_columns=None, user_name=None,
                       timeout_sec=1800, retry_wait_sec=300, max_retries=3, long_wait_sec=3600):
    """bigdatalake `getData`를 타임아웃 + 재시도로 감싼다(성공할 때까지 반환하지 않음).

    - 쿼리 1회가 timeout_sec(기본 30분)를 넘으면 강제 중단(daemon thread 포기)하고 재시도.
    - 에러/타임아웃 시: 처음 max_retries(3)회까지는 retry_wait_sec(5분) 간격으로 재시도.
    - 3회 모두 실패하면 그 이후로는 long_wait_sec(1시간) 간격으로 계속 재시도(성공할 때까지).
    → API 오류로 터미널이 멈추지 않고 자동 복구한다.
    """
    import threading
    import time as _time
    from bigdataquery import getData   # 지연 import — 실제 쿼리 시점에만 bigdataquery 로드

    attempt = 0
    while True:
        attempt += 1
        _res = {}

        def _run():
            try:
                _res['data'] = getData(params, custom_columns=custom_columns, user_name=user_name)
            except Exception as _e:
                _res['err'] = _e

        _t = threading.Thread(target=_run, daemon=True)
        _t.start()
        _t.join(timeout_sec)

        if _t.is_alive():
            print(f"[WARN] bigdatalake 쿼리 타임아웃({timeout_sec // 60}분 초과) — 강제 중단 후 재시도 (attempt {attempt})")
        elif 'err' in _res:
            print(f"[WARN] bigdatalake 쿼리 실패: {_res['err']} (attempt {attempt})")
        else:
            if attempt > 1:
                print(f"[INFO] bigdatalake 쿼리 재시도 성공 (attempt {attempt})")
            return _res['data']

        # 실패/타임아웃 → 대기 후 재시도 (3회까지는 5분, 이후 1시간 간격)
        if attempt < max_retries:
            _w = retry_wait_sec
            print(f"[WARN] {_w // 60}분 후 재시도합니다 ({attempt}/{max_retries})")
        else:
            _w = long_wait_sec
            print(f"[ERROR] bigdatalake 쿼리 {max_retries}회 연속 실패 — 이후 {_w // 60}분 간격으로 계속 재시도합니다")
        _time.sleep(_w)


def etdata_query():
    """ET(Electrical Test) 측정 데이터를 빅데이터 서버에서 쿼리하여 일별 Hive-파티셔닝 parquet로 저장.

    GLOBAL_CONFIG에서 읽어오는 주요 설정:
        - vehicle: 대상 차종(마스크) 이름
        - DB_et_daily: 일별 파켓 저장 루트 경로
        - QueryTimeSpan: 전체 쿼리 기간 (일)
        - SplitTimeSpan: 쿼리 분할 단위 (일)
        - process_id, line_id: 공정/라인 필터
        - et_custom_columns, user_name: 쿼리 파라미터
        - et_log_path: Lot 로그 CSV 경로

    저장 구조 (Hive Partitioning):
        ``{DB_et_daily}/date={YYYY-MM-DD}/data.parquet``
    """
    try : 

        item_et = pd.read_csv(f'reformatter/{GLOBAL_CONFIG.get("vehicle")}_reformatter.csv') 

        # 경로 생성 (DB_et_daily만 — LOTWF 디렉토리는 더 이상 사용하지 않음)
        if not os.path.exists(GLOBAL_CONFIG.get("DB_et_daily")):
            os.makedirs(GLOBAL_CONFIG.get("DB_et_daily"))

        current_time = datetime.now()
        formatted_time = current_time.strftime("%Y-%m-%d %H:%M")


        sub_datetime_now = datetime.now()
        if GLOBAL_CONFIG.get("now_minus") :
            sub_to_date_time = sub_datetime_now - timedelta(days = GLOBAL_CONFIG.get("now_minus"))
            sub_from_date_time = sub_datetime_now - timedelta(days = GLOBAL_CONFIG.get("QueryTimeSpan"))
        else :
            sub_to_date_time = sub_datetime_now 
            sub_from_date_time = sub_datetime_now - timedelta(days = GLOBAL_CONFIG.get("QueryTimeSpan"))
        dateTo = sub_to_date_time.strftime('%Y-%m-%d')
        dateFrom = sub_from_date_time.strftime('%Y-%m-%d')

        date_format = "%Y-%m-%d"
        start_date = datetime.strptime(dateFrom, date_format)
        end_date = datetime.strptime(dateTo, date_format)
        interval = timedelta(days=GLOBAL_CONFIG.get("SplitTimeSpan"))

        date_list = []
        current_date = start_date
        while current_date < end_date:
            date_list.append(current_date.strftime(date_format))
            current_date += interval

        if date_list[-1] != end_date:
            date_list.append(end_date.strftime(date_format))

        paired_date_list = [[date_list[i], date_list[i+1]] for i in range(len(date_list)-1)]
        print(f"[Date Ranges] {paired_date_list}")

        for paired_dates in paired_date_list :
            dateFrom = (datetime.strptime(paired_dates[0], date_format) + timedelta(days=1)).strftime(date_format)
            dateTo = paired_dates[1]
            print("\n" + "="*60)
            print(f"[Query Setting] {GLOBAL_CONFIG.get('SplitTimeSpan')}일치 Query")
            print(f"[Query Period] {dateFrom} ~ {dateTo}")

            is_real = item_et['CATEGORY'] == 'REAL' 
            is_addp = item_et['CATEGORY'] == 'ADDP' 
            real = item_et[is_real]
            addp = item_et[is_addp]
            real = real.loc[:,['ITEMID', 'ALIAS', 'SCALE FACTOR','ABSOLUTE']]
            ITEMID_List=real.ITEMID.tolist()
            addp = addp.loc[:,['ALIAS', 'ADDP FORM','SCALE FACTOR']]
            addp['addpscale'] = addp['SCALE FACTOR'].astype(str) + '*(' + addp['ADDP FORM'] + ')'
            ALIAS = list(map(str,addp.ALIAS)) #ADDP FORM ALIAS
            FORMULA = list(map(str,addp.addpscale)) #ADDP FORM FOMULA

            start_time = time.time()

            print(f"[Query Start] {GLOBAL_CONFIG.get('vehicle')} {GLOBAL_CONFIG.get('SplitTimeSpan')}일치 Query 시작")
            params = {
                        'table_name': 'eds.f_et_test',
                        'dateFrom': dateFrom, 
                        'dateTo': dateTo, 
                        'process_id' : GLOBAL_CONFIG.get("process_id"),
                        'line_id': GLOBAL_CONFIG.get("line_id"),
                        'item_id' : ITEMID_List,  #ITEM ID_REAL
                        'not_like_conditions': {'subitem_id' : ['Q%','AVG','MAX','MIN','RANGE','STD']}, #불필요 통계치
                        'like_conditions': {'step_seq' : GLOBAL_CONFIG.get("setting_stepseq")}
                        }

            Query_Table_tmp = getData_with_retry(params, custom_columns=GLOBAL_CONFIG.get("et_custom_columns"), user_name=GLOBAL_CONFIG.get("user_name"))

            print(f"[Query Complete] {GLOBAL_CONFIG.get('vehicle')} {dateFrom} ~ {dateTo} 데이터 추출 완료")

            Query_Table_tmp['tkout_time'] = pd.to_datetime(Query_Table_tmp['tkout_time'])
            Query_Table_tmp['et_value'] = pd.to_numeric(Query_Table_tmp['et_value'], errors='coerce')
            Query_Table_tmp['temperature'] = pd.to_numeric(Query_Table_tmp['temperature'], errors='coerce')
            Query_Table_tmp['et_value'] = Query_Table_tmp['et_value'].abs() #전체 absolute 적용
            daily_groups = Query_Table_tmp.groupby(Query_Table_tmp['tkout_time'].dt.date)

            #et_log 파일생성
            if os.path.exists(GLOBAL_CONFIG.get("et_log_path")):
                existing_lot_log = pd.read_csv(GLOBAL_CONFIG.get("et_log_path"))
                #형변환
                existing_lot_log['wafer_id'] = existing_lot_log['wafer_id'] .apply(eval)
                existing_lot_log['step_seq'] = existing_lot_log['step_seq'] .apply(eval)
                existing_lot_log['total_site_cnt'] = existing_lot_log['total_site_cnt'] .apply(eval)
            else:
                existing_lot_log = pd.DataFrame()

            lot_log = Query_Table_tmp.copy()
            lot_log['mask'] = GLOBAL_CONFIG.get("vehicle")
            lot_log['lot_id6'] = lot_log['lot_id'].str.split('_').str[0]
            Query_Table_tmp['wafer_id'] = pd.to_numeric(Query_Table_tmp['wafer_id'], errors='coerce')
            Query_Table_tmp['total_site_cnt'] = pd.to_numeric(Query_Table_tmp['total_site_cnt'], errors='coerce')
            lot_log['wafer_id'] = lot_log['wafer_id'].astype(int)
            lot_log['total_site_cnt'] = lot_log['total_site_cnt'].astype(int)
            lot_log['tkout_time'] = lot_log['tkout_time'].astype(str)
            lot_log['tkout_time'] = pd.to_datetime(lot_log['tkout_time'], errors='coerce')
            lot_log['prime_key'] = lot_log['mask'].astype(str) + '_' + lot_log['fab_lot_id'].astype(str) + '_' + lot_log['step_id'].astype(str)

            lot_log_unique = lot_log.groupby('prime_key').agg({
                'wafer_id': lambda x: x.unique().tolist(),  
                'step_seq': lambda x: x.unique().tolist(),  
                'total_site_cnt': lambda x: x.unique().tolist(),  
                'tkout_time': 'max'  
            }).reset_index()

            combined_lot_log = pd.DataFrame()
            final_lot_log = pd.DataFrame()
            combined_lot_log = pd.concat([existing_lot_log, lot_log_unique])
            combined_lot_log['tkout_time'] = pd.to_datetime(combined_lot_log['tkout_time'], errors='coerce')
            final_lot_log = combined_lot_log.groupby('prime_key').agg({
                'wafer_id': lambda x: list(sorted(set(sum(x, [])))),
                'step_seq': lambda x: list(sorted(set(sum(x, [])))),
                'total_site_cnt': lambda x: list(sorted(set(sum(x, [])))),
                'tkout_time': 'max' 
            }).reset_index()
            final_lot_log.to_csv(GLOBAL_CONFIG.get("et_log_path"), index = False)

            # Hive 파티셔닝 구조로 일별 parquet 저장
            # 저장 경로: {DB_et_daily}/date={YYYY-MM-DD}/data.parquet
            DB_et_daily = GLOBAL_CONFIG.get("DB_et_daily")
            for date_val, group in daily_groups:
                partition_dir = os.path.join(DB_et_daily, f'date={date_val}')
                os.makedirs(partition_dir, exist_ok=True)
                group.to_parquet(os.path.join(partition_dir, 'data.parquet'), index=False)

            end_time = time.time()
            elapsed_time = end_time - start_time

            print(f"[ET Query Complete] {GLOBAL_CONFIG.get('vehicle')} ET Query 완료 (소요시간: {elapsed_time:.2f}초)")
            print("="*60 + "\n") 

    except Exception as e:
        print(f"[ERROR] etdata_query 실패: {e}")
        traceback.print_exc()

# 사용중 Inline Query
def inlinedata_query(root_lot_id):
    """Inline(FAB 공정 내) 측정 데이터를 쿼리하여 CSV로 저장 후 DataFrame 반환.

    Parameters
    ----------
    root_lot_id : str
        대상 루트 Lot ID.

    Returns
    -------
    pd.DataFrame
        Inline 측정 데이터. 에러 시 빈 DataFrame 반환.
    """
    try: 
        setting_file = pd.read_excel(GLOBAL_CONFIG.get("inline_file_path"), sheet_name=None)
        Inline1 = setting_file[GLOBAL_CONFIG.get("inline_file_sheet")]
        Inline1_step_id_list = Inline1['STEP_DESC'].tolist()

        current_time_fab = datetime.now()
        date_To = current_time_fab 
        date_From= current_time_fab - timedelta(days=GLOBAL_CONFIG.get("inline_QueryTimeSpan"))

        params = {
                    'table_name': 'fab.f_fab_wf_met',
                    'dateFrom': date_From.strftime("%Y-%m-%d"),
                    'dateTo': date_To.strftime("%Y-%m-%d"),
                    'process_id' : GLOBAL_CONFIG.get("process_id"),
                    'line_id': GLOBAL_CONFIG.get("line_id"),
                    'root_lot_id' : root_lot_id,
                    'step_seq' : Inline1_step_id_list,
                    'not_like_conditions' : {'subitem_id' : ['SLOTID','MIN','MAX','Q2','RANGE','AVG','STD']}
                    }

        Query_Table = getData_with_retry(params, custom_columns=GLOBAL_CONFIG.get("inline_custom_columns"), user_name=GLOBAL_CONFIG.get("user_name"))

        Query_Table['fab_value'] = pd.to_numeric(Query_Table['fab_value'], errors='coerce')
        rows_to_multiply = Query_Table[Query_Table['item_id'].str.match(r'^CD\d+$')].index
        Query_Table.loc[rows_to_multiply, 'fab_value'] *= 1000

        Query_Table['spc_ctrl_spec_high'] = pd.to_numeric(Query_Table['spc_ctrl_spec_high'], errors='coerce')
        Query_Table['spc_ctrl_spec_limit'] = pd.to_numeric(Query_Table['spc_ctrl_spec_limit'], errors='coerce')
        Query_Table['spc_ctrl_spec_low'] = pd.to_numeric(Query_Table['spc_ctrl_spec_low'], errors='coerce')

        Query_Table.loc[rows_to_multiply, 'spc_ctrl_spec_high'] *= 1000
        Query_Table.loc[rows_to_multiply, 'spc_ctrl_spec_limit'] *= 1000
        Query_Table.loc[rows_to_multiply, 'spc_ctrl_spec_low'] *= 1000

        Query_Table.rename(columns={'step_seq': 'STEP_DESC'}, inplace=True)
        Query_Table = pd.merge(Query_Table, Inline1, on='STEP_DESC', how='left')

        Query_Table.to_csv(GLOBAL_CONFIG.get('DB') + f"{GLOBAL_CONFIG.get('vehicle')}_inline_table.csv", encoding='utf-8-sig', index=False)

        return Query_Table
    
    except Exception as e:
        print(f"inlinedata_query 에러가 발생했습니다: {e}")
        traceback.print_exc()
        return pd.DataFrame()

def wipdata_query():
    """WIP(Work In Progress) 현황 데이터를 쿼리하여 CSV로 저장.

    GLOBAL_CONFIG에서 읽어오는 주요 설정:
        - process_id, line_id: 공정/라인 필터
        - wip_custom_columns, user_name: 쿼리 파라미터
        - DB: 저장 디렉토리 경로
        - vehicle: 대상 차종(마스크) 이름
    """
    try : 
        params = {
                    'table_name': 'fab.f_wip_current',
                    'process_id' : GLOBAL_CONFIG.get("process_id"),
                    'line_id': GLOBAL_CONFIG.get("line_id"),
                    }

        Query_Table_tmp = getData_with_retry(params, custom_columns=GLOBAL_CONFIG.get("wip_custom_columns"), user_name=GLOBAL_CONFIG.get("user_name"))
        Query_Table_tmp['lot_id6'] = Query_Table_tmp['lot_id'].str.split('.').str[0]
        Query_Table_tmp.rename(columns={'step_seq': 'step_id'}, inplace=True)

        Query_Table_tmp.to_csv(GLOBAL_CONFIG.get('DB') + f"{GLOBAL_CONFIG.get('vehicle')}_wip_current.csv", index = False, encoding='cp949')
            
        print('wip data 추출완료')
    
    except Exception as e:
        print(f"wipdata_query 에러가 발생했습니다: {e}")
        traceback.print_exc()
